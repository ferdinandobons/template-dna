# SPDX-License-Identifier: MIT
"""Deterministic builder for the COMPLEX synthetic PPTX fixture.

Produces ``tests/fixtures/complex/acme_complex.pptx``: a 100% synthetic
(made-up "Acme Corp" brand, never proprietary) deck that stresses the
brand-pptx extractor/generator across as many PowerPoint component types as
python-pptx (plus a little raw lxml) can author from the default template:

  * A multi-PLACEHOLDER cover/title slide built on the default ``Title Slide``
    layout - title + subtitle + date + footer + slide-number placeholders all
    filled (the multi-placeholder cover the extractor inventories as cover
    anchors).
  * Slides on several DISTINCT layouts: ``Title Slide`` (cover),
    ``Section Header`` (agenda / section list), ``Title and Content``
    (content-text and content-table), ``Title Only`` (chart + the demo slide),
    ``Picture with Caption`` (picture), and a closing ``Section Header`` slide.
  * A real ``p14:sectionLst`` (PowerPoint *sections*) injected into
    ``presentation.xml`` via lxml - python-pptx 1.x does not model it - with
    THREE named sections ("Overview", "Financials", "Closing") whose
    ``p14:sldId`` lists point at real slides. The agenda slide's body lists
    those exact section names, so the section list and the agenda text agree.
  * A NATIVE table (``graphicFrame``/``a:tbl``) via ``shapes.add_table`` with a
    header row + body rows + a totals row.
  * A NATIVE chart (``graphicFrame``/``c:chart``) via ``shapes.add_chart`` with
    a small clustered-bar dataset (the embedded chart workbook is part of the
    package, exercising the relationship/parts walker).
  * A PICTURE: a small deterministic PNG generated in-process (no external /
    proprietary asset on disk) placed both as a free ``add_picture`` shape and
    as a "logo" mark on the cover.
  * A logo-like GROUPED-SHAPE mark + auto-shapes tinted with the synthetic Acme
    theme colors (approximating SmartArt - see NOTE below).
  * A DEMO / sample-content slide whose ONLY text equals a layout placeholder
    prompt (``"Click to edit Master title style"``) so the engine's
    demo-detection classifies it as a clearable demo region.

NOTE - what is APPROXIMATED / SKIPPED (python-pptx limits, by design):
  * SmartArt cannot be authored by python-pptx; it is APPROXIMATED with a
    grouped set of connected auto-shapes (a "process" row of boxes + arrows).
  * New custom slide masters / layouts cannot be created from scratch by
    python-pptx; this deck REUSES the default template's master + layouts (the
    task allows exactly this). Theme COLORS/FONTS are the ones that ship with
    the default template; an Acme palette is applied at the shape level.

Reproducibility: no randomness, no wall-clock in the deck itself. The
core-properties timestamps are pinned to a fixed instant and the embedded PNG
bytes are computed deterministically, so two rebuilds in the SAME environment
are identical. The output is content-reproducible within a fixed library set;
the committed binary is the source of truth, and rebuilds may differ
byte-for-byte across python-pptx / openpyxl / lxml versions (benign
serialization noise - including the embedded chart workbook's own timestamped
core.xml). Equality is therefore asserted STRUCTURALLY, not by raw bytes (see
tests/test_fixture_determinism.py).

Run:
    PYTHONPATH=scripts .venv/bin/python tests/fixtures/builders/build_complex_pptx.py
"""

from __future__ import annotations

import struct
import zlib
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Emu, Pt

OUT = Path(__file__).resolve().parents[1] / "complex" / "acme_complex.pptx"

# ---------------------------------------------------------------------------
# Synthetic "Acme Corp" brand palette (made-up; never proprietary).
# ---------------------------------------------------------------------------
ACME_NAVY = RGBColor(0x1F, 0x38, 0x64)
ACME_TEAL = RGBColor(0x2E, 0x8B, 0x8B)
ACME_AMBER = RGBColor(0xE8, 0xA3, 0x3D)
ACME_LIGHT = RGBColor(0xEA, 0xF0, 0xF6)
ACME_SLATE = RGBColor(0x55, 0x66, 0x77)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Default-template layout indices (verified against python-pptx 1.x default).
LO_TITLE = 0  # Title Slide      : CENTER_TITLE + SUBTITLE + date/footer/num
LO_TITLE_CONTENT = 1  # Title and Content: TITLE + OBJECT body
LO_SECTION = 2  # Section Header   : TITLE + BODY
LO_TITLE_ONLY = 5  # Title Only       : TITLE only
LO_PIC_CAPTION = 8  # Picture w/Caption: TITLE + PICTURE + BODY

# The exact default-template prompt string that the BODY/TITLE placeholders
# carry. A demo slide's only text must EQUAL one of these for the extractor's
# language-invariant demo-detection to fire. Captured here as a constant so the
# builder stays self-documenting; it is re-read live from the layout below to
# stay correct even if the bundled template ever changes.
TITLE_PROMPT = "Click to edit Master title style"

# Named PowerPoint sections to inject (name -> 0-based slide indices it spans).
# Filled in build() once the real slides (and their sldId values) exist.
SECTION_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"
SECTION_EXT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"


# ---------------------------------------------------------------------------
# A tiny synthetic PNG generated in-process (no external/proprietary asset).
# ---------------------------------------------------------------------------
def _synthetic_png() -> bytes:
    """Return bytes of a deterministic 96x40 RGBA PNG (an 'Acme' brand block).

    Navy field with an amber stripe and a teal corner - a purely decorative,
    made-up mark so the fixture carries a real ``<p:pic>`` / image part without
    committing any external image. Pixels are computed, never read from disk.
    """
    w, h = 96, 40
    navy = (0x1F, 0x38, 0x64, 0xFF)
    amber = (0xE8, 0xA3, 0x3D, 0xFF)
    teal = (0x2E, 0x8B, 0x8B, 0xFF)
    raw = bytearray()
    for y in range(h):
        raw.append(0)  # PNG filter type 0 (None) per scanline
        for x in range(w):
            if x < 14 and y < 14:
                px = teal
            elif 16 <= y < 24 and 6 <= x < 90:
                px = amber
            else:
                px = navy
            raw.extend(px)

    def _chunk(tag: bytes, data: bytes) -> bytes:
        return (
            struct.pack(">I", len(data))
            + tag
            + data
            + struct.pack(">I", zlib.crc32(tag + data) & 0xFFFFFFFF)
        )

    sig = b"\x89PNG\r\n\x1a\n"
    ihdr = struct.pack(">IIBBBBB", w, h, 8, 6, 0, 0, 0)  # 8-bit RGBA
    idat = zlib.compress(bytes(raw), 9)
    return sig + _chunk(b"IHDR", ihdr) + _chunk(b"IDAT", idat) + _chunk(b"IEND", b"")


# ---------------------------------------------------------------------------
# Small text helpers.
# ---------------------------------------------------------------------------
def _set_text(placeholder, text: str, *, size=None, bold=None, color=None) -> None:
    """Set placeholder/shape text in a single run with optional run formatting."""
    tf = placeholder.text_frame
    tf.text = text
    run = tf.paragraphs[0].runs[0]
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color


def _bullets(shape, items, *, size=18, color=ACME_NAVY) -> None:
    """Fill a shape's text frame with one bullet paragraph per (text, level)."""
    text_frame = shape.text_frame
    text_frame.clear()
    for i, (txt, level) in enumerate(items):
        para = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
        para.text = txt
        para.level = level
        run = para.runs[0]
        run.font.size = Pt(size)
        run.font.color.rgb = color


# ---------------------------------------------------------------------------
# Slide builders. Each returns the created slide so build() can collect sldIds.
# ---------------------------------------------------------------------------
def _add_cover(prs, png_path: Path):
    """Multi-PLACEHOLDER cover: title + subtitle + date + footer + slide-number.

    Every placeholder the ``Title Slide`` layout exposes is filled, so the
    extractor's cover-anchor inventory sees a genuine multi-placeholder cover.
    A small 'logo' picture is dropped in the top-left corner.
    """
    layout = prs.slide_layouts[LO_TITLE]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(
        ph[0],
        "Acme Corp Quarterly Business Review",
        size=40,
        bold=True,
        color=ACME_NAVY,
    )
    _set_text(
        ph[1], "FY2026 - Performance, Outlook & Initiatives", size=20, color=ACME_TEAL
    )
    if 10 in ph:  # DATE placeholder
        _set_text(ph[10], "January 15, 2026", size=12, color=ACME_SLATE)
    if 11 in ph:  # FOOTER placeholder
        _set_text(
            ph[11],
            "Acme Corp - Confidential (synthetic sample)",
            size=12,
            color=ACME_SLATE,
        )
    if 12 in ph:  # SLIDE_NUMBER placeholder
        _set_text(ph[12], "1", size=12, color=ACME_SLATE)
    # Logo-like picture in the corner (the synthetic Acme mark).
    slide.shapes.add_picture(
        str(png_path), Emu(457200), Emu(457200), height=Emu(457200)
    )
    return slide


def _add_agenda(prs, section_names):
    """Agenda / section-list slide: body lists the deck's real section names.

    Built on the ``Section Header`` layout (TITLE + BODY). The body text is the
    exact list of ``p14:section`` names injected later, so the agenda and the
    real section list agree (the extractor surfaces both).
    """
    layout = prs.slide_layouts[LO_SECTION]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(ph[0], "Agenda", size=36, bold=True, color=ACME_NAVY)
    _bullets(
        ph[1],
        [(f"{i + 1}. {name}", 0) for i, name in enumerate(section_names)],
        size=20,
    )
    return slide


def _add_content_text(prs):
    """Content-text slide: title + multi-level bulleted body (Title and Content)."""
    layout = prs.slide_layouts[LO_TITLE_CONTENT]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(ph[0], "Executive Summary", size=32, bold=True, color=ACME_NAVY)
    _bullets(
        ph[1],
        [
            ("Acme Corp grew net revenue 18% YoY across all four regions.", 0),
            ("Gross margin expanded 230 bps on supply-chain efficiencies.", 0),
            ("North region led growth; East rebounded after Q1 softness.", 1),
            ("FY2026 outlook raised on a stronger services pipeline.", 0),
            ("Key risk: input-cost volatility in the Gadget line.", 1),
        ],
        size=18,
    )
    return slide


def _add_content_table(prs):
    """Content-table slide: title + a NATIVE table (graphicFrame/a:tbl).

    A 5-row x 5-col table with a styled header band, body rows tinted with the
    Acme palette, and a totals row - exercises the table extractor path.
    """
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(
        ph[0], "Regional Revenue (USD thousands)", size=28, bold=True, color=ACME_NAVY
    )

    rows, cols = 5, 5
    left, top, width, height = Emu(457200), Emu(1828800), Emu(8229600), Emu(3200400)
    gtable = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = gtable.table
    headers = ["Region", "Q1", "Q2", "Q3", "Q4"]
    body = [
        ["North", "320", "351", "372", "410"],
        ["South", "210", "228", "241", "265"],
        ["East", "175", "168", "199", "232"],
    ]
    totals = ["Total", "705", "747", "812", "907"]

    for c, label in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = label
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACME_NAVY
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(16)
        run.font.color.rgb = WHITE
    for r, datarow in enumerate(body, start=1):
        for c, val in enumerate(datarow):
            cell = table.cell(r, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACME_LIGHT if r % 2 else WHITE
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size = Pt(14)
            run.font.color.rgb = ACME_NAVY
    for c, val in enumerate(totals):
        cell = table.cell(rows - 1, c)
        cell.text = val
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACME_TEAL
        run = cell.text_frame.paragraphs[0].runs[0]
        run.font.bold = True
        run.font.size = Pt(15)
        run.font.color.rgb = WHITE
    return slide


def _add_chart(prs):
    """Chart slide: title + a NATIVE clustered-bar chart (graphicFrame/c:chart)."""
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(
        ph[0], "Quarterly Net Revenue by Region", size=28, bold=True, color=ACME_NAVY
    )

    chart_data = CategoryChartData()
    chart_data.categories = ["Q1", "Q2", "Q3", "Q4"]
    chart_data.add_series("North", (320, 351, 372, 410))
    chart_data.add_series("South", (210, 228, 241, 265))
    chart_data.add_series("East", (175, 168, 199, 232))

    left, top, width, height = Emu(457200), Emu(1828800), Emu(8229600), Emu(4114800)
    gframe = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
    )
    chart = gframe.chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Acme Corp - FY2026 Revenue"
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    return slide


def _add_picture_slide(prs, png_path: Path):
    """Picture slide on ``Picture with Caption``: TITLE + PICTURE ph + caption."""
    layout = prs.slide_layouts[LO_PIC_CAPTION]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(ph[0], "The Acme Brand Mark", size=28, bold=True, color=ACME_NAVY)
    # Fill the PICTURE placeholder (idx 1) with the synthetic mark...
    if 1 in ph:
        try:
            ph[1].insert_picture(str(png_path))
        except Exception:
            slide.shapes.add_picture(
                str(png_path), Emu(457200), Emu(1828800), height=Emu(914400)
            )
    else:
        slide.shapes.add_picture(
            str(png_path), Emu(457200), Emu(1828800), height=Emu(914400)
        )
    # ...and write a caption in the BODY placeholder (idx 2).
    if 2 in ph:
        _bullets(
            ph[2],
            [
                ("A synthetic, generated mark - no proprietary asset.", 0),
                ("Navy field, amber stripe, teal corner.", 0),
            ],
            size=16,
        )
    # A second, free-floating copy via add_picture to exercise that path too.
    slide.shapes.add_picture(
        str(png_path), Emu(5486400), Emu(1828800), height=Emu(914400)
    )
    return slide


def _add_smartart_approx(prs):
    """Closing/SmartArt-approximation slide: a grouped 'process' of shapes.

    python-pptx CANNOT author SmartArt, so this APPROXIMATES a 3-step process
    diagram with grouped rounded rectangles connected by arrows, tinted with the
    Acme palette. Documented as an approximation in the module docstring.
    """
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(
        ph[0],
        "Our Approach (SmartArt-style process)",
        size=28,
        bold=True,
        color=ACME_NAVY,
    )

    steps = [("Discover", ACME_NAVY), ("Build", ACME_TEAL), ("Scale", ACME_AMBER)]
    box_w, box_h = Emu(2286000), Emu(1143000)
    top = Emu(2743200)
    gap = Emu(457200)
    x = Emu(685800)
    centers_y = int(top) + int(box_h) // 2
    prev_right = None
    for label, color in steps:
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, top, box_w, box_h)
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.color.rgb = ACME_SLATE
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.text = label
        run = p.runs[0]
        run.font.size = Pt(20)
        run.font.bold = True
        run.font.color.rgb = WHITE
        if prev_right is not None:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, prev_right, Emu(centers_y), x, Emu(centers_y)
            )
            conn.line.color.rgb = ACME_SLATE
            conn.line.width = Pt(2.5)
        prev_right = int(x) + int(box_w)
        x = Emu(int(x) + int(box_w) + int(gap))
    return slide


def _add_demo_slide(prs):
    """DEMO slide: its only text EQUALS a layout placeholder prompt.

    The engine classifies a slide as a clearable 'demo' region when every text
    run on it equals a layout/master placeholder prompt. We build a Title Only
    slide and set the title to the exact default-template title prompt, so the
    slide's only text is an unedited prompt -> demo-detection fires.
    """
    layout = prs.slide_layouts[LO_TITLE_ONLY]
    slide = prs.slides.add_slide(layout)
    title = slide.placeholders[0]
    # Re-read the live prompt so this stays correct if the template changes.
    layout_title = layout.placeholders[0]
    prompt = (
        (layout_title.text or "").strip()
        if getattr(layout_title, "has_text_frame", False)
        else ""
    )
    title.text_frame.text = prompt or TITLE_PROMPT
    return slide


def _add_closing(prs):
    """Closing slide on ``Section Header`` (a distinct closing layout use)."""
    layout = prs.slide_layouts[LO_SECTION]
    slide = prs.slides.add_slide(layout)
    ph = {p.placeholder_format.idx: p for p in slide.placeholders}
    _set_text(ph[0], "Thank You", size=40, bold=True, color=ACME_NAVY)
    _bullets(
        ph[1],
        [
            ("Questions? hello@acme.example", 0),
            ("Acme Corp - synthetic sample deck.", 0),
        ],
        size=18,
    )
    return slide


# ---------------------------------------------------------------------------
# Section list (p14:sectionLst) injection via raw lxml (python-pptx can't reach).
# ---------------------------------------------------------------------------
def _inject_sections(prs, sections) -> None:
    """Inject a real ``p14:sectionLst`` into ``presentation.xml``.

    ``sections`` is a list of ``(name, [slide_index, ...])``. We read the live
    ``p:sldIdLst`` to map slide indices to their ``r:id``/``id`` and emit a
    ``<p:ext uri="{...}"><p14:sectionLst>...</p14:sectionLst></p:ext>`` block
    under ``p:extLst`` - the exact shape the extractor's ``detect_sections``
    walks. Deterministic section GUIDs keep the output reproducible.
    """
    pres = prs.part._element  # <p:presentation>
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"

    # Map slide index -> the sldId 'id' attribute (PowerPoint section sldId uses it).
    sld_id_lst = pres.find(f"{{{p_ns}}}sldIdLst")
    sld_ids = list(sld_id_lst) if sld_id_lst is not None else []
    index_to_id = {i: sld.get("id") for i, sld in enumerate(sld_ids)}

    # Build (or find) the presentation-level extLst, AFTER sldSz/notesSz per schema
    # order; appending at the end of <p:presentation> is schema-valid for extLst.
    ext_lst = pres.find(f"{{{p_ns}}}extLst")
    if ext_lst is None:
        ext_lst = etree.SubElement(pres, f"{{{p_ns}}}extLst")

    ext = etree.SubElement(ext_lst, f"{{{p_ns}}}ext")
    ext.set("uri", SECTION_EXT_URI)
    # Bind the ``p14`` prefix EXPLICITLY on the sectionLst element via nsmap so
    # serialization is deterministic. Relying on lxml's global
    # ``register_namespace`` registry yields a process-dependent prefix
    # (``ns0`` vs ``p14``), which would break byte-reproducibility.
    section_lst = etree.SubElement(
        ext, f"{{{SECTION_NS}}}sectionLst", nsmap={"p14": SECTION_NS}
    )
    # Pin a deterministic GUID per section (index-derived) for reproducibility.
    for si, (name, slide_indices) in enumerate(sections):
        section = etree.SubElement(section_lst, f"{{{SECTION_NS}}}section")
        section.set("name", name)
        guid = "{%08X-0000-4000-8000-%012X}" % (si + 1, si + 1)
        section.set("id", guid)
        sld_id_lst_el = etree.SubElement(section, f"{{{SECTION_NS}}}sldIdLst")
        for idx in slide_indices:
            sid = index_to_id.get(idx)
            if sid is None:
                continue
            sld = etree.SubElement(sld_id_lst_el, f"{{{SECTION_NS}}}sldId")
            sld.set("id", sid)


# ---------------------------------------------------------------------------
# Build.
# ---------------------------------------------------------------------------
def build(out: Path = OUT) -> Path:
    out = Path(out)
    png_bytes = _synthetic_png()

    # Materialize the synthetic PNG to a temp file (python-pptx wants a path or
    # a stream; we use a BytesIO-backed temp via NamedTemporaryFile-free path).
    import tempfile

    tmp_png = Path(tempfile.gettempdir()) / "acme_complex_mark.png"
    tmp_png.write_bytes(png_bytes)

    prs = Presentation()  # default template (master + 11 layouts + theme)
    # Drop nothing: the default presentation ships with zero slides.

    # The three named sections and the agenda must agree, so define names first.
    section_names = ["Overview", "Financials", "Closing"]

    # --- Build slides in deck order; collect their 0-based indices. -----------
    s0 = _add_cover(prs, tmp_png)  # 0 cover
    _add_agenda(prs, section_names)  # 1 agenda / section list
    s2 = _add_content_text(prs)  # 2 content-text
    s3 = _add_content_table(prs)  # 3 content-table (native table)
    s4 = _add_chart(prs)  # 4 native chart
    s5 = _add_picture_slide(prs, tmp_png)  # 5 picture
    s6 = _add_smartart_approx(prs)  # 6 grouped-shape "SmartArt"
    s7 = _add_demo_slide(prs)  # 7 DEMO (prompt-only text)
    s8 = _add_closing(prs)  # 8 closing
    _ = (s0, s2, s3, s4, s5, s6, s7, s8)

    # --- Inject the real PowerPoint section list (lxml). ----------------------
    # Section spans (0-based slide indices):
    #   Overview   -> cover, agenda, exec-summary           (0,1,2)
    #   Financials -> table, chart                          (3,4)
    #   Closing    -> picture, smartart, demo, thank-you    (5,6,7,8)
    sections = [
        ("Overview", [0, 1, 2]),
        ("Financials", [3, 4]),
        ("Closing", [5, 6, 7, 8]),
    ]
    _inject_sections(prs, sections)

    # --- Pin core-properties timestamps for byte-reproducibility. -------------
    fixed = datetime(2026, 1, 15, 0, 0, 0, tzinfo=timezone.utc)
    cp = prs.core_properties
    cp.author = "Acme Corp (synthetic fixture)"
    cp.title = "Acme Corp Quarterly Business Review"
    cp.subject = "Synthetic test fixture - no proprietary content"
    cp.created = fixed
    cp.modified = fixed
    cp.last_modified_by = "Acme Corp (synthetic fixture)"
    cp.revision = 1

    out.parent.mkdir(parents=True, exist_ok=True)
    # Save via a buffer first so we never leave a half-written file on error.
    buf = BytesIO()
    prs.save(buf)
    out.write_bytes(buf.getvalue())
    return out


if __name__ == "__main__":
    path = build()
    print(f"built {path}")
