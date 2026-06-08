# SPDX-License-Identifier: MIT
"""Regression tests for the PPTX correctness + brand-guarantee fixes.

Covers the confirmed findings:
  - C3  the extractor DERIVES roles from the real parsed layouts: every
        resolver points at a layout.name that actually exists in the deck, and
        when no suitable layout exists the role is an honest ``stub`` (no
        fabricated ``"Title Slide"`` / ``"Title and Content"``). The legacy
        literal ``surface.pptx.role_layout_map`` is gone.
  - M6  generation resolves layouts from the profile's real role data, never a
        hardcoded layout name or positional ``slide_layouts[0]/[1]`` fallback.
  - M7  slides are built from the IR block stream - one slide per heading (its
        own runs as the title), following blocks as that slide's body - with no
        flattening and without dropping tables / quotes / captions / lists.

All decks are synthesized in a temp dir with python-pptx and never committed.
"""

from __future__ import annotations

import sys
import tempfile
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

import json

import lxml.etree as ET
from pptx import Presentation
from pptx.util import Inches

from docx import Document

from brandkit.formats.docx import generate as docx_generate
from brandkit.formats.pptx import extract as px
from brandkit.formats.pptx import generate as pg
from brandkit.formats.pptx import structure as ps
from brandkit.ir import components as ir_components
from brandkit.ir import model as ir
from brandkit.ir.model import parse_idoc
from brandkit.profile import schema, store
from brandkit.qa import checks_deterministic as cd
from brandkit.qa.model import Finding


# ---------------------------------------------------------------------------
# Synthetic deck builders (temp only - NEVER committed)
# ---------------------------------------------------------------------------
def _branded_template(path: Path) -> None:
    """A deck whose layouts carry NON-default names, to prove the extractor and
    generator never depend on literal ``"Title Slide"`` / ``"Title and Content"``.

    python-pptx's default index-0 layout has a CENTER_TITLE + SUBTITLE (a cover)
    and index-1 has a TITLE + OBJECT body (a content layout); we only rename them.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.slide_layouts[0].name = "BrandCover"
    prs.slide_layouts[1].name = "BrandContent"
    prs.save(path)


def _placeholderless_template(path: Path) -> None:
    """A deck whose layouts expose NO title/body placeholders at all (forces the
    honest ``stub`` path)."""
    prs = Presentation()
    for layout in prs.slide_layouts:
        sp_tree = layout.shapes._spTree
        for ph in list(layout.placeholders):
            sp_tree.remove(ph._element)
    prs.save(path)


def _multibody_template(path: Path) -> None:
    """A deck whose CONTENT layout exposes TWO body placeholders at distinct idxs.

    python-pptx's default layout index 3 ("Two Content") carries a title (idx 0)
    plus two OBJECT body placeholders (idx 1 and idx 2). Renaming it lets a test
    prove the body-content path writes into the profile-NAMED ``ph_idx`` (idx 2,
    the second body) rather than the positional first body placeholder (idx 1).
    Index 0 ("Title Slide") is renamed as the cover.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.slide_layouts[0].name = "BrandCover"
    prs.slide_layouts[3].name = "BrandTwoBody"
    prs.save(path)


def _extract_profile(template: Path, name: str = "deck") -> dict:
    """Build a profile the way ``extract`` does, in-memory (no disk store)."""
    prs = Presentation(template)
    layouts = px._layouts(prs)
    roles = px._roles(prs, layouts)
    profile = schema.build_envelope(
        "pptx",
        {"name": name, "display_name": name},
        theme=px._theme(),
        roles=roles,
        surface={
            "pptx": {
                "slide_size_emu": {
                    "w": int(prs.slide_width),
                    "h": int(prs.slide_height),
                },
                "layouts": layouts,
                "safe_area_emu": {"l": 457200, "t": 457200, "r": 457200, "b": 457200},
            }
        },
    )
    profile["capabilities"] = px._capabilities()
    return profile


def _slide_body(slide) -> str:
    for shape in slide.placeholders:
        if shape != slide.shapes.title and shape.has_text_frame:
            return shape.text
    return ""


def _all_text(prs: Presentation) -> str:
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append("\t".join(cell.text for cell in row.cells))
    return "\n".join(parts)


# ---------------------------------------------------------------------------
# C3 - roles derived from the REAL deck, not fabricated
# ---------------------------------------------------------------------------
class C3RolesFromRealLayouts(unittest.TestCase):
    def test_resolver_layouts_actually_exist_in_deck(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            template = Path(td) / "branded.pptx"
            _branded_template(template)
            profile = _extract_profile(template)

            layout_names = set(profile["surface"]["pptx"]["layouts"])
            roles = profile["roles"]
            for rid in roles["_index"]:
                layout = roles[rid]["resolver"].get("layout")
                # Every non-stub role names a layout the deck actually has.
                if layout is not None:
                    self.assertIn(
                        layout,
                        layout_names,
                        f"role {rid} points at fabricated layout {layout!r}",
                    )

            # The specific renamed real layouts are the ones chosen.
            self.assertEqual(roles["cover.title"]["resolver"]["layout"], "BrandCover")
            self.assertEqual(roles["heading.1"]["resolver"]["layout"], "BrandContent")
            self.assertEqual(roles["paragraph"]["resolver"]["layout"], "BrandContent")

    def test_no_fabricated_literal_names(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            template = Path(td) / "branded.pptx"
            _branded_template(template)
            profile = _extract_profile(template)

            roles = profile["roles"]
            for rid in roles["_index"]:
                layout = roles[rid]["resolver"].get("layout")
                self.assertNotIn(layout, {"Title Slide", "Title and Content"})

            # The legacy literal map must be gone from the surface.
            self.assertNotIn("role_layout_map", profile["surface"]["pptx"])

    def test_profile_validates_clean(self) -> None:
        # The intra-profile consistency check (Core) would flag a fabricated
        # layout; a derived profile must pass with zero problems.
        with tempfile.TemporaryDirectory() as td:
            template = Path(td) / "branded.pptx"
            _branded_template(template)
            profile = _extract_profile(template)
            self.assertEqual(schema.validate(profile), [])

    def test_placeholder_idx_is_real(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            template = Path(td) / "branded.pptx"
            _branded_template(template)
            profile = _extract_profile(template)
            layouts = profile["surface"]["pptx"]["layouts"]
            roles = profile["roles"]
            for rid in roles["_index"]:
                resolver = roles[rid]["resolver"]
                layout = resolver.get("layout")
                ph_idx = resolver.get("ph_idx")
                if layout is None:
                    continue
                idxs = {ph["idx"] for ph in layouts[layout]["placeholders"]}
                self.assertIn(
                    ph_idx, idxs, f"{rid}: ph_idx {ph_idx} absent from {layout}"
                )

    def test_honest_stub_when_no_suitable_layout(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            template = Path(td) / "bare.pptx"
            _placeholderless_template(template)
            profile = _extract_profile(template)

            roles = profile["roles"]
            for rid in ("cover.title", "heading.1", "paragraph"):
                self.assertEqual(roles[rid]["status"], schema.Status.STUB.value)
                self.assertIsNone(roles[rid]["resolver"]["layout"])
                self.assertEqual(roles[rid]["confidence"], 0.0)
                self.assertFalse(roles[rid]["verified"])

            # A stub profile still validates (no fabricated target to contradict).
            self.assertEqual(schema.validate(profile), [])

    def test_body_text_degrades_loudly_when_no_body_placeholder(self) -> None:
        # A placeholderless layout has nowhere to put body text. The engine must
        # DEGRADE LOUDLY (block_degraded), never silently drop the content.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "bare.pptx"
            _placeholderless_template(template)
            profile = _extract_profile(template)
            idoc = ir.IntermediateDocument(
                blocks=[
                    ir.Heading(level=1, runs=[{"t": "Section"}]),
                    ir.Paragraph(runs=[{"t": "Body text with nowhere to go."}]),
                ]
            )
            out = tp / "out.pptx"
            sink: list[Finding] = []
            pg.generate(profile, template, idoc, out, findings=sink)
            self.assertTrue(out.is_file())
            degraded = [f for f in sink if f.check == "block_degraded"]
            self.assertTrue(
                any("paragraph" in f.message for f in degraded),
                f"body text must degrade loudly, not vanish: {[f.message for f in sink]}",
            )

    def test_real_extract_reports_cover_anchor_honestly(self) -> None:
        # Drive the real on-disk extract() to confirm anchors track reality:
        # present for a deck with a cover layout, absent for a placeholder-less one.
        import json

        for builder, expect_present in (
            (_branded_template, True),
            (_placeholderless_template, False),
        ):
            with tempfile.TemporaryDirectory() as td:
                tp = Path(td)
                template = tp / "tpl.pptx"
                builder(template)
                profile_path = px.extract(template, "deck", scope="project", cwd=tp)
                profile = json.loads(Path(profile_path).read_text())
                cover = profile["anchors"]["cover"]
                anchors = profile["surface"]["pptx"]["cover_anchors"]
                if expect_present:
                    # M-i-7: the cover is MULTI-placeholder - every placeholder on the
                    # cover layout surfaces as its own anchor (title + subtitle + date
                    # + footer + ...), not a single hardcoded slot. The count tracks
                    # the real placeholder count and a title-family anchor is present.
                    self.assertGreaterEqual(cover["slots_found"], 1)
                    self.assertEqual(cover["slots_found"], len(anchors))
                    self.assertTrue(
                        any(a.get("family") == "title" for a in anchors), anchors
                    )
                    self.assertEqual(cover["kind"], schema.AnchorKind.PLACEHOLDER.value)
                else:
                    self.assertEqual(cover["slots_found"], 0)
                    self.assertEqual(anchors, [])
                    self.assertEqual(cover["kind"], schema.AnchorKind.NONE.value)
                # The on-disk profile must also validate clean either way.
                self.assertEqual(schema.validate(profile), [])


# ---------------------------------------------------------------------------
# M6 / M7 - generation uses real layouts + the IR block stream
# ---------------------------------------------------------------------------
class M6M7GenerateFromRealLayouts(unittest.TestCase):
    def _generate(self, template: Path, idoc_dict: dict, out: Path):
        profile = _extract_profile(template)
        idoc = parse_idoc(idoc_dict)
        pg.generate(profile, template, idoc, out)
        return Presentation(out)

    def test_generation_uses_only_real_layouts(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            out = tp / "out.pptx"
            res = self._generate(
                template,
                {
                    "cover": {"title": "Board Update", "subtitle": "Q2"},
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "Section A"},
                        {"type": "paragraph", "text": "Body."},
                    ],
                },
                out,
            )
            used = {s.slide_layout.name for s in res.slides}
            # Only the real renamed layouts - never a fabricated literal.
            self.assertTrue(used <= {"BrandCover", "BrandContent"}, used)
            self.assertNotIn("Title Slide", used)
            self.assertNotIn("Title and Content", used)

    def test_one_slide_per_heading_with_distinct_titles(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            out = tp / "out.pptx"
            res = self._generate(
                template,
                {
                    "cover": {"title": "Deck"},
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "Section A"},
                        {"type": "paragraph", "text": "Alpha."},
                        {"type": "heading", "level": 1, "text": "Section B"},
                        {"type": "paragraph", "text": "Beta."},
                    ],
                },
                out,
            )
            titles = [
                s.shapes.title.text
                for s in res.slides
                if s.shapes.title and s.shapes.title.text
            ]
            # Both headings are their OWN slide titles (not "Section A (2)").
            self.assertIn("Section A", titles)
            self.assertIn("Section B", titles)

    def test_heading_text_not_duplicated_into_body(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            out = tp / "out.pptx"
            res = self._generate(
                template,
                {
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "Section A"},
                        {"type": "paragraph", "text": "Just the body here."},
                    ]
                },
                out,
            )
            for slide in res.slides:
                body = _slide_body(slide)
                self.assertNotIn("Section A", body)

    def test_tables_quotes_captions_lists_not_dropped(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            out = tp / "out.pptx"
            res = self._generate(
                template,
                {
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "Mixed"},
                        {
                            "type": "list",
                            "items": [
                                {"text": "first bullet"},
                                {"text": "second bullet"},
                            ],
                        },
                        {
                            "type": "table",
                            "columns": ["Area", "Status"],
                            "rows": [["Pipeline", "Healthy"]],
                        },
                        {
                            "type": "quote",
                            "text": "A pithy remark.",
                            "attribution": "Anon",
                        },
                        {"type": "caption", "text": "Figure 1. A caption."},
                        {"type": "callout", "intent": "info", "text": "Mind the gap."},
                    ]
                },
                out,
            )
            text = _all_text(res)
            for needle in (
                "first bullet",
                "second bullet",
                "Pipeline",
                "Healthy",
                "pithy remark",
                "Anon",
                "Figure 1",
                "Mind the gap",
            ):
                self.assertIn(needle, text, f"content {needle!r} was dropped")

    def test_long_section_splits_into_continuation_slides(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            out = tp / "out.pptx"
            res = self._generate(
                template,
                {
                    "cover": {"title": "Long"},
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "Long Section"},
                        {"type": "paragraph", "text": " ".join(["capacity"] * 420)},
                    ],
                },
                out,
            )
            # cover + at least two content slides for the over-capacity paragraph.
            self.assertGreaterEqual(len(res.slides), 3)
            content_titles = [
                s.shapes.title.text
                for s in res.slides
                if s.shapes.title and s.shapes.title.text
            ]
            # The continuation slide carries the SAME section title (suffixed).
            self.assertIn("Long Section", content_titles)
            self.assertTrue(any(t.startswith("Long Section (") for t in content_titles))

    def test_no_cover_slide_when_idoc_has_no_cover(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            out = tp / "out.pptx"
            res = self._generate(
                template,
                {"blocks": [{"type": "heading", "level": 1, "text": "Only Content"}]},
                out,
            )
            self.assertTrue(
                all(s.slide_layout.name != "BrandCover" for s in res.slides)
            )


# ---------------------------------------------------------------------------
# Body placeholder is chosen by the profile-resolved ph_idx (not positionally)
# ---------------------------------------------------------------------------
def _placeholder_text_by_idx(slide, idx: int) -> str:
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx and shape.has_text_frame:
            return shape.text
    return ""


class BodyPlaceholderHonorsResolvedIdx(unittest.TestCase):
    """Body content is written into the placeholder the profile RESOLVED (its
    ``paragraph`` role's ``ph_idx``), not merely the positional first body
    placeholder. The two diverge only on a multi-body layout; this proves the
    schema's resolved ``ph_idx`` is honored (compute-validate-USE, not
    compute-validate-ignore), while the positional fallback still covers a stub /
    absent idx so the brand guarantee never fabricates a placeholder.
    """

    def _multibody_profile(self, template: Path, *, body_idx: int) -> dict:
        """Extract a profile from the multi-body deck, then point BOTH content roles
        at the "BrandTwoBody" layout with ``paragraph`` resolving to ``body_idx``
        (the second body placeholder). This is the source-of-truth contract the
        generator must honor."""
        profile = _extract_profile(template)
        roles = profile["roles"]
        roles["heading.1"]["resolver"] = {
            "type": "placeholder",
            "layout": "BrandTwoBody",
            "ph_idx": 0,
            "ph_type": "title",
        }
        roles["paragraph"]["resolver"] = {
            "type": "placeholder",
            "layout": "BrandTwoBody",
            "ph_idx": body_idx,
            "ph_type": "body",
        }
        return profile

    def test_body_written_into_named_idx_not_first(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "twobody.pptx"
            _multibody_template(template)
            # The deck's content layout exposes TWO body placeholders (idx 1 + 2);
            # the profile names the SECOND (idx 2) for body content.
            profile = self._multibody_profile(template, body_idx=2)
            out = tp / "out.pptx"
            idoc = parse_idoc(
                {
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "Section A"},
                        {"type": "paragraph", "text": "Body into the named idx."},
                    ]
                }
            )
            pg.generate(profile, template, idoc, out)
            res = Presentation(out)
            content = next(
                s for s in res.slides if s.slide_layout.name == "BrandTwoBody"
            )
            # The content landed in the NAMED placeholder (idx 2), not the first
            # body placeholder (idx 1, which stays empty).
            self.assertIn(
                "Body into the named idx.", _placeholder_text_by_idx(content, 2)
            )
            self.assertEqual(_placeholder_text_by_idx(content, 1), "")

    def test_first_body_used_when_named_idx_is_first(self) -> None:
        # Control: when the profile names the FIRST body idx (idx 1), the resolved
        # selection and the positional fallback agree - content lands in idx 1.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "twobody.pptx"
            _multibody_template(template)
            profile = self._multibody_profile(template, body_idx=1)
            out = tp / "out.pptx"
            idoc = parse_idoc(
                {
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "Section A"},
                        {"type": "paragraph", "text": "Body into the first body."},
                    ]
                }
            )
            pg.generate(profile, template, idoc, out)
            res = Presentation(out)
            content = next(
                s for s in res.slides if s.slide_layout.name == "BrandTwoBody"
            )
            self.assertIn(
                "Body into the first body.", _placeholder_text_by_idx(content, 1)
            )
            self.assertEqual(_placeholder_text_by_idx(content, 2), "")

    def test_falls_back_to_first_body_when_idx_absent(self) -> None:
        # When the profile names an idx the slide does NOT carry (stub-like), the
        # body-placeholder selection falls back to the positional first body
        # placeholder - never fabricating a placeholder, never dropping content.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "twobody.pptx"
            _multibody_template(template)
            # idx 99 is absent from the layout -> fallback to the first body (idx 1).
            profile = self._multibody_profile(template, body_idx=99)
            out = tp / "out.pptx"
            idoc = parse_idoc(
                {
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "Section A"},
                        {"type": "paragraph", "text": "Fallback body text."},
                    ]
                }
            )
            pg.generate(profile, template, idoc, out)
            res = Presentation(out)
            content = next(
                s for s in res.slides if s.slide_layout.name == "BrandTwoBody"
            )
            self.assertIn("Fallback body text.", _placeholder_text_by_idx(content, 1))


# ---------------------------------------------------------------------------
# M-i-7 - fact enrichment + reconcile-not-rebuild
# ---------------------------------------------------------------------------
_P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
_SECT_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"


def _q(ns: str, tag: str) -> str:
    return f"{{{ns}}}{tag}"


def _add_section_list(prs, *, sections: list[tuple[str, list[int]]]) -> None:
    """Attach a p14:sectionLst to a deck. ``sections`` maps name -> slide indices."""
    pres = prs.part._element
    ids = [sid.get("id") for sid in prs.slides._sldIdLst]
    ext_lst = ET.SubElement(pres, _q(_P_NS, "extLst"))
    ext = ET.SubElement(ext_lst, _q(_P_NS, "ext"))
    ext.set("uri", _SECT_URI)
    sect_lst = ET.SubElement(ext, _q(_P14, "sectionLst"))
    for n, (name, idxs) in enumerate(sections):
        sec = ET.SubElement(sect_lst, _q(_P14, "section"))
        sec.set("name", name)
        sec.set("id", "{%032d}" % n)
        sld_lst = ET.SubElement(sec, _q(_P14, "sldIdLst"))
        for i in idxs:
            e = ET.SubElement(sld_lst, _q(_P14, "sldId"))
            e.set("id", ids[i])


def _reconcile_deck(path: Path) -> str:
    """A deck with: a cover slide, an agenda/section-list slide (NON-English title,
    body listing the section names), a real structural slide (authored text), a demo
    slide (body == layout prompt), and a two-section section list. Returns the demo
    slide's layout body prompt so the caller can assert it is purged.

    The agenda slide's title is deliberately NON-English ("Sommario") so the
    reconcile must carry the template's OWN word forward, not inject "Agenda".
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.slide_layouts[0].name = "BrandCover"
    prs.slide_layouts[1].name = "BrandContent"
    content_layout = prs.slide_layouts[1]
    body_prompt = ""
    for ph in content_layout.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame and ph.text:
            body_prompt = ph.text.strip()
            break

    prs.slides.add_slide(prs.slide_layouts[0])  # 0: cover (prompts only)
    sa = prs.slides.add_slide(content_layout)  # 1: agenda (NON-English)
    sa.shapes.title.text = "Sommario"
    pg._first_body_placeholder(sa).text = "Intro\nBody"
    s2 = prs.slides.add_slide(content_layout)  # 2: structural (authored)
    s2.shapes.title.text = "Real Structural Slide"
    pg._first_body_placeholder(s2).text = "Authored content that must be preserved."
    s3 = prs.slides.add_slide(content_layout)  # 3: demo (prompts only)
    s3.shapes.title.text = (
        content_layout.placeholders[0].text.strip() or "Click to edit"
    )
    if body_prompt:
        pg._first_body_placeholder(s3).text = body_prompt

    _add_section_list(prs, sections=[("Intro", [0]), ("Body", [1, 2, 3])])
    prs.save(path)
    return body_prompt


def _reconcile_deck_demo_not_last(path: Path) -> None:
    """A reconcile deck whose DEMO slide (index 3) is NOT the highest-indexed part:
    two more structural slides follow it. Clearing it must not orphan a higher slide
    part that the subsequent add_slide then collides with (duplicate ZIP part name).
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.slide_layouts[0].name = "BrandCover"
    prs.slide_layouts[1].name = "BrandContent"
    content = prs.slide_layouts[1]
    body_prompt = ""
    for ph in content.placeholders:
        if ph.placeholder_format.idx != 0 and ph.has_text_frame and ph.text:
            body_prompt = ph.text.strip()
            break
    prs.slides.add_slide(prs.slide_layouts[0])  # 0 cover
    sa = prs.slides.add_slide(content)  # 1 agenda
    sa.shapes.title.text = "Sommario"
    pg._first_body_placeholder(sa).text = "Intro\nBody"
    s2 = prs.slides.add_slide(content)  # 2 structural
    s2.shapes.title.text = "Structural One"
    pg._first_body_placeholder(s2).text = "Authored content one."
    s3 = prs.slides.add_slide(content)  # 3 DEMO (body == layout prompt)
    s3.shapes.title.text = content.placeholders[0].text.strip() or "Click to edit"
    if body_prompt:
        pg._first_body_placeholder(s3).text = body_prompt
    # Structural slides AFTER the demo -> the demo is not the highest part index.
    s4 = prs.slides.add_slide(content)  # 4 structural
    s4.shapes.title.text = "Structural Two"
    pg._first_body_placeholder(s4).text = "Authored content two."
    s5 = prs.slides.add_slide(content)  # 5 structural
    s5.shapes.title.text = "Structural Three"
    pg._first_body_placeholder(s5).text = "Authored content three."
    _add_section_list(prs, sections=[("Intro", [0]), ("Body", [1, 2, 3, 4, 5])])
    prs.save(path)


def _extract_on_disk(template: Path, tp: Path, name: str = "deck") -> dict:
    profile_path = px.extract(template, name, scope="project", cwd=tp)
    return json.loads(Path(profile_path).read_text())


def _present_comp(prof: dict, comp: dict, *, confidence: float = 0.9) -> None:
    prof.setdefault("provenance", {}).setdefault("shell", {})
    sha = prof["provenance"]["shell"]["sha256"]
    block = schema.empty_comprehension()
    block["status"] = schema.ComprehensionStatus.PRESENT.value
    block["source_shell_sha256"] = sha
    block.update(comp)
    # The explicit kwarg wins over any confidence carried in ``comp`` (e.g. the
    # default _comp_for value), so a test can force a below-floor confidence.
    block["confidence"] = confidence
    prof["comprehension"] = block


def _titles(prs) -> list[str]:
    return [
        s.shapes.title.text
        for s in prs.slides
        if s.shapes.title and s.shapes.title.text
    ]


class MI7FactEnrichmentTest(unittest.TestCase):
    """The extractor surfaces the format-uniform inventory + drops the literals."""

    def test_multi_placeholder_cover_anchors_surfaced(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            prof = _extract_on_disk(template, tp)
            anchors = prof["surface"]["pptx"]["cover_anchors"]
            # The cover is MULTI-placeholder: more than one anchor, each keyed by
            # its layout+ph idx, with a title-family anchor present.
            self.assertGreater(len(anchors), 1, anchors)
            self.assertTrue(all(a["id"].startswith("ph.") for a in anchors))
            self.assertTrue(any(a["family"] == "title" for a in anchors))
            # Every anchor carries the captured layout prompt (demo value) + ph type.
            for a in anchors:
                self.assertIn("placeholder", a)
                self.assertIn("ph_type", a)
            self.assertEqual(schema.validate(prof), [])

    def test_real_section_list_surfaces_as_field(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "sectioned.pptx"
            _reconcile_deck(template)
            prof = _extract_on_disk(template, tp)
            s = prof["surface"]["pptx"]
            self.assertEqual([f["id"] for f in s["fields"]], ["field.sections"])
            self.assertEqual(s["fields"][0]["section_count"], 2)
            self.assertEqual([sec["name"] for sec in s["sections"]], ["Intro", "Body"])
            self.assertTrue(prof["anchors"]["sections"]["present"])

    def test_demo_region_and_toc_literals_are_dropped(self) -> None:
        # A deck with NO slides must not report demo_region present (the old
        # ``bool(prs.slides)``) nor a toc-always-False stub.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "empty.pptx"
            _branded_template(template)  # 0 slides
            prof = _extract_on_disk(template, tp)
            self.assertFalse(prof["anchors"]["demo_region"]["present"])
            self.assertNotIn("toc", prof["anchors"])

    def test_demo_slide_detected_by_prompt_equality(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "deck.pptx"
            _reconcile_deck(template)
            prof = _extract_on_disk(template, tp)
            regions = {r["id"]: r["kind"] for r in prof["surface"]["pptx"]["regions"]}
            # slide 0 cover, slide 1 agenda (structural), slide 2 structural
            # (authored), slide 3 demo (prompts).
            self.assertEqual(regions["region.slide.0"], "cover")
            self.assertEqual(regions["region.slide.1"], "structural")
            self.assertEqual(regions["region.slide.2"], "structural")
            self.assertEqual(regions["region.slide.3"], "demo")
            self.assertTrue(prof["anchors"]["demo_region"]["present"])


class MI7ReconcileNotRebuildTest(unittest.TestCase):
    """Comprehension-steered generation keeps structural slides, fills the cover in
    place, clears only demo slides, and regenerates the agenda from new headings."""

    def _idoc(self) -> ir.IntermediateDocument:
        return ir.IntermediateDocument(
            blocks=[
                ir.Heading(level=1, runs=[{"t": "New Section One"}]),
                ir.Paragraph(runs=[{"t": "Fresh body text."}]),
                ir.Heading(level=1, runs=[{"t": "New Section Two"}]),
                ir.Paragraph(runs=[{"t": "More fresh body."}]),
            ],
            cover=ir.Cover(
                title=[{"t": "Recon Title"}], subtitle=[{"t": "Recon Subtitle"}]
            ),
        )

    def _comp_for(self, prof: dict, *, demo_slide_ref: str = "region.slide.3") -> dict:
        anchors = prof["surface"]["pptx"]["cover_anchors"]
        title_anchor = next(a["id"] for a in anchors if a["family"] == "title")
        sub_anchor = next((a["id"] for a in anchors if a["family"] == "subtitle"), None)
        slots = {
            title_anchor: {
                "binds_to": "title",
                "fill_rule": "in_place",
                "demo_value": "",
            }
        }
        if sub_anchor:
            slots[sub_anchor] = {
                "binds_to": "subtitle",
                "fill_rule": "in_place",
                "demo_value": "",
            }
        return {
            # A confident classification (clears the destructive floor): a demo clear
            # represents a confident "this is boilerplate" judgement. The CLI path
            # round-trips this through comprehension.merge, which propagates it.
            "confidence": 0.9,
            "cover_slots": slots,
            "conventions": {
                "indexes": [
                    {
                        "index_ref": "field.sections",
                        "reconcile": "regenerate",
                        "seq_id": None,
                    }
                ],
                "sections": [],
            },
            "demo_classification": {
                "regions": [{"region_ref": demo_slide_ref, "verdict": "demo"}]
            },
        }

    def test_keeps_structural_clears_demo_fills_cover_regens_agenda(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "deck.pptx"
            body_prompt = _reconcile_deck(template)
            prof = _extract_on_disk(template, tp)
            _present_comp(prof, self._comp_for(prof))
            self.assertTrue(store.comprehension_is_present(prof))

            out = tp / "out.pptx"
            findings: list[Finding] = []
            pg.generate(prof, template, self._idoc(), out, findings=findings)
            res = Presentation(out)
            titles = _titles(res)
            bodies = "\n".join(
                pg._first_body_placeholder(s).text
                for s in res.slides
                if pg._first_body_placeholder(s) is not None
            )

            # Cover filled IN PLACE on the cover layout (not recreated, no duplicate).
            cover_slides = [
                s for s in res.slides if s.slide_layout.name == "BrandCover"
            ]
            self.assertEqual(len(cover_slides), 1)
            self.assertEqual(cover_slides[0].shapes.title.text, "Recon Title")
            self.assertEqual(titles.count("Recon Title"), 1)

            # Structural slide KEPT (its authored text survives).
            self.assertIn("Real Structural Slide", titles)
            self.assertIn("Authored content that must be preserved.", bodies)

            # Demo slide CLEARED: its layout prompt text is gone from the deck.
            if body_prompt:
                self.assertNotIn(body_prompt, bodies)

            # New body content appended.
            self.assertIn("New Section One", titles)
            self.assertIn("New Section Two", titles)

            # Agenda refreshed IN PLACE, carrying the template's OWN (non-English)
            # list title forward - NOT an injected "Agenda" literal. Proves
            # language-invariance: the Italian "Sommario" survives, "Agenda" never
            # appears, and there is exactly ONE agenda page (no stale duplicate): the
            # existing "Sommario" slide's body is rewritten to the NEW headings.
            self.assertIn("Sommario", titles)
            self.assertNotIn("Agenda", titles)
            self.assertEqual(titles.count("Sommario"), 1)
            sommario_bodies = [
                pg._first_body_placeholder(s).text.strip()
                for s in res.slides
                if s.shapes.title
                and s.shapes.title.text == "Sommario"
                and pg._first_body_placeholder(s) is not None
            ]
            self.assertEqual(sommario_bodies, ["New Section One\nNew Section Two"])
            # The stale section list ("Intro"/"Body") is gone.
            self.assertNotIn("Intro\nBody", bodies)
            self.assertTrue(any(f.check == "agenda_regenerated" for f in findings))

            # No net structure loss (the demo clear was corroborated).
            self.assertFalse(any(f.check == "no_net_structure_loss" for f in findings))

    def test_demo_clear_downgraded_when_not_corroborated(self) -> None:
        # The model tags the STRUCTURAL slide (authored text) demo: determinism does
        # not corroborate (its text is not a layout prompt) -> KEEP + WARNING.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "deck.pptx"
            _reconcile_deck(template)
            prof = _extract_on_disk(template, tp)
            comp = self._comp_for(prof, demo_slide_ref="region.slide.1")
            _present_comp(prof, comp)

            out = tp / "out.pptx"
            findings: list[Finding] = []
            pg.generate(prof, template, self._idoc(), out, findings=findings)
            res = Presentation(out)
            # The wrongly-tagged structural slide is KEPT, a WARNING is recorded, and
            # nothing was removed without corroboration (no net-loss ERROR).
            self.assertIn("Real Structural Slide", _titles(res))
            self.assertTrue(any(f.check == "demo_clear_downgraded" for f in findings))
            self.assertFalse(
                any(
                    f.check == "no_net_structure_loss" and f.severity == "ERROR"
                    for f in findings
                )
            )

    def test_demo_clear_downgraded_when_below_confidence_floor(self) -> None:
        # The demo slide IS corroborated by determinism (body == layout prompt), but
        # the model's confidence is below the destructive floor (0.3 < 0.5): the
        # slide is KEPT + WARNING, uniform with the docx/xlsx/cover sites. A wrong
        # delete is not recoverable, so low confidence keeps it.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "deck.pptx"
            body_prompt = _reconcile_deck(template)
            prof = _extract_on_disk(template, tp)
            _present_comp(prof, self._comp_for(prof), confidence=0.3)

            out = tp / "out.pptx"
            findings: list[Finding] = []
            pg.generate(prof, template, self._idoc(), out, findings=findings)
            res = Presentation(out)
            bodies = "\n".join(
                pg._first_body_placeholder(s).text
                for s in res.slides
                if pg._first_body_placeholder(s) is not None
            )
            # The corroborated demo slide is KEPT (its prompt body survives) because
            # confidence did not clear the floor, and the downgrade is surfaced.
            if body_prompt:
                self.assertIn(body_prompt, bodies)
            self.assertTrue(
                any(
                    f.check == "demo_clear_downgraded" and "confidence" in f.message
                    for f in findings
                )
            )
            # KEEP-on-low-confidence is not a net loss, so no ERROR.
            self.assertFalse(
                any(
                    f.check == "no_net_structure_loss" and f.severity == "ERROR"
                    for f in findings
                )
            )

    def test_clearing_a_non_last_demo_slide_yields_no_duplicate_parts(self) -> None:
        # Regression (synthetic, isolated): a demo clear where the demo slide is NOT
        # the highest-indexed part must not orphan a higher slide part that the next
        # add_slide then collides with (duplicate ZIP part name -> corrupt OPC
        # package). Complements the example-template reconcile test; this builds the
        # exact mid-index-demo geometry the original 4-slide fixture (demo last) missed.
        import zipfile
        from collections import Counter

        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "deck.pptx"
            _reconcile_deck_demo_not_last(template)
            prof = _extract_on_disk(template, tp)
            _present_comp(prof, self._comp_for(prof, demo_slide_ref="region.slide.3"))
            out = tp / "out.pptx"
            pg.generate(prof, template, self._idoc(), out)
            names = zipfile.ZipFile(out).namelist()
            dups = [n for n, c in Counter(names).items() if c > 1]
            self.assertEqual(dups, [], f"duplicate package parts: {dups}")
            Presentation(out)  # reopens cleanly

    def test_generate_twice_is_idempotent(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "deck.pptx"
            _reconcile_deck(template)
            prof = _extract_on_disk(template, tp)
            _present_comp(prof, self._comp_for(prof))

            o1 = tp / "o1.pptx"
            o2 = tp / "o2.pptx"
            pg.generate(prof, template, self._idoc(), o1)
            pg.generate(prof, template, self._idoc(), o2)

            def sig(p: Path):
                r = Presentation(p)
                return [
                    (
                        s.slide_layout.name,
                        s.shapes.title.text if s.shapes.title else "",
                        pg._first_body_placeholder(s).text
                        if pg._first_body_placeholder(s)
                        else "",
                    )
                    for s in r.slides
                ]

            self.assertEqual(sig(o1), sig(o2))

    def test_existing_cover_slide_is_filled_not_duplicated(self) -> None:
        # When the deck already SHIPS a cover slide, the reconcile fills it in place
        # (one cover slide, no duplicate) rather than adding a second one.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "deck.pptx"
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)
            prs.slide_layouts[0].name = "BrandCover"
            prs.slide_layouts[1].name = "BrandContent"
            cover = prs.slides.add_slide(prs.slide_layouts[0])
            cover.shapes.title.text = "OLD COVER TITLE"
            prs.save(template)
            prof = _extract_on_disk(template, tp)
            anchors = prof["surface"]["pptx"]["cover_anchors"]
            title_anchor = next(a["id"] for a in anchors if a["family"] == "title")
            _present_comp(
                prof,
                {
                    "cover_slots": {
                        title_anchor: {
                            "binds_to": "title",
                            "fill_rule": "in_place",
                            "demo_value": "",
                        }
                    }
                },
            )
            out = tp / "out.pptx"
            findings: list[Finding] = []
            pg.generate(prof, template, self._idoc(), out, findings=findings)
            res = Presentation(out)
            cover_slides = [
                s for s in res.slides if s.slide_layout.name == "BrandCover"
            ]
            self.assertEqual(len(cover_slides), 1)
            self.assertEqual(cover_slides[0].shapes.title.text, "Recon Title")
            self.assertNotIn("OLD COVER TITLE", _titles(res))

    def test_absent_comprehension_uses_deterministic_rebuild(self) -> None:
        # With comprehension absent, the deck is blind-rebuilt (today's behavior):
        # the structural slide is NOT preserved (all slides cleared), proving the
        # reconcile path is gated on a present, sha-current comprehension.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "deck.pptx"
            _reconcile_deck(template)
            prof = _extract_on_disk(template, tp)
            self.assertEqual(prof["comprehension"]["status"], "absent")

            out = tp / "out.pptx"
            pg.generate(prof, template, self._idoc(), out)
            res = Presentation(out)
            titles = _titles(res)
            # Deterministic rebuild: cover from the IR + one slide per heading; the
            # template's own "Real Structural Slide" is NOT carried over.
            self.assertNotIn("Real Structural Slide", titles)
            self.assertIn("Recon Title", titles)
            self.assertIn("New Section One", titles)


class MI7CliFindingsWiringTest(unittest.TestCase):
    """The CLI generate path folds pptx reconciliation/destructive findings into the
    ONE QA report (parity with docx/xlsx). Drives ``cli.main`` end to end, not the
    generator directly, so the CLI wiring itself is covered."""

    def _idoc_json(self, tmp: Path) -> Path:
        p = tmp / "idoc.json"
        p.write_text(
            json.dumps(
                {
                    "cover": {"title": "Recon Title", "subtitle": "Recon Subtitle"},
                    "blocks": [
                        {"type": "heading", "level": 1, "text": "New Section One"},
                        {"type": "paragraph", "text": "Fresh body text."},
                        {"type": "heading", "level": 1, "text": "New Section Two"},
                        {"type": "paragraph", "text": "More fresh body."},
                    ],
                }
            ),
            encoding="utf-8",
        )
        return p

    def _run_cli_reconcile(self, tmp: Path, *, demo_slide_ref: str):
        """Extract a reconcile deck via the CLI, attach a present comprehension via
        the CLI ``comprehend`` command, then CLI-generate. Returns (rc, stdout)."""
        import io
        import os
        from contextlib import redirect_stdout

        from brandkit.cli import main

        template = tmp / "deck.pptx"
        _reconcile_deck(template)
        old = Path.cwd()
        os.chdir(tmp)
        try:
            self.assertEqual(
                main(
                    [
                        "extract",
                        "--name",
                        "deck",
                        "--template",
                        str(template),
                        "--scope",
                        "project",
                    ]
                ),
                0,
            )
            prof_path = tmp / "brand-kit" / "deck" / "profile.json"
            prof = json.loads(prof_path.read_text())
            comp = MI7ReconcileNotRebuildTest()._comp_for(
                prof, demo_slide_ref=demo_slide_ref
            )
            comp_path = tmp / "comp.json"
            comp_path.write_text(json.dumps(comp), encoding="utf-8")
            self.assertEqual(
                main(
                    [
                        "comprehend",
                        "--name",
                        "deck",
                        "--input",
                        str(comp_path),
                        "--scope",
                        "project",
                    ]
                ),
                0,
            )
            out = tmp / "out.pptx"
            buf = io.StringIO()
            with redirect_stdout(buf):
                rc = main(
                    [
                        "generate",
                        "--name",
                        "deck",
                        "--input",
                        str(self._idoc_json(tmp)),
                        "--output",
                        str(out),
                        "--scope",
                        "project",
                        "--qa",
                        "fast",
                    ]
                )
            return rc, buf.getvalue(), out
        finally:
            os.chdir(old)

    def test_sanctioned_reconcile_surfaces_findings_in_cli_report(self) -> None:
        # A sanctioned reconcile (demo clear corroborated): the generator's INFO/
        # WARNING reconciliation findings must appear in the CLI-printed report and
        # the run must still PASS (rc 0). Before the wiring fix these were dropped.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            rc, stdout, out = self._run_cli_reconcile(
                tp, demo_slide_ref="region.slide.3"
            )
            self.assertEqual(rc, 0, stdout)
            self.assertTrue(out.is_file())
            # The agenda-regenerated INFO finding (emitted ONLY on the reconcile path)
            # reached the printed report - proof the findings out-param is wired.
            self.assertIn("agenda_regenerated", stdout)

    def test_unsanctioned_demo_clear_downgraded_surfaces_in_cli_report(self) -> None:
        # The model tags the agenda slide (authored "Intro\nBody", NOT a layout
        # prompt) demo: determinism does not corroborate -> KEEP + WARNING. The
        # WARNING must reach the CLI report (it is silently dropped without the fix).
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            rc, stdout, out = self._run_cli_reconcile(
                tp, demo_slide_ref="region.slide.1"
            )
            # No corroborated removal -> still passes, but the WARNING is surfaced.
            self.assertEqual(rc, 0, stdout)
            self.assertIn("demo_clear_downgraded", stdout)


class MI7SharedResolverTest(unittest.TestCase):
    """pptx generation routes layout resolution through the SHARED resolver spine."""

    def test_layout_resolution_uses_shared_resolver(self) -> None:
        from brandkit.profile.resolver import ProfileResolver

        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            prof = _extract_on_disk(template, tp)
            prs = Presentation(template)
            resolver = ProfileResolver(prof)
            cover = pg._layout_for_role(prs, resolver, "cover.title")
            content = pg._layout_for_role(prs, resolver, "heading.1")
            self.assertEqual(cover.name, "BrandCover")
            self.assertEqual(content.name, "BrandContent")

    def test_resolver_refuses_foreign_resolver_type(self) -> None:
        # A docx-style named_style smuggled into a pptx role must NOT resolve to a
        # layout (the shared spine refuses a type illegal for kind pptx).
        from brandkit.profile.resolver import ProfileResolver

        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            prof = _extract_on_disk(template, tp)
            prof["roles"]["cover.title"]["resolver"] = {
                "type": "named_style",
                "style_id": "Title",
            }
            prs = Presentation(template)
            resolver = ProfileResolver(prof)
            self.assertIsNone(pg._layout_for_role(prs, resolver, "cover.title"))


_COMPLEX_PPTX = (
    Path(__file__).resolve().parents[1]
    / "tests"
    / "fixtures"
    / "complex"
    / "acme_complex.pptx"
)


class PptxCheapFidelityTest(unittest.TestCase):
    """Regression coverage for the cheap PPTX fidelity fixes (P1/P4/P5/X7/Q14),
    grounded on the committed complex fixture (native table + chart + pictures +
    multi-level lists). Native table authoring is covered here; chart/SmartArt
    writers remain deferred and visible through degradation/survival warnings."""

    def _branded(self, td: Path) -> Path:
        template = td / "branded.pptx"
        _branded_template(template)
        return template

    def _gen(
        self, template: Path, idoc: ir.IntermediateDocument, out: Path, findings=None
    ):
        profile = _extract_profile(template)
        pg.generate(profile, template, idoc, out, findings=findings)
        return Presentation(out)

    # P1 - list items become REAL body paragraphs carrying paragraph.level --------
    def test_list_items_are_real_paragraphs_with_levels(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = self._branded(tp)
            out = tp / "out.pptx"
            idoc = ir.IntermediateDocument(
                blocks=[
                    ir.Heading(level=1, runs=[{"t": "Topics"}]),
                    ir.ListBlock(
                        items=[
                            ir.ListItem(
                                runs=[{"t": "Topic A"}],
                                level=0,
                                items=[
                                    ir.ListItem(runs=[{"t": "Sub A1"}], level=1),
                                ],
                            ),
                            ir.ListItem(runs=[{"t": "Topic B"}], level=0),
                        ]
                    ),
                ]
            )
            res = self._gen(template, idoc, out)
            topic_slide = next(
                s
                for s in res.slides
                if s.shapes.title
                and s.shapes.title.text
                and s.shapes.title.text.startswith("Topics")
            )
            body = pg._first_body_placeholder(topic_slide)
            paras = [(p.text, p.level) for p in body.text_frame.paragraphs]
            # Each item is its OWN paragraph at its real level - NOT a string-joined
            # "    • " blob in one paragraph.
            self.assertIn(("Topic A", 0), paras)
            self.assertIn(("Sub A1", 1), paras)
            self.assertIn(("Topic B", 0), paras)
            # No string-joined bullet glyph / indentation prefix leaked into text.
            for text, _ in paras:
                self.assertNotIn("•", text)
                self.assertFalse(text.startswith("    "))

    # X7 - generate-twice is byte-identical (pinned package modified time) ---------
    def test_generate_twice_is_byte_identical_on_complex_shell(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            prof = _extract_on_disk(_COMPLEX_PPTX, tp, name="acme")
            idoc = ir.IntermediateDocument(
                blocks=[
                    ir.Heading(level=1, runs=[{"t": "Update"}]),
                    ir.Paragraph(runs=[{"t": "Body text."}]),
                ],
                cover=ir.Cover(title=[{"t": "Cover"}]),
            )
            o1 = tp / "o1.pptx"
            o2 = tp / "o2.pptx"
            pg.generate(prof, _COMPLEX_PPTX, idoc, o1)
            pg.generate(prof, _COMPLEX_PPTX, idoc, o2)
            self.assertEqual(o1.read_bytes(), o2.read_bytes())

    # P5 - typed component inventory in extraction + survival baseline -------------
    def test_extractor_surfaces_typed_component_inventory(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            prof = _extract_on_disk(_COMPLEX_PPTX, tp, name="acme")
            slides = prof["artifact_catalog"]["slides"]
            self.assertTrue(all("components" in s for s in slides))
            totals = {"table": 0, "chart": 0, "picture": 0}
            for s in slides:
                for fam, n in s["components"].items():
                    totals[fam] += n
            # The fixture ships a native table, a native chart and pictures.
            self.assertGreaterEqual(totals["table"], 1)
            self.assertGreaterEqual(totals["chart"], 1)
            self.assertGreaterEqual(totals["picture"], 1)

    def test_inventory_components_matches_structure_helper(self) -> None:
        prs = Presentation(_COMPLEX_PPTX)
        totals = ps.inventory_components(prs)
        self.assertGreaterEqual(totals["table"], 1)
        self.assertGreaterEqual(totals["chart"], 1)
        self.assertGreaterEqual(totals["picture"], 1)

    # Q14 - table blocks author real native PPTX tables ---------------------------
    def test_table_block_authors_native_table_without_degradation(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = self._branded(tp)
            out = tp / "out.pptx"
            idoc = ir.IntermediateDocument(
                blocks=[
                    ir.Heading(level=1, runs=[{"t": "Data"}]),
                    ir.Table(
                        columns=[{"t": "Region"}, {"t": "Rev"}],
                        rows=[
                            [
                                ir.TableCell(runs=[{"t": "North"}]),
                                ir.TableCell(runs=[{"t": "100"}]),
                            ]
                        ],
                    ),
                ]
            )
            findings: list[Finding] = []
            prs = self._gen(template, idoc, out, findings=findings)
            tables = [
                shape.table
                for slide in prs.slides
                for shape in slide.shapes
                if getattr(shape, "has_table", False)
            ]
            self.assertEqual(len(tables), 1)
            self.assertEqual(tables[0].cell(0, 0).text, "Region")
            self.assertEqual(tables[0].cell(1, 0).text, "North")
            degraded = [f for f in findings if f.check == "block_degraded"]
            self.assertFalse(any("table" in f.message for f in degraded), findings)

    def test_generated_native_table_satisfies_table_component_survival(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            prof = _extract_on_disk(_COMPLEX_PPTX, tp, name="acme")
            idoc = ir.IntermediateDocument(
                blocks=[
                    ir.Heading(level=1, runs=[{"t": "Data"}]),
                    ir.Table(
                        columns=[{"t": "Region"}, {"t": "Rev"}],
                        rows=[
                            [
                                ir.TableCell(runs=[{"t": "North"}]),
                                ir.TableCell(runs=[{"t": "100"}]),
                            ]
                        ],
                    ),
                ]
            )
            out = tp / "out.pptx"
            sink: list[Finding] = []
            pg.generate(prof, _COMPLEX_PPTX, idoc, out, findings=sink)
            # Single source of truth: the GENERATOR no longer emits component_survival
            # (it used to, with drop-to-zero semantics, duplicating the QA gate).
            self.assertEqual(
                [f for f in sink if f.check == "component_survival"],
                [],
                "generator must not emit component_survival; the QA gate owns it",
            )
            # The authored native table survives, so the QA gate (the authority) reports
            # NO tables-survival drop (the unauthored chart/picture legitimately drop -
            # asserted by the sibling test below).
            survival = cd.check_component_survival(_COMPLEX_PPTX, out, prof)
            self.assertEqual([f for f in survival if "tables" in f.message], [])

    # CC-3(b) - native component lost from shell -> component_survival WARNING -----
    def test_component_survival_warns_when_native_table_lost(self) -> None:
        # The complex shell carries a native table/chart/picture; a deterministic
        # rebuild that emits only text down-renders them -> the QA gate (single
        # source of truth) WARNs for each lost family.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            prof = _extract_on_disk(_COMPLEX_PPTX, tp, name="acme")
            idoc = ir.IntermediateDocument(
                blocks=[
                    ir.Heading(level=1, runs=[{"t": "Text only"}]),
                    ir.Paragraph(runs=[{"t": "No native components here."}]),
                ]
            )
            out = tp / "out.pptx"
            pg.generate(prof, _COMPLEX_PPTX, idoc, out)
            findings = cd.check_component_survival(_COMPLEX_PPTX, out, prof)
            self.assertTrue(findings, "lost native components must be WARNed")
            self.assertTrue(all(f.check == "component_survival" for f in findings))
            self.assertTrue(all(f.severity == "WARNING" for f in findings))
            # The table that the shell carried is gone -> reported by family.
            self.assertTrue(any("tables" in f.message for f in findings))

    # P4 - filled cover placeholder keeps its run formatting (rPr re-assertion) ----
    def test_cover_fill_preserves_run_formatting(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = self._branded(tp)
            out = tp / "out.pptx"
            idoc = ir.IntermediateDocument(
                blocks=[ir.Heading(level=1, runs=[{"t": "S"}])],
                cover=ir.Cover(title=[{"t": "Branded Title"}], subtitle=[{"t": "Sub"}]),
            )
            res = self._gen(template, idoc, out)
            cover = next(s for s in res.slides if s.slide_layout.name == "BrandCover")
            self.assertEqual(cover.shapes.title.text, "Branded Title")
            # The title text frame has exactly one paragraph with a single run holding
            # the new text (rPr-preserving fill, not a clobbered text frame).
            tf = cover.shapes.title.text_frame
            self.assertEqual(len([p for p in tf.paragraphs]), 1)


class SilentDropFix(unittest.TestCase):
    """Component/Section/Toc/Divider must not vanish silently in pptx.

    ``_sections`` routes every non-heading / non-pagebreak block into a slide body,
    which feeds ``_body_lines``. Before the fix these four block types fell through
    with no output and no finding - a silent content drop, asymmetric with the docx
    leg (which expands Component/Section) and with the pptx kpi/chart/smartart/image
    blocks (which loudly ``_degrade``). Each now records a ``block_degraded``
    WARNING so the unrendered block is visible in QA, honoring the engine's
    "never drop content silently" invariant.
    """

    def test_unhandled_blocks_surface_as_block_degraded(self) -> None:
        blocks = [
            ir.Paragraph.from_dict({"type": "paragraph", "text": "kept"}),
            ir.Component.from_dict({"type": "component", "ref": "hero"}),
            ir.Section.from_dict({"type": "section", "ref": "intro"}),
            ir.Toc.from_dict({"type": "toc", "title": "Contents"}),
            ir.Divider.from_dict({"type": "divider"}),
        ]
        sink: list[Finding] = []
        lines = pg._body_lines(blocks, sink)
        # The real paragraph still renders.
        self.assertIn("kept", [line.text for line in lines])
        # All four unhandled types surfaced, none silently dropped.
        self.assertEqual(len(sink), 4)
        self.assertTrue(all(f.check == "block_degraded" for f in sink))
        self.assertTrue(all(f.severity == "WARNING" for f in sink))
        kinds = sorted(f.message.split("'")[1] for f in sink)
        self.assertEqual(kinds, ["component", "divider", "section", "toc"])


class ComponentExpansionSymmetryTest(unittest.TestCase):
    """Profile-defined component/section fragments expand the SAME on both IID legs.

    The docx leg already wired ``components.expand_components`` at the top of
    ``generate``; this proves the pptx leg now does too (TASK #8). A defined ref
    becomes its primitive sub-blocks in BOTH outputs (and the ``component``/
    ``section`` block does not survive to a writer); an undefined ref RAISES
    ``ComponentExpansionError`` on the pptx leg, symmetric with docx and fail-closed.

    Scope: STATIC, profile-defined expansion. Registry auto-POPULATION and ``slots``
    PARAMETERIZATION remain deferred milestones, intentionally not exercised here.
    """

    _COMPONENT_IDOC = {
        "blocks": [
            {"type": "component", "ref": "intro"},
            {"type": "heading", "level": 1, "text": "Real Section"},
            {"type": "paragraph", "text": "Authored body."},
        ]
    }

    @staticmethod
    def _component_registry() -> dict:
        return {
            "intro": {
                "blocks": [
                    {"type": "heading", "level": 1, "text": "Intro Heading"},
                    {"type": "paragraph", "text": "Intro body from the fragment."},
                ]
            }
        }

    def test_pptx_expands_defined_component(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            out = tp / "out.pptx"
            profile = _extract_profile(template)
            profile["components"] = self._component_registry()
            idoc = parse_idoc(self._COMPONENT_IDOC)
            pg.generate(profile, template, idoc, out)
            text = _all_text(Presentation(out))
            # The fragment's PRIMITIVE content is present in the deck...
            self.assertIn("Intro Heading", text)
            self.assertIn("Intro body from the fragment", text)
            # ...alongside the authored content, and nothing was dropped silently.
            self.assertIn("Real Section", text)
            self.assertIn("Authored body", text)

    def test_docx_expands_defined_component(self) -> None:
        """The SAME profile + idoc expands on the docx leg (cross-format symmetry)."""
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = tp / "shell.docx"
            Document().save(shell)
            out = tp / "out.docx"
            profile = schema.build_envelope("docx", {"name": "t"})
            profile["surface"] = {"docx": {}}
            profile["roles"] = {"_index": []}
            profile["components"] = self._component_registry()
            idoc = parse_idoc(self._COMPONENT_IDOC)
            docx_generate.generate(profile, shell, idoc, out)
            text = "\n".join(p.text for p in Document(out).paragraphs)
            self.assertIn("Intro Heading", text)
            self.assertIn("Intro body from the fragment", text)
            self.assertIn("Real Section", text)
            self.assertIn("Authored body", text)

    def test_component_block_does_not_survive_expansion(self) -> None:
        """After expansion the document carries primitives only - no ``component``."""
        profile = schema.build_envelope("pptx", {"name": "deck"})
        profile["components"] = self._component_registry()
        idoc = parse_idoc(self._COMPONENT_IDOC)
        expanded = ir_components.expand_components(idoc, profile)
        types = [b.TYPE for b in expanded.blocks]
        self.assertNotIn("component", types)
        self.assertEqual(types, ["heading", "paragraph", "heading", "paragraph"], types)

    def test_pptx_undefined_ref_raises(self) -> None:
        """An undefined component ref RAISES on the pptx leg (symmetric with docx)."""
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            out = tp / "out.pptx"
            profile = _extract_profile(template)  # empty components registry
            idoc = parse_idoc({"blocks": [{"type": "component", "ref": "ghost"}]})
            with self.assertRaises(ir_components.ComponentExpansionError):
                pg.generate(profile, template, idoc, out)

    def test_pptx_undefined_section_ref_raises(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            template = tp / "branded.pptx"
            _branded_template(template)
            out = tp / "out.pptx"
            profile = _extract_profile(template)  # empty sections registry
            idoc = parse_idoc({"blocks": [{"type": "section", "ref": "ghost"}]})
            with self.assertRaises(ir_components.ComponentExpansionError):
                pg.generate(profile, template, idoc, out)


class FragmentSlotSubstitutionTest(unittest.TestCase):
    """A referencing block's ``slots`` fill ``{{name}}`` tokens in the fragment
    template at expansion time, without mutating the shared profile registry.

    Slot parameterization (TASK: fragments) is now wired: ``expand_components``
    substitutes ``{{name}}`` -> ``slots[name]`` at the dict level (a deep copy)
    before parsing, so it works for every text-bearing field. An unfilled token
    resolves to the empty string and is never leaked as literal ``{{...}}``.
    """

    @staticmethod
    def _profile_with_slotted_fragment() -> dict:
        prof = schema.build_envelope("pptx", {"name": "deck"})
        prof["components"] = {
            "note": {
                "blocks": [
                    {"type": "heading", "level": 2, "runs": [{"t": "{{title}}"}]},
                    {"type": "paragraph", "runs": [{"t": "Body: {{body}}"}]},
                ]
            }
        }
        return prof

    @staticmethod
    def _runs_text(expanded) -> str:
        return " ".join(
            r.get("t", "")
            for b in expanded.blocks
            for r in (getattr(b, "runs", []) or [])
        )

    def test_slots_substituted_into_runs(self) -> None:
        prof = self._profile_with_slotted_fragment()
        idoc = parse_idoc(
            {
                "blocks": [
                    {
                        "type": "component",
                        "ref": "note",
                        "slots": {"title": "Hello", "body": "World"},
                    }
                ]
            }
        )
        joined = self._runs_text(ir_components.expand_components(idoc, prof))
        self.assertIn("Hello", joined)
        self.assertIn("Body: World", joined)
        self.assertNotIn("{{title}}", joined)
        self.assertNotIn("{{body}}", joined)

    def test_unfilled_slot_resolves_empty_not_leaked(self) -> None:
        prof = self._profile_with_slotted_fragment()
        idoc = parse_idoc(
            {"blocks": [{"type": "component", "ref": "note", "slots": {"title": "X"}}]}
        )
        joined = self._runs_text(ir_components.expand_components(idoc, prof))
        self.assertNotIn("{{body}}", joined)  # token not leaked into output
        self.assertIn("Body:", joined)  # surrounding literal text preserved

    def test_no_slots_reference_strips_tokens(self) -> None:
        # A reference with NO slots must still strip tokens (-> "") rather than leak
        # the literal {{name}} into the output.
        prof = self._profile_with_slotted_fragment()
        idoc = parse_idoc({"blocks": [{"type": "component", "ref": "note"}]})
        joined = self._runs_text(ir_components.expand_components(idoc, prof))
        self.assertNotIn("{{title}}", joined)
        self.assertNotIn("{{body}}", joined)
        self.assertIn("Body:", joined)  # surrounding literal text preserved

    def test_none_slot_value_becomes_empty_not_none_string(self) -> None:
        prof = self._profile_with_slotted_fragment()
        idoc = parse_idoc(
            {
                "blocks": [
                    {
                        "type": "component",
                        "ref": "note",
                        "slots": {"title": None, "body": "ok"},
                    }
                ]
            }
        )
        joined = self._runs_text(ir_components.expand_components(idoc, prof))
        self.assertNotIn("None", joined)  # a null slot -> "" not the string "None"
        self.assertIn("Body: ok", joined)

    def test_registry_not_mutated_by_slot_fill(self) -> None:
        prof = self._profile_with_slotted_fragment()
        idoc = parse_idoc(
            {
                "blocks": [
                    {
                        "type": "component",
                        "ref": "note",
                        "slots": {"title": "Z", "body": "Q"},
                    }
                ]
            }
        )
        ir_components.expand_components(idoc, prof)
        # The template still carries the unfilled tokens: substitution deep-copies,
        # it never mutates the profile registry the fragment lives in.
        self.assertEqual(
            prof["components"]["note"]["blocks"][0]["runs"][0]["t"], "{{title}}"
        )


class FragmentRegistryValidationTest(unittest.TestCase):
    """``schema.validate`` reports a malformed components/sections registry entry.

    Well-formedness only: each entry must be a dict carrying a ``blocks`` list. A
    malformed entry is surfaced (fail-closed) rather than blowing up later inside
    the expander. Block CONTENTS are not re-checked here (``block_from_dict`` owns
    that contract).
    """

    def test_component_entry_missing_blocks_is_reported(self) -> None:
        prof = schema.build_envelope("pptx", {"name": "deck"})
        prof["components"] = {"intro": {"description": "no blocks here"}}
        problems = schema.validate(prof)
        self.assertTrue(any("components.intro.blocks" in p for p in problems), problems)

    def test_section_entry_not_an_object_is_reported(self) -> None:
        prof = schema.build_envelope("pptx", {"name": "deck"})
        prof["sections"] = {"hero": ["not", "a", "dict"]}
        problems = schema.validate(prof)
        self.assertTrue(any("sections.hero" in p for p in problems), problems)

    def test_registry_not_a_map_is_reported(self) -> None:
        prof = schema.build_envelope("pptx", {"name": "deck"})
        prof["components"] = ["not", "a", "map"]
        problems = schema.validate(prof)
        self.assertTrue(any(p.startswith("components:") for p in problems), problems)

    def test_well_formed_registry_validates_clean(self) -> None:
        prof = schema.build_envelope("pptx", {"name": "deck"})
        prof["components"] = {
            "intro": {"blocks": [{"type": "paragraph", "text": "ok"}]}
        }
        prof["sections"] = {"hero": {"blocks": []}}
        problems = schema.validate(prof)
        self.assertFalse(
            any(p.startswith(("components", "sections")) for p in problems), problems
        )


class AgendaDetectionHonorsBodyIdx(unittest.TestCase):
    """On a multi-body layout, agenda DETECTION reads the same profile-resolved
    ``ph_idx`` body the refresh write path targets - so detect and rewrite never
    point at different placeholders. Before the fix _existing_agenda_slide read the
    positional first body while _regenerate_agenda wrote into the named idx.
    """

    def test_existing_agenda_detected_via_named_idx_not_first_body(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            template = Path(td) / "deck.pptx"
            _multibody_template(template)
            prs = Presentation(template)
            two_body = next(
                lay for lay in prs.slide_layouts if lay.name == "BrandTwoBody"
            )
            # Slide 0: a cover-ish slide so the section sldIdLst has ids to point at.
            prs.slides.add_slide(prs.slide_layouts[0])
            # Slide 1: the agenda. Its section list lives in the SECOND body (idx 2),
            # NOT the first (idx 1), which carries unrelated text.
            agenda = prs.slides.add_slide(two_body)
            agenda.shapes.title.text = "Sommario"
            agenda.placeholders[1].text = "unrelated first-body text"
            agenda.placeholders[2].text = "Alpha\nBeta"
            _add_section_list(prs, sections=[("Alpha", [0]), ("Beta", [1])])
            prs.save(template)

            prs = Presentation(template)
            # Detect honoring the NAMED idx (2) finds the agenda (its idx-2 body
            # lists exactly the section names).
            self.assertIsNotNone(pg._existing_agenda_slide(prs, 2))
            # Reading the positional first body (idx 1) would NOT match (it carries
            # unrelated text), so the old positional detection missed this agenda.
            self.assertIsNone(pg._existing_agenda_slide(prs, 1))


# ---------------------------------------------------------------------------
# Native PPTX charts (ir.Chart -> real graphicFrame/c:chart, theme-colored)
# ---------------------------------------------------------------------------
class NativeChartTest(unittest.TestCase):
    """A Chart block is authored as a REAL PowerPoint chart (not flattened to body
    text): correct type/series/categories/title, on-brand by theme inheritance,
    byte-idempotent (the embedded data workbook's wall-clock timestamps are
    normalized), and graceful on empty/unknown input - never a crash or silent drop.
    """

    _IDOC = {
        "blocks": [
            {"type": "heading", "level": 1, "runs": [{"t": "Ricavi"}]},
            {
                "type": "chart",
                "chart_type": "bar",
                "title": "Ricavi (M)",
                "categories": ["Q1", "Q2", "Q3"],
                "series": [
                    {"name": "A", "values": [1, 2, 3]},
                    {"name": "B", "values": [3, 2, 1]},
                ],
            },
            {"type": "heading", "level": 1, "runs": [{"t": "Quota"}]},
            {
                "type": "chart",
                "chart_type": "pie",
                "title": "Quota",
                "categories": ["X", "Y"],
                "series": [{"name": "S", "values": [60, 40]}],
            },
        ]
    }

    def _charts(self, prs):
        return [sh.chart for s in prs.slides for sh in s.shapes if sh.has_chart]

    def test_chart_blocks_become_native_charts(self):
        from pptx.enum.chart import XL_CHART_TYPE

        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            template = tmp / "t.pptx"
            _branded_template(template)
            profile = _extract_profile(template)
            out = tmp / "out.pptx"
            sink: list[Finding] = []
            pg.generate(profile, template, parse_idoc(self._IDOC), out, findings=sink)

            charts = self._charts(Presentation(out))
            self.assertEqual(len(charts), 2, "both chart blocks should be native")
            kinds = {c.chart_type for c in charts}
            self.assertIn(XL_CHART_TYPE.COLUMN_CLUSTERED, kinds)  # "bar" -> columns
            self.assertIn(XL_CHART_TYPE.PIE, kinds)
            col = next(
                c for c in charts if c.chart_type == XL_CHART_TYPE.COLUMN_CLUSTERED
            )
            self.assertEqual(len(col.plots[0].series), 2)
            self.assertEqual(list(col.plots[0].categories), ["Q1", "Q2", "Q3"])
            self.assertEqual(col.chart_title.text_frame.text, "Ricavi (M)")
            # Series NAMES and VALUES round-trip intact (a writer that reversed,
            # duplicated or garbled the data would be caught here, not just by count).
            s0, s1 = col.plots[0].series[0], col.plots[0].series[1]
            self.assertEqual((s0.name, tuple(s0.values)), ("A", (1.0, 2.0, 3.0)))
            self.assertEqual((s1.name, tuple(s1.values)), ("B", (3.0, 2.0, 1.0)))
            # On-brand by construction: no explicit series fill is set, so the chart
            # inherits the deck theme's accent colors (a regression injecting a literal
            # color would give the fill a concrete .type instead of None).
            self.assertIsNone(s0.format.fill.type, "series fill must inherit the theme")
            # Chart is NOT degraded (it is native, not flattened to text).
            self.assertFalse(
                any(f.check == "block_degraded" and "chart" in f.message for f in sink)
            )

    def _embedded_workbook_cores(self, path: Path) -> list[str]:
        """docProps/core.xml of every embedded ``.xlsx`` workbook in the deck."""
        import io
        import zipfile

        cores: list[str] = []
        with zipfile.ZipFile(path) as outer:
            for name in outer.namelist():
                if "embeddings" in name and name.endswith(".xlsx"):
                    with zipfile.ZipFile(io.BytesIO(outer.read(name))) as inner:
                        cores.append(inner.read("docProps/core.xml").decode("utf-8"))
        return cores

    def test_chart_generation_is_byte_idempotent(self):
        # python-pptx embeds a data workbook whose core.xml carries WALL-CLOCK
        # dcterms timestamps; repack_fixed_timestamps must normalize the nested
        # package so two generations are byte-identical.
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            template = tmp / "t.pptx"
            _branded_template(template)
            profile = _extract_profile(template)
            a, b = tmp / "a.pptx", tmp / "b.pptx"
            pg.generate(profile, template, parse_idoc(self._IDOC), a)
            pg.generate(profile, template, parse_idoc(self._IDOC), b)
            self.assertEqual(
                a.read_bytes(),
                b.read_bytes(),
                "native-chart deck is not byte-idempotent",
            )
            # The chart actually produced an embedded workbook (so the test exercises
            # the nested-package path, not a no-op)...
            cores = self._embedded_workbook_cores(a)
            self.assertTrue(cores, "no embedded chart workbook was produced")
            self.assertTrue(self._charts(Presentation(a)), "no native chart authored")
            # ...and its wall-clock dcterms timestamps were PINNED to the fixed epoch.
            # This proves the fix RAN (the two in-process builds alone could coincide
            # within one wall-clock second and pass even without normalization).
            for core in cores:
                self.assertIn("1980-01-01T00:00:00Z", core)
                self.assertNotIn("2026", core)  # no surviving wall-clock year

    def test_unknown_type_falls_back_and_empty_chart_degrades(self):
        idoc = {
            "blocks": [
                {"type": "heading", "level": 1, "runs": [{"t": "Odd"}]},
                {
                    "type": "chart",
                    "chart_type": "nonsense",
                    "categories": ["A", "B"],
                    "series": [{"name": "S", "values": [1, 2]}],
                },
                {"type": "heading", "level": 1, "runs": [{"t": "Empty"}]},
                {"type": "chart", "chart_type": "bar", "categories": [], "series": []},
            ]
        }
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            template = tmp / "t.pptx"
            _branded_template(template)
            profile = _extract_profile(template)
            out = tmp / "out.pptx"
            sink: list[Finding] = []
            pg.generate(profile, template, parse_idoc(idoc), out, findings=sink)
            # The unknown type still renders a (fallback) native chart, and the
            # fallback is specifically a clustered column chart (not just "some" chart).
            from pptx.enum.chart import XL_CHART_TYPE

            rendered = self._charts(Presentation(out))
            self.assertEqual(len(rendered), 1)
            self.assertEqual(rendered[0].chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
            self.assertTrue(
                any(f.check == "chart_type_fallback" for f in sink),
                "unknown chart_type should surface an INFO fallback, not be silent",
            )
            # ...and the empty chart degrades loudly (never a silent drop).
            self.assertTrue(
                any(f.check == "block_degraded" and "chart" in f.message for f in sink)
            )

    def test_multi_series_pie_warns_and_keeps_first_series(self):
        # A pie renders only its first series; the dropped series are surfaced as a
        # WARNING (data loss is visible), and exactly one series is plotted.
        idoc = {
            "blocks": [
                {"type": "heading", "level": 1, "runs": [{"t": "Pie"}]},
                {
                    "type": "chart",
                    "chart_type": "pie",
                    "categories": ["A", "B"],
                    "series": [
                        {"name": "First", "values": [1, 2]},
                        {"name": "Second", "values": [3, 4]},
                    ],
                },
            ]
        }
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            template = tmp / "t.pptx"
            _branded_template(template)
            profile = _extract_profile(template)
            out = tmp / "out.pptx"
            sink: list[Finding] = []
            pg.generate(profile, template, parse_idoc(idoc), out, findings=sink)
            self.assertTrue(
                any(f.check == "chart_series_truncated" for f in sink),
                "a multi-series pie must surface the dropped series, not hide them",
            )
            charts = self._charts(Presentation(out))
            self.assertEqual(len(charts), 1)
            self.assertEqual(len(charts[0].plots[0].series), 1)


# ---------------------------------------------------------------------------
# Native PPTX SmartArt (ir.SmartArt -> brand-themed autoshapes, not flattened)
# ---------------------------------------------------------------------------
class NativeSmartArtTest(unittest.TestCase):
    """A SmartArt block is authored as REAL brand-themed autoshapes (a chevron row
    for a process, a stacked box list otherwise), not flattened to body text:
    correct shape count, node text + children preserved, byte-idempotent, empty
    degrades loudly."""

    _IDOC = {
        "blocks": [
            {"type": "heading", "level": 1, "runs": [{"t": "Flow"}]},
            {
                "type": "smartart",
                "diagram": "process",
                "nodes": [{"text": "Plan"}, {"text": "Build"}, {"text": "Ship"}],
            },
            {"type": "heading", "level": 1, "runs": [{"t": "Pillars"}]},
            {
                "type": "smartart",
                "diagram": "list",
                "nodes": [
                    {"text": "Quality", "children": [{"text": "tests"}]},
                    {"text": "Speed"},
                ],
            },
        ]
    }

    def _autoshapes(self, prs):
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        return [
            sh
            for s in prs.slides
            for sh in s.shapes
            if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
        ]

    def test_smartart_becomes_native_shapes(self):
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            template = tmp / "t.pptx"
            _branded_template(template)
            profile = _extract_profile(template)
            out = tmp / "out.pptx"
            sink: list[Finding] = []
            pg.generate(profile, template, parse_idoc(self._IDOC), out, findings=sink)
            shapes = self._autoshapes(Presentation(out))
            # 3 process chevrons + 2 list boxes = 5 native shapes.
            self.assertEqual(len(shapes), 5)
            texts = "\n".join(sh.text_frame.text for sh in shapes)
            self.assertIn("Plan", texts)
            self.assertIn("Quality", texts)
            self.assertIn("tests", texts)  # a child is preserved as a sub-line
            self.assertFalse(
                any(
                    f.check == "block_degraded" and "smartart" in f.message
                    for f in sink
                )
            )

    def test_smartart_is_byte_idempotent(self):
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            template = tmp / "t.pptx"
            _branded_template(template)
            profile = _extract_profile(template)
            a, b = tmp / "a.pptx", tmp / "b.pptx"
            pg.generate(profile, template, parse_idoc(self._IDOC), a)
            pg.generate(profile, template, parse_idoc(self._IDOC), b)
            self.assertEqual(a.read_bytes(), b.read_bytes())

    def test_empty_smartart_degrades(self):
        idoc = {"blocks": [{"type": "smartart", "diagram": "process", "nodes": []}]}
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            template = tmp / "t.pptx"
            _branded_template(template)
            profile = _extract_profile(template)
            out = tmp / "out.pptx"
            sink: list[Finding] = []
            pg.generate(profile, template, parse_idoc(idoc), out, findings=sink)
            self.assertEqual(len(self._autoshapes(Presentation(out))), 0)
            self.assertTrue(
                any(
                    f.check == "block_degraded" and "smartart" in f.message
                    for f in sink
                )
            )


# ---------------------------------------------------------------------------
# Native PPTX table cell merge (colspan/rowspan parity with the docx writer)
# ---------------------------------------------------------------------------
class PptxTableMergeTest(unittest.TestCase):
    """A table cell with colspan/rowspan merges the spanned grid cells in the native
    PowerPoint table (it used to render as a full ungrouped grid, unlike docx)."""

    def test_colspan_banner_merges_cells(self):
        idoc = {
            "blocks": [
                {"type": "heading", "level": 1, "runs": [{"t": "T"}]},
                {
                    "type": "table",
                    "columns": ["A", "B"],
                    "rows": [
                        [{"runs": [{"t": "Banner"}], "colspan": 2}],
                        [{"runs": [{"t": "x"}]}, {"runs": [{"t": "y"}]}],
                    ],
                },
            ]
        }
        with tempfile.TemporaryDirectory() as td:
            tmp = Path(td)
            template = tmp / "t.pptx"
            _branded_template(template)
            profile = _extract_profile(template)
            out = tmp / "out.pptx"
            pg.generate(profile, template, parse_idoc(idoc), out)
            tbl = next(
                sh.table
                for s in Presentation(out).slides
                for sh in s.shapes
                if sh.has_table
            )
            # header row 0 (A|B); the banner is grid row 1, col 0, spanning 2 columns.
            origin = tbl.cell(1, 0)
            self.assertTrue(origin.is_merge_origin, "banner cell did not merge")
            self.assertEqual(origin.span_width, 2)
            self.assertEqual(origin.text, "Banner")
            # the normal row below keeps two distinct cells
            self.assertFalse(tbl.cell(2, 0).is_merge_origin)
            self.assertEqual(tbl.cell(2, 1).text, "y")


class FragmentFanoutBudgetTest(unittest.TestCase):
    """Runaway component/section FAN-OUT (within the depth cap) fails LOUD via the
    node-count budget instead of hanging / OOM-ing the generator."""

    def test_runaway_fanout_raises_component_expansion_error(self) -> None:
        profile = schema.build_envelope("pptx", {"name": "deck"})
        # A leaf of 200 primitives, referenced 300x -> 60_000 primitives > the cap.
        profile["components"] = {
            "leaf": {"blocks": [{"type": "paragraph", "text": "x"}] * 200}
        }
        idoc = parse_idoc({"blocks": [{"type": "component", "ref": "leaf"}] * 300})
        with self.assertRaises(ir_components.ComponentExpansionError):
            ir_components.expand_components(idoc, profile)

    def test_normal_fanout_within_budget_expands(self) -> None:
        profile = schema.build_envelope("pptx", {"name": "deck"})
        profile["components"] = {
            "leaf": {"blocks": [{"type": "paragraph", "text": "x"}] * 5}
        }
        idoc = parse_idoc({"blocks": [{"type": "component", "ref": "leaf"}] * 5})
        expanded = ir_components.expand_components(idoc, profile)
        self.assertEqual(len(expanded.blocks), 25)  # 5 x 5, well under the cap


if __name__ == "__main__":  # pragma: no cover
    unittest.main()
