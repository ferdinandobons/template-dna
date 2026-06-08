# SPDX-License-Identifier: MIT
"""End-to-end "kitchen sink" regression against the committed example templates.

Extracts each examples/templates/branddocs_template.{docx,pptx,xlsx}, generates from
a content document that exercises EVERY IR block type plus inline rich-run emphasis
(bold/italic/underline/hyperlink), and runs the QA gate. This is both the loop that
tests the skill against the showcase templates and a permanent regression net:

  * generation never raises and QA never FAILS (degradations are warnings, not errors);
  * inline emphasis + hyperlinks SURVIVE on docx (they used to be silently flattened);
  * `divider` is now a native docx artifact (no longer a block_degraded warning);
  * `chart` is a native PowerPoint chart on the pptx vertical (docx chart is still
    deferred); the still-deferred native writers degrade LOUDLY (a visible
    block_degraded finding), never silently.

The templates are 100% synthetic BrandDocs showcases; no proprietary file is used.
"""

from __future__ import annotations

import hashlib
import os
import sys
import tempfile
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

from brandkit.formats.docx import extract as docx_extract
from brandkit.formats.docx import generate as docx_generate
from brandkit.formats.pptx import extract as pptx_extract
from brandkit.formats.pptx import generate as pptx_generate
from brandkit.formats.xlsx import extract as xlsx_extract
from brandkit.formats.xlsx import generate as xlsx_generate
from brandkit.grid.model import GridDocument
from brandkit.ir.model import parse_idoc
from brandkit.profile import store
from brandkit.qa.gate import run_qa

TEMPLATES = Path(__file__).resolve().parents[1] / "examples" / "templates"

_RICH = [
    {"t": "Plain "},
    {"t": "bold", "b": True},
    {"t": " "},
    {"t": "italic", "i": True},
    {"t": " "},
    {"t": "underline", "u": True},
    {"t": " and a "},
    {"t": "link", "link": "https://example.com"},
    {"t": "."},
]

# Content exercising every block type. component/section are omitted (they require a
# profile-defined fragment registry and otherwise raise fail-closed by design).
KITCHEN = {
    "cover": {
        "title": "Kitchen Sink Report",
        "subtitle": "Every component",
        "fields": {
            "doc_id": "KS-2026-001",
            "date": "2026-06-08",
            "author": "BrandDocs",
        },
    },
    "blocks": [
        {
            "type": "heading",
            "level": 1,
            "runs": [{"t": "Executive "}, {"t": "Summary", "b": True}],
        },
        {"type": "paragraph", "runs": _RICH},
        {"type": "heading", "level": 2, "text": "Details"},
        {
            "type": "list",
            "ordered": False,
            "items": [
                {
                    "runs": [{"t": "First ", "b": True}, {"t": "bullet"}],
                    "level": 0,
                    "items": [{"text": "nested", "level": 1}],
                },
                {"text": "Second"},
            ],
        },
        {
            "type": "list",
            "ordered": True,
            "items": [{"text": "Step one"}, {"text": "Step two"}],
        },
        {
            "type": "callout",
            "intent": "info",
            "title": "Note",
            "runs": [{"t": "Info with "}, {"t": "emphasis", "i": True}],
        },
        {"type": "callout", "intent": "warning", "text": "Careful"},
        {
            "type": "table",
            "columns": ["Area", "Status"],
            "rows": [["Pipeline", "Healthy"], ["Delivery", "Green"]],
            "caption": "Status table",
        },
        {
            "type": "quote",
            "runs": [{"t": "A memorable quote."}],
            "attribution": "Jane Doe",
        },
        {"type": "caption", "target": "figure", "text": "Figure caption"},
        {
            "type": "kpi",
            "items": [
                {"label": "Revenue", "value": "$1.2M", "delta": "+8%"},
                {"label": "Churn", "value": "2.1%", "delta": "-0.3%"},
            ],
        },
        {
            "type": "chart",
            "chart_type": "bar",
            "title": "Revenue",
            "categories": ["Q1", "Q2"],
            "series": [{"name": "Net", "values": [10, 12]}],
        },
        {
            "type": "smartart",
            "diagram": "process",
            "nodes": [{"text": "Plan"}, {"text": "Build"}, {"text": "Ship"}],
        },
        {
            "type": "image",
            "src": "missing.png",
            "alt": "A figure",
            "caption": "An image",
        },
        {"type": "divider"},
        {"type": "toc", "title": "Contents", "max_level": 3},
        {"type": "pagebreak"},
        {"type": "heading", "level": 1, "text": "Appendix"},
        {"type": "paragraph", "text": "End."},
    ],
}


def _degraded_kinds(findings) -> set[str]:
    kinds = set()
    for f in findings:
        if f.check == "block_degraded":
            # message form: "'<kind>' block ..."
            parts = f.message.split("'")
            if len(parts) > 1:
                kinds.add(parts[1])
    return kinds


class _Base(unittest.TestCase):
    KIND = ""
    EXTRACT = None
    GENERATE = None

    @classmethod
    def setUpClass(cls):
        cls.template = TEMPLATES / f"branddocs_template.{cls.KIND}"
        if not cls.template.exists():
            raise unittest.SkipTest(f"missing example template {cls.template}")

    def _extract(self, td: Path):
        old = os.getcwd()
        os.chdir(td)
        try:
            self.EXTRACT.extract(self.template, "ks", scope="project", cwd=td)
            return store.load_profile("ks", "project")
        finally:
            os.chdir(old)


class DocxKitchenSink(_Base):
    KIND = "docx"
    EXTRACT = docx_extract
    GENERATE = docx_generate

    def test_kitchen_sink_generates_with_inline_and_native_divider(self):
        import copy

        from PIL import Image as PILImage

        with tempfile.TemporaryDirectory() as t:
            td = Path(t)
            loaded = self._extract(td)
            # A real image source so the native image writer places it (the shared
            # KITCHEN points at a missing file to exercise graceful degradation).
            png = td / "fig.png"
            PILImage.new("RGB", (40, 20), (10, 30, 60)).save(png)
            data = copy.deepcopy(KITCHEN)
            for b in data["blocks"]:
                if b.get("type") == "image":
                    b["src"] = str(png)

            out = td / "out.docx"
            sink = []
            docx_generate.generate(
                loaded.profile, loaded.shell_path, parse_idoc(data), out, findings=sink
            )
            self.assertTrue(out.is_file())
            report = run_qa(
                out,
                loaded.profile,
                shell=loaded.shell_path,
                extra_findings=list(sink),
                qa="fast",
            )
            self.assertNotEqual(
                report.verdict,
                "failed",
                [f.message for f in report.findings if f.severity == "ERROR"],
            )

            from docx import Document
            import lxml.etree as ET

            doc = Document(out)
            runs = [r for p in doc.paragraphs for r in p.runs]
            # Inline emphasis survives (was silently flattened before the run-aware writer).
            self.assertTrue(any(r.bold for r in runs), "bold run lost")
            self.assertTrue(any(r.italic for r in runs), "italic run lost")
            self.assertTrue(any(r.underline for r in runs), "underline run lost")
            body_xml = ET.tostring(doc.element.body).decode().lower()
            self.assertIn("hyperlink", body_xml, "hyperlink lost")
            # Divider is native (a paragraph border), image is native (a:blip drawing),
            # KPI is native (a brand table), chart is native (a c:chart part) - none
            # should be a degraded warning.
            self.assertIn("pbdr", body_xml, "native divider rule missing")
            self.assertIn("a:blip", body_xml, "native image not placed")
            self.assertIn("c:chart", body_xml, "native chart drawing not placed")
            degraded = _degraded_kinds(sink)
            self.assertEqual(
                degraded & {"divider", "image", "kpi", "chart"},
                set(),
                f"these should be native now: {degraded}",
            )
            # Genuinely-deferred native writers still degrade loudly (never silently).
            self.assertTrue({"smartart"} <= degraded)
            # Output fidelity on the SHOWCASE profile (different brand style names than
            # the synthetic fixture): lists carry real numbering and the table carries
            # the brand table style, so role nomination is exercised end-to-end.
            self.assertIn("numpr", body_xml, "list numbering (w:numPr) missing")
            self.assertIn("tblstyle", body_xml, "brand table style missing")

    def test_generate_is_byte_idempotent(self):
        with tempfile.TemporaryDirectory() as t:
            td = Path(t)
            loaded = self._extract(td)
            a, b = td / "a.docx", td / "b.docx"
            docx_generate.generate(
                loaded.profile, loaded.shell_path, parse_idoc(KITCHEN), a
            )
            docx_generate.generate(
                loaded.profile, loaded.shell_path, parse_idoc(KITCHEN), b
            )
            self.assertEqual(
                hashlib.sha256(a.read_bytes()).hexdigest(),
                hashlib.sha256(b.read_bytes()).hexdigest(),
            )


class PptxKitchenSink(_Base):
    KIND = "pptx"
    EXTRACT = pptx_extract
    GENERATE = pptx_generate

    def test_kitchen_sink_generates_native_kpi_image_and_degrades_loudly(self):
        import copy

        from PIL import Image as PILImage
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as t:
            td = Path(t)
            loaded = self._extract(td)
            png = td / "fig.png"
            PILImage.new("RGB", (40, 20), (10, 30, 60)).save(png)
            data = copy.deepcopy(KITCHEN)
            for b in data["blocks"]:
                if b.get("type") == "image":
                    b["src"] = str(png)
            out = td / "out.pptx"
            sink = []
            pptx_generate.generate(
                loaded.profile, loaded.shell_path, parse_idoc(data), out, findings=sink
            )
            self.assertTrue(out.is_file())
            report = run_qa(
                out,
                loaded.profile,
                shell=loaded.shell_path,
                extra_findings=list(sink),
                qa="fast",
            )
            self.assertNotEqual(
                report.verdict,
                "failed",
                [f.message for f in report.findings if f.severity == "ERROR"],
            )
            degraded = _degraded_kinds(sink)
            # KPI (native table), Image (native picture from a real src) and Chart
            # (native PowerPoint chart) no longer degrade; only smartart still has no
            # native writer and degrades loudly.
            self.assertEqual(
                degraded & {"kpi", "image", "chart"},
                set(),
                f"kpi/image/chart should be native now: {degraded}",
            )
            self.assertTrue({"smartart"} <= degraded)
            # A picture shape, a (KPI) table shape and a chart shape are authored.
            prs = Presentation(out)
            has_pic = any(
                sh.shape_type == 13
                for s in prs.slides
                for sh in s.shapes  # PICTURE
            )
            has_tbl = any(sh.has_table for s in prs.slides for sh in s.shapes)
            has_chart = any(sh.has_chart for s in prs.slides for sh in s.shapes)
            self.assertTrue(has_pic, "native picture not placed")
            self.assertTrue(has_tbl, "native KPI/table shape not placed")
            self.assertTrue(has_chart, "native chart shape not placed")

    def test_reconcile_path_no_duplicate_parts(self):
        # Reconcile/comprehension path against the SHOWCASE deck: clearing a demo
        # slide that is NOT the highest-indexed one must not leave an orphaned slide
        # part that the next add_slide collides with (duplicate ZIP part name -> a
        # corrupt OPC package PowerPoint would repair). Regression for the bug the
        # deterministic-only kitchen-sink and the highest-index reconcile fixture
        # both missed.
        import zipfile
        from collections import Counter

        from brandkit.formats.pptx import structure as ps
        from brandkit.profile import schema, store

        with tempfile.TemporaryDirectory() as t:
            td = Path(t)
            loaded = self._extract(td)
            prof = loaded.profile
            from pptx import Presentation

            demo = sorted(ps.demo_slide_indices(Presentation(loaded.shell_path)))
            if not demo:
                self.skipTest("example deck has no detectable demo slide")
            block = schema.empty_comprehension()
            block["status"] = "present"
            block["source_shell_sha256"] = prof["provenance"]["shell"]["sha256"]
            block["confidence"] = 0.9
            block["demo_classification"] = {
                "regions": [
                    {"region_ref": f"region.slide.{demo[0]}", "verdict": "demo"}
                ]
            }
            prof["comprehension"] = block
            self.assertTrue(store.comprehension_is_present(prof))

            out = td / "out.pptx"
            sink = []
            pptx_generate.generate(
                prof,
                loaded.shell_path,
                parse_idoc(
                    {
                        "cover": {"title": "T"},
                        "blocks": [
                            {"type": "heading", "level": 1, "text": "Intro"},
                            {"type": "paragraph", "text": "body"},
                        ],
                    }
                ),
                out,
                findings=sink,
            )
            names = zipfile.ZipFile(out).namelist()
            dups = [n for n, c in Counter(names).items() if c > 1]
            self.assertEqual(dups, [], f"duplicate package parts: {dups}")
            Presentation(out)  # reopens cleanly
            report = run_qa(
                out, prof, shell=loaded.shell_path, extra_findings=list(sink), qa="fast"
            )
            self.assertNotEqual(
                report.verdict,
                "failed",
                [f.message for f in report.findings if f.severity == "ERROR"],
            )

    def test_generate_is_byte_idempotent(self):
        with tempfile.TemporaryDirectory() as t:
            td = Path(t)
            loaded = self._extract(td)
            a, b = td / "a.pptx", td / "b.pptx"
            pptx_generate.generate(
                loaded.profile, loaded.shell_path, parse_idoc(KITCHEN), a
            )
            pptx_generate.generate(
                loaded.profile, loaded.shell_path, parse_idoc(KITCHEN), b
            )
            self.assertEqual(
                hashlib.sha256(a.read_bytes()).hexdigest(),
                hashlib.sha256(b.read_bytes()).hexdigest(),
            )


class XlsxKitchenSink(_Base):
    KIND = "xlsx"
    EXTRACT = xlsx_extract
    GENERATE = xlsx_generate

    def test_kitchen_sink_fills_named_regions_clean(self):
        with tempfile.TemporaryDirectory() as t:
            td = Path(t)
            loaded = self._extract(td)
            prof = loaded.profile
            regions = ((prof.get("surface") or {}).get("xlsx") or {}).get(
                "named_regions"
            ) or {}
            self.assertTrue(regions, "example xlsx has no named regions to fill")
            cells = {
                n: "X"
                for n, g in regions.items()
                if g.get("cardinality") == "single_cell"
            }
            regs = {
                n: [
                    ["a"] * max(g.get("cols", 1), 1)
                    for _ in range(min(max(g.get("rows", 1), 1), 2))
                ]
                for n, g in regions.items()
                if g.get("cardinality") != "single_cell"
            }
            out = td / "out.xlsx"
            sink = []
            xlsx_generate.generate(
                prof,
                loaded.shell_path,
                GridDocument(cells=cells, regions=regs),
                out,
                findings=sink,
            )
            self.assertTrue(out.is_file())
            report = run_qa(
                out, prof, shell=loaded.shell_path, extra_findings=list(sink), qa="fast"
            )
            errors = [f.message for f in report.findings if f.severity == "ERROR"]
            self.assertEqual(errors, [], errors)

    def test_generate_is_byte_idempotent(self):
        with tempfile.TemporaryDirectory() as t:
            td = Path(t)
            loaded = self._extract(td)
            grid = GridDocument(cells={}, regions={})
            a, b = td / "a.xlsx", td / "b.xlsx"
            xlsx_generate.generate(loaded.profile, loaded.shell_path, grid, a)
            xlsx_generate.generate(loaded.profile, loaded.shell_path, grid, b)
            self.assertEqual(
                hashlib.sha256(a.read_bytes()).hexdigest(),
                hashlib.sha256(b.read_bytes()).hexdigest(),
            )

    def test_kitchen_sink_authors_native_chart(self):
        # Exercise the NATIVE xlsx chart writer against the example template: chart
        # the showcase's own 'Data' sheet numbers (D2:D4, labels A2:A4). A real chart
        # part must be authored, with no degradation and a non-failed QA verdict.
        import zipfile

        with tempfile.TemporaryDirectory() as t:
            td = Path(t)
            loaded = self._extract(td)
            grid = GridDocument(
                charts=[
                    {
                        "sheet": "Data",
                        "type": "bar",
                        "title": "Volumi",
                        "anchor": "F2",
                        "data": "D2:D4",
                        "categories": "A2:A4",
                        "data_titles": False,
                    }
                ]
            )
            out = td / "out.xlsx"
            sink = []
            xlsx_generate.generate(
                loaded.profile, loaded.shell_path, grid, out, findings=sink
            )
            with zipfile.ZipFile(out) as z:
                chart_parts = [
                    n
                    for n in z.namelist()
                    if "/charts/chart" in n and n.endswith(".xml")
                ]
            self.assertTrue(chart_parts, "native chart not authored into the template")
            self.assertFalse([f for f in sink if f.check == "block_degraded"])
            report = run_qa(
                out,
                loaded.profile,
                shell=loaded.shell_path,
                extra_findings=list(sink),
                qa="fast",
            )
            self.assertNotEqual(
                report.verdict,
                "failed",
                [f.message for f in report.findings if f.severity == "ERROR"],
            )


if __name__ == "__main__":  # pragma: no cover
    unittest.main()
