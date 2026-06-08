# SPDX-License-Identifier: MIT
from __future__ import annotations

import json
import os
import sys
import tempfile
import unittest
from contextlib import redirect_stdout
from io import StringIO
from pathlib import Path
from unittest.mock import patch

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

from docx import Document
from docx.enum.style import WD_STYLE_TYPE
from docx.shared import Pt, RGBColor
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, PatternFill
from openpyxl.workbook.defined_name import DefinedName
from pptx import Presentation
from pptx.util import Inches

from brandkit import doctor
from brandkit.cli import main
from brandkit.profile import store


def _synthetic_template(path: Path) -> None:
    doc = Document()

    styles = doc.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(11)

    h1 = styles["Heading 1"]
    h1.font.name = "Aptos Display"
    h1.font.size = Pt(18)
    h1.font.bold = True
    h1.font.color.rgb = RGBColor(0x00, 0x5A, 0xAB)

    callout = styles.add_style("ACME Callout Info", WD_STYLE_TYPE.PARAGRAPH)
    callout.base_style = styles["Normal"]
    callout.font.bold = True
    callout.font.color.rgb = RGBColor(0x00, 0x5A, 0xAB)

    table_style = styles.add_style("ACME Table", WD_STYLE_TYPE.TABLE)
    table_style.base_style = styles["Table Grid"]

    doc.add_paragraph("{{title}}", style="Title")
    doc.add_paragraph("Example first-level title", style="Heading 1")
    doc.add_paragraph("General instructions: replace this demo text.", style="Normal")
    doc.add_paragraph("Brand note", style="ACME Callout Info")
    table = doc.add_table(rows=1, cols=2)
    table.style = "ACME Table"
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    doc.save(path)


def _synthetic_pptx_template(path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "{{title}}"
    slide.placeholders[1].text = "Example slide instructions"
    prs.save(path)


def _synthetic_xlsx_template(path: Path) -> None:
    wb = Workbook()
    ws = wb.active
    ws.title = "Report"
    ws["A1"] = "{{title}}"
    ws["A1"].font = Font(bold=True, color="005AAB")
    ws["A3"] = "Metric"
    ws["B3"] = "Value"
    ws["A4"] = "Example row"
    ws["B4"] = 1
    ws["B5"] = "=SUM(B4:B4)"
    ws["A3"].fill = PatternFill("solid", fgColor="D9EAF7")
    ws["B3"].fill = PatternFill("solid", fgColor="D9EAF7")
    wb.defined_names.add(DefinedName("data_region", attr_text="'Report'!$A$4:$B$4"))
    wb.defined_names.add(DefinedName("title_cell", attr_text="'Report'!$A$1"))
    wb.save(path)


class M1SmokeTest(unittest.TestCase):
    def test_m1_extract_verify_generate_docx(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tmp_path = Path(td)
            old_cwd = Path.cwd()
            os.chdir(tmp_path)
            try:
                template = tmp_path / "synthetic-template.docx"
                _synthetic_template(template)

                self.assertEqual(
                    main(
                        [
                            "extract",
                            "--name",
                            "acme",
                            "--template",
                            str(template),
                            "--scope",
                            "project",
                        ]
                    ),
                    0,
                )

                profile_path = tmp_path / "brand-kit" / "acme" / "profile.json"
                profile = json.loads(profile_path.read_text())
                self.assertEqual(profile["kind"], "docx")
                self.assertEqual(
                    profile["roles"]["heading.1"]["resolver"]["type"], "named_style"
                )
                self.assertTrue(profile["provenance"]["shell"]["sha256"])
                self.assertIn("artifact_catalog", profile)
                self.assertIn(
                    "ACME Callout Info",
                    profile["artifact_catalog"]["styles"]["paragraph"],
                )
                self.assertIn(
                    "word/styles.xml", profile["artifact_catalog"]["ooxml_parts"]
                )
                self.assertIn("capabilities", profile)
                self.assertTrue(profile["capabilities"]["extracts_all_ooxml_parts"])

                self.assertEqual(
                    main(
                        [
                            "verify",
                            "--name",
                            "acme",
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    0,
                )

                idoc = tmp_path / "idoc.json"
                idoc.write_text(
                    json.dumps(
                        {
                            "cover": {"title": "Quarterly Review"},
                            "blocks": [
                                {"type": "heading", "level": 1, "text": "Highlights"},
                                {
                                    "type": "paragraph",
                                    "text": "Revenue grew without markdown literals.",
                                },
                                {
                                    "type": "callout",
                                    "intent": "info",
                                    "text": "Use the brand callout style.",
                                },
                                {
                                    "type": "quote",
                                    "text": "A short quotation uses the quote role when available.",
                                },
                                {
                                    "type": "caption",
                                    "text": "Figure 1. Branded caption.",
                                },
                                {
                                    "type": "list",
                                    "items": [{"text": "Keep the shell brand."}],
                                },
                                {
                                    "type": "table",
                                    "columns": ["Area", "Status"],
                                    "rows": [
                                        ["Pipeline", "Healthy"],
                                        ["Delivery", "Green"],
                                    ],
                                },
                            ],
                        }
                    ),
                    encoding="utf-8",
                )
                out = tmp_path / "out.docx"

                self.assertEqual(
                    main(
                        [
                            "generate",
                            "--name",
                            "acme",
                            "--input",
                            str(idoc),
                            "--output",
                            str(out),
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    0,
                )
                self.assertTrue(out.is_file())

                generated = Document(out)
                text = "\n".join(p.text for p in generated.paragraphs)
                self.assertIn("Quarterly Review", text)
                self.assertIn("Highlights", text)
                self.assertIn(
                    "A short quotation uses the quote role when available.", text
                )
                self.assertIn("Figure 1. Branded caption.", text)
                self.assertNotIn("Example first-level title", text)
                self.assertNotIn("General instructions", text)
            finally:
                os.chdir(old_cwd)

    def test_profile_extra_files_cannot_escape_profile_dir(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            root = Path(td) / "profile"
            profile = {
                "kind": "docx",
                "provenance": {},
            }
            with self.assertRaises(store.ProfileStoreError):
                store.save_profile(
                    root, profile, b"fake", extra_files={"../escape.txt": "nope"}
                )

    def test_m2_extract_verify_generate_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tmp_path = Path(td)
            old_cwd = Path.cwd()
            os.chdir(tmp_path)
            try:
                template = tmp_path / "synthetic-template.pptx"
                _synthetic_pptx_template(template)

                self.assertEqual(
                    main(
                        [
                            "extract",
                            "--name",
                            "deck",
                            "--template",
                            str(template),
                            "--scope",
                            "project",
                        ]
                    ),
                    0,
                )
                profile = json.loads(
                    (tmp_path / "brand-kit" / "deck" / "profile.json").read_text()
                )
                self.assertEqual(profile["kind"], "pptx")
                self.assertEqual(
                    profile["roles"]["cover.title"]["resolver"]["type"], "placeholder"
                )
                self.assertIn("artifact_catalog", profile)
                self.assertIn("slide_layouts", profile["artifact_catalog"])
                self.assertTrue(profile["artifact_catalog"]["slide_layouts"])
                self.assertIn("capabilities", profile)
                self.assertTrue(profile["capabilities"]["extracts_layout_geometry"])

                idoc = tmp_path / "deck-idoc.json"
                idoc.write_text(
                    json.dumps(
                        {
                            "cover": {"title": "Board Update"},
                            "blocks": [
                                {"type": "heading", "level": 1, "text": "Highlights"},
                                {"type": "paragraph", "text": "Pipeline is healthy."},
                            ],
                        }
                    ),
                    encoding="utf-8",
                )
                out = tmp_path / "out.pptx"
                self.assertEqual(
                    main(
                        [
                            "verify",
                            "--name",
                            "deck",
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    0,
                )
                self.assertEqual(
                    main(
                        [
                            "generate",
                            "--name",
                            "deck",
                            "--input",
                            str(idoc),
                            "--output",
                            str(out),
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    0,
                )
                prs = Presentation(out)
                text = "\n".join(
                    shape.text
                    for slide in prs.slides
                    for shape in slide.shapes
                    if hasattr(shape, "text")
                )
                self.assertIn("Board Update", text)
                self.assertIn("Highlights", text)
                self.assertNotIn("Example slide instructions", text)
            finally:
                os.chdir(old_cwd)

    def test_pptx_long_content_splits_across_slides(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tmp_path = Path(td)
            old_cwd = Path.cwd()
            os.chdir(tmp_path)
            try:
                template = tmp_path / "synthetic-template.pptx"
                _synthetic_pptx_template(template)
                self.assertEqual(
                    main(
                        [
                            "extract",
                            "--name",
                            "deck",
                            "--template",
                            str(template),
                            "--scope",
                            "project",
                        ]
                    ),
                    0,
                )
                idoc = tmp_path / "deck-long.json"
                idoc.write_text(
                    json.dumps(
                        {
                            "cover": {"title": "Long Update"},
                            "blocks": [
                                {"type": "heading", "level": 1, "text": "Long Section"},
                                {
                                    "type": "paragraph",
                                    "text": " ".join(["capacity"] * 420),
                                },
                            ],
                        }
                    ),
                    encoding="utf-8",
                )
                out = tmp_path / "long.pptx"
                self.assertEqual(
                    main(
                        [
                            "generate",
                            "--name",
                            "deck",
                            "--input",
                            str(idoc),
                            "--output",
                            str(out),
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    0,
                )
                prs = Presentation(out)
                self.assertGreaterEqual(len(prs.slides), 3)
            finally:
                os.chdir(old_cwd)

    def test_m2_extract_verify_generate_xlsx(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tmp_path = Path(td)
            old_cwd = Path.cwd()
            os.chdir(tmp_path)
            try:
                template = tmp_path / "synthetic-template.xlsx"
                _synthetic_xlsx_template(template)

                self.assertEqual(
                    main(
                        [
                            "extract",
                            "--name",
                            "model",
                            "--template",
                            str(template),
                            "--scope",
                            "project",
                        ]
                    ),
                    0,
                )
                profile = json.loads(
                    (tmp_path / "brand-kit" / "model" / "profile.json").read_text()
                )
                self.assertEqual(profile["kind"], "xlsx")
                # De-literalized: no privileged "title" role. Every named range is a
                # generic ``named_range`` role keyed by its slugified own name. The
                # author's range names ("title_cell"/"data_region") are DATA carried
                # as ids, never code-side matching literals.
                self.assertEqual(
                    profile["roles"]["region.titlecell"]["resolver"]["type"],
                    "named_range",
                )
                self.assertEqual(
                    profile["roles"]["region.titlecell"]["resolver"]["name"],
                    "title_cell",
                )
                self.assertEqual(
                    profile["roles"]["region.dataregion"]["resolver"]["type"],
                    "named_range",
                )
                # Format-uniform comprehension inventories (geometry evidence only).
                surface = profile["surface"]["xlsx"]
                anchors_by_name = {a["name"]: a for a in surface["cover_anchors"]}
                self.assertEqual(
                    anchors_by_name["title_cell"]["cardinality"], "single_cell"
                )
                self.assertEqual(
                    anchors_by_name["title_cell"]["demo_value"], "{{title}}"
                )
                self.assertEqual(
                    anchors_by_name["data_region"]["cardinality"], "multi_cell"
                )
                self.assertEqual(
                    surface["fields"], []
                )  # legal-empty xlsx field inventory
                region_ids = {r["id"] for r in surface["regions"]}
                self.assertIn("region.dataregion", region_ids)  # multi-cell sample-data
                self.assertTrue(any(r["kind"] == "sheet" for r in surface["regions"]))
                self.assertIn("artifact_catalog", profile)
                self.assertIn("formulas", profile["artifact_catalog"])
                self.assertIn("Report!B5", profile["artifact_catalog"]["formulas"])
                self.assertIn(
                    "data_region", profile["artifact_catalog"]["named_ranges"]
                )
                self.assertIn("capabilities", profile)
                self.assertTrue(profile["capabilities"]["preserves_formulas_in_shell"])

                grid = tmp_path / "grid.json"
                grid.write_text(
                    json.dumps(
                        {
                            "cells": {"title_cell": "Quarterly Model"},
                            "regions": {"data_region": [["Pipeline", 42]]},
                        }
                    ),
                    encoding="utf-8",
                )
                out = tmp_path / "out.xlsx"
                self.assertEqual(
                    main(
                        [
                            "verify",
                            "--name",
                            "model",
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    0,
                )
                self.assertEqual(
                    main(
                        [
                            "generate",
                            "--name",
                            "model",
                            "--input",
                            str(grid),
                            "--output",
                            str(out),
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    0,
                )
                wb = load_workbook(out, data_only=False)
                ws = wb["Report"]
                self.assertEqual(ws["A1"].value, "Quarterly Model")
                self.assertEqual(ws["A4"].value, "Pipeline")
                self.assertEqual(ws["B4"].value, 42)
                # Formula preserved verbatim (never re-authored).
                self.assertEqual(ws["B5"].value, "=SUM(B4:B4)")
                # Recalc requested so Excel recomputes preserved formulas on open.
                self.assertTrue(wb.calculation.fullCalcOnLoad)
            finally:
                os.chdir(old_cwd)

    def test_xlsx_region_fill_cannot_overrun_named_range(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tmp_path = Path(td)
            old_cwd = Path.cwd()
            os.chdir(tmp_path)
            try:
                template = tmp_path / "synthetic-template.xlsx"
                _synthetic_xlsx_template(template)
                self.assertEqual(
                    main(
                        [
                            "extract",
                            "--name",
                            "model",
                            "--template",
                            str(template),
                            "--scope",
                            "project",
                        ]
                    ),
                    0,
                )
                grid = tmp_path / "grid-overrun.json"
                grid.write_text(
                    json.dumps({"regions": {"data_region": [["A", 1], ["B", 2]]}}),
                    encoding="utf-8",
                )
                out = tmp_path / "overrun.xlsx"
                self.assertEqual(
                    main(
                        [
                            "generate",
                            "--name",
                            "model",
                            "--input",
                            str(grid),
                            "--output",
                            str(out),
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    1,
                )
                self.assertFalse(out.exists())
            finally:
                os.chdir(old_cwd)

    def test_xlsx_unknown_named_range_is_error(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            tmp_path = Path(td)
            old_cwd = Path.cwd()
            os.chdir(tmp_path)
            try:
                template = tmp_path / "synthetic-template.xlsx"
                _synthetic_xlsx_template(template)
                self.assertEqual(
                    main(
                        [
                            "extract",
                            "--name",
                            "model",
                            "--template",
                            str(template),
                            "--scope",
                            "project",
                        ]
                    ),
                    0,
                )
                grid = tmp_path / "grid-unknown.json"
                grid.write_text(
                    json.dumps({"cells": {"missing_cell": "Nope"}}), encoding="utf-8"
                )
                out = tmp_path / "unknown.xlsx"
                self.assertEqual(
                    main(
                        [
                            "generate",
                            "--name",
                            "model",
                            "--input",
                            str(grid),
                            "--output",
                            str(out),
                            "--scope",
                            "project",
                            "--qa",
                            "fast",
                        ]
                    ),
                    1,
                )
                self.assertFalse(out.exists())
            finally:
                os.chdir(old_cwd)


class DoctorPreflightTest(unittest.TestCase):
    """The ``doctor`` command is a real preflight gate: exit code reflects the
    REQUIRED python deps, ``--json`` emits the verbatim probe(), and ``--fast``
    skips the slow soffice render smoke-test."""

    def test_doctor_returns_zero_when_required_deps_present(self) -> None:
        # All REQUIRED python deps are installed in this venv, so the gate passes.
        buf = StringIO()
        with redirect_stdout(buf):
            self.assertEqual(main(["doctor"]), 0)

    def test_doctor_returns_nonzero_when_required_dep_missing(self) -> None:
        full = doctor.probe(skip_visual_pipeline=True)
        broken = dict(full)
        broken["python_deps"] = {**full["python_deps"], "lxml": False}

        def fake_probe(*, skip_visual_pipeline=False):
            return broken

        buf = StringIO()
        with patch.object(doctor, "probe", fake_probe), redirect_stdout(buf):
            self.assertEqual(main(["doctor"]), 1)

    def test_doctor_missing_optional_does_not_gate(self) -> None:
        full = doctor.probe(skip_visual_pipeline=True)
        no_renderers = dict(full)
        no_renderers["binaries"] = {"soffice": False, "pdftoppm": False}
        no_renderers["visual_qa"] = False
        no_renderers["ocr_binaries"] = {"tesseract": False}
        no_renderers["ocr_qa"] = False

        def fake_probe(*, skip_visual_pipeline=False):
            return no_renderers

        buf = StringIO()
        with patch.object(doctor, "probe", fake_probe), redirect_stdout(buf):
            self.assertEqual(main(["doctor"]), 0)

    def test_doctor_json_emits_parseable_probe(self) -> None:
        expected = doctor.probe(skip_visual_pipeline=True)

        def fake_probe(*, skip_visual_pipeline=False):
            return expected

        buf = StringIO()
        with patch.object(doctor, "probe", fake_probe), redirect_stdout(buf):
            self.assertEqual(main(["doctor", "--json", "--fast"]), 0)

        out = buf.getvalue()
        # Human report lines must NOT be present in --json mode.
        self.assertNotIn("python:docx:", out)
        parsed = json.loads(out)
        self.assertEqual(parsed, expected)

    def test_doctor_fast_does_not_run_render_probes(self) -> None:
        calls: list = []

        def spy_pipeline(*args, **kwargs):
            calls.append((args, kwargs))
            return True, None

        buf = StringIO()
        with (
            patch.object(doctor, "_probe_visual_pipeline", spy_pipeline),
            redirect_stdout(buf),
        ):
            self.assertEqual(main(["doctor", "--fast"]), 0)

        self.assertEqual(calls, [])
        self.assertIn("visual QA: not probed", buf.getvalue())

    def test_probe_fast_marks_visual_not_probed(self) -> None:
        status = doctor.probe(skip_visual_pipeline=True)
        self.assertIsNone(status["visual_qa"])
        self.assertFalse(status["visual_qa_probed"])


if __name__ == "__main__":
    unittest.main()
