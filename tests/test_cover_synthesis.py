# SPDX-License-Identifier: MIT
"""Cluster E4 - universal cover synthesis for ``AnchorKind.NONE``.

A template whose cover-anchor detection recorded the STRUCTURAL FACT
``anchors.cover.kind == NONE`` has no slot the in-place cover machinery can ever
fill. E4 lets generation BUILD a minimal cover from the profile's own RESOLVABLE
``cover.*`` roles through ``resolve_role`` (the single brand chokepoint) - never
from literals - and no-ops byte-identically when nothing resolves.

The invariants pinned here:

  - kind==NONE + resolvable cover.* roles  -> cover synthesized from the roles,
    audited by an INFO ``cover_synthesized`` finding (role ids only, never brand
    text). DOCX: role-styled paragraphs before the first toc/body child. PPTX:
    a cover slide on the role-resolved layout (including on the reconcile path,
    where the empty cover_slots inventory previously dropped the cover).
  - kind != NONE (or the fact never recorded: ``anchors == {}``)  -> the entire
    existing cover path runs untouched, byte-identical, no new finding.
  - kind==NONE without a resolvable role / without cover content  -> exactly
    today's output, byte-identical, no new finding.
  - an authored title is NEVER dropped: title content + unresolvable
    ``cover.title`` declines synthesis entirely (deterministic fallback places it).
  - generate twice  -> byte-identical (idempotency).
  - XLSX: N/A by design - cover anchors are named ranges with no page/slide
    concept; generation never emits ``cover_synthesized`` for xlsx.
  - the cover_slots membership gate is unaffected: a ref into an empty
    cover_anchors inventory still rejects fail-closed at merge.
"""

from __future__ import annotations

import hashlib
import sys
import tempfile
import unittest
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[1] / "scripts"))

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from openpyxl import Workbook
from pptx import Presentation

from brandkit.formats.docx import generate as docx_generate
from brandkit.formats.docx import structure as docx_structure
from brandkit.formats.docx.structure import _local_name, w
from brandkit.formats.pptx import generate as pptx_generate
from brandkit.formats.xlsx import generate as xlsx_generate
from brandkit.grid.model import GridDocument
from brandkit.ir import model as ir
from brandkit.profile import comprehension as comp_mod
from brandkit.profile import schema, store
from brandkit.qa.model import Finding


# ---------------------------------------------------------------------------
# helpers
# ---------------------------------------------------------------------------
def _sha(path: Path) -> str:
    return hashlib.sha256(Path(path).read_bytes()).hexdigest()


def _add_toc_field(doc, instr='TOC \\o "1-3" \\h \\z \\u'):
    """Append a real TOC complex field paragraph (begin/instrText/separate/end)."""
    p = doc.add_paragraph()
    r = p.add_run()
    fb = OxmlElement("w:fldChar")
    fb.set(qn("w:fldCharType"), "begin")
    r._r.append(fb)
    r2 = p.add_run()
    it = OxmlElement("w:instrText")
    it.text = instr
    r2._r.append(it)
    r3 = p.add_run()
    fs = OxmlElement("w:fldChar")
    fs.set(qn("w:fldCharType"), "separate")
    r3._r.append(fs)
    p.add_run("entry .... 1")
    r6 = p.add_run()
    fe = OxmlElement("w:fldChar")
    fe.set(qn("w:fldCharType"), "end")
    r6._r.append(fe)
    return p


def _docx_shell(td: Path) -> Path:
    """A shell with a TOC + body but NO cover anchor (no SDT, no cover slot)."""
    shell = td / "shell.docx"
    doc = Document()
    _add_toc_field(doc)
    doc.add_paragraph("Body Heading", style="Heading 1")
    doc.add_paragraph("shell body text")
    doc.save(shell)
    return shell


def _docx_profile(*, kind=None, roles=None) -> dict:
    prof = schema.build_envelope("docx", {"name": "t"})
    prof["surface"] = {"docx": {}}
    prof["roles"] = roles if roles is not None else {"_index": []}
    if kind is not None:
        prof["anchors"] = {"cover": {"kind": kind, "slots_found": 0}}
    return prof


def _docx_cover_roles(*, title_style="Title", subtitle_style="Subtitle") -> dict:
    """cover.* roles resolving to styles python-docx's default template carries."""
    roles: dict = {"_index": []}
    if title_style is not None:
        roles["_index"].append("cover.title")
        roles["cover.title"] = {
            "resolver": {
                "type": "named_style",
                "style_id": title_style,
                "style_name": title_style,
            }
        }
    if subtitle_style is not None:
        roles["_index"].append("cover.subtitle")
        roles["cover.subtitle"] = {
            "resolver": {
                "type": "named_style",
                "style_id": subtitle_style,
                "style_name": subtitle_style,
            }
        }
    return roles


def _docx_generate(prof, shell, idoc, out) -> list[Finding]:
    findings: list[Finding] = []
    docx_generate.generate(prof, shell, idoc, out, findings=findings)
    return findings


def _para_regions(gen) -> tuple[dict, list]:
    """Map paragraph text -> classified region, plus document order of texts."""
    cls = docx_structure.classify_body_children(gen)
    children = list(gen.element.body)
    regions: dict = {}
    order: list = []
    for c in cls:
        el = children[c["index"]]
        if _local_name(el.tag) == "p":
            txt = "".join(t.text or "" for t in el.iter(w("t")))
            if txt:
                regions[txt] = c["region"]
                order.append(txt)
    return regions, order


def _pptx_shell(td: Path) -> Path:
    """A deck on python-pptx's default template (its layout 0 has title+subtitle)."""
    shell = td / "deck.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[1])
    prs.save(shell)
    return shell


def _pptx_cover_layout_name() -> str:
    return Presentation().slide_layouts[0].name


def _pptx_profile(*, kind=None, layout=None) -> dict:
    """A pptx profile whose cover.title role names ``layout`` (None = honest stub)."""
    prof = schema.build_envelope("pptx", {"name": "deck"})
    prof["surface"] = {"pptx": {"cover_anchors": []}}
    prof["roles"] = {
        "_index": ["cover.title"],
        "cover.title": {
            "resolver": {
                "type": "placeholder",
                "layout": layout,
                "ph_idx": 0 if layout else None,
                "ph_type": "title" if layout else None,
            }
        },
    }
    if kind is not None:
        prof["anchors"] = {"cover": {"kind": kind, "slots_found": 0}}
    return prof


def _present_empty_comp(prof: dict) -> None:
    """Stamp a PRESENT, sha-current comprehension with NO cover_slots (the only
    comprehension a kind==NONE shell can carry: its cover inventory is empty)."""
    prof["provenance"]["shell"]["sha256"] = "shellsha"
    block = schema.empty_comprehension()
    block["status"] = schema.ComprehensionStatus.PRESENT.value
    block["source_shell_sha256"] = "shellsha"
    block["confidence"] = 0.9
    prof["comprehension"] = block


def _pptx_generate(prof, shell, idoc, out) -> list[Finding]:
    findings: list[Finding] = []
    pptx_generate.generate(prof, shell, idoc, out, findings=findings)
    return findings


_NONE = schema.AnchorKind.NONE.value


# ---------------------------------------------------------------------------
# DOCX - synthesis fires only on the recorded fact + resolvable roles
# ---------------------------------------------------------------------------
class DocxCoverSynthesisTest(unittest.TestCase):
    def _idoc(self, *, title="Brand Title", subtitle="Brand Subtitle"):
        return ir.IntermediateDocument(
            blocks=[],
            cover=ir.Cover(
                title=[{"t": title}] if title else None,
                subtitle=[{"t": subtitle}] if subtitle else None,
            ),
        )

    def test_kind_none_with_resolvable_roles_synthesizes(self):
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _docx_shell(tp)
            prof = _docx_profile(kind=_NONE, roles=_docx_cover_roles())
            out = tp / "out.docx"
            findings = _docx_generate(prof, shell, self._idoc(), out)

            gen = Document(out)
            regions, order = _para_regions(gen)
            # Both slots synthesized ON THE COVER (before the TOC), title first.
            self.assertEqual(regions.get("Brand Title"), "cover")
            self.assertEqual(regions.get("Brand Subtitle"), "cover")
            self.assertLess(order.index("Brand Title"), order.index("Brand Subtitle"))
            # Each paragraph carries its ROLE-RESOLVED style (never a literal).
            styles = {p.text: p.style.style_id for p in gen.paragraphs if p.text}
            self.assertEqual(styles.get("Brand Title"), "Title")
            self.assertEqual(styles.get("Brand Subtitle"), "Subtitle")
            # Audit finding: INFO, names the fact + role ids, NEVER brand text.
            synth = [f for f in findings if f.check == "cover_synthesized"]
            self.assertEqual(len(synth), 1)
            self.assertEqual(synth[0].severity, schema.Severity.INFO.value)
            self.assertIn("kind=NONE", synth[0].message)
            self.assertIn("cover.title", synth[0].message)
            self.assertIn("cover.subtitle", synth[0].message)
            self.assertNotIn("Brand Title", synth[0].message)
            self.assertNotIn("Brand Subtitle", synth[0].message)
            # The old append-fallback WARNING is not ALSO recorded.
            self.assertFalse(
                any(
                    f.check == "cover_degraded"
                    and f.severity == schema.Severity.WARNING.value
                    for f in findings
                )
            )

    def test_kind_none_without_resolvable_roles_is_byte_identical(self):
        # No cover.* roles at all: synthesis must decline and the output must be
        # EXACTLY today's (the fact-absent profile is the pre-E4 ground truth).
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _docx_shell(tp)
            out_e4 = tp / "e4.docx"
            out_today = tp / "today.docx"
            f_e4 = _docx_generate(
                _docx_profile(kind=_NONE), shell, self._idoc(), out_e4
            )
            f_today = _docx_generate(_docx_profile(), shell, self._idoc(), out_today)
            self.assertEqual(_sha(out_e4), _sha(out_today))
            self.assertFalse(any(f.check == "cover_synthesized" for f in f_e4))
            # Today's degraded WARNING (append fallback) still fires on both.
            for fs in (f_e4, f_today):
                self.assertTrue(
                    any(
                        f.check == "cover_degraded"
                        and f.severity == schema.Severity.WARNING.value
                        for f in fs
                    )
                )

    def test_kind_none_with_stub_roles_is_byte_identical(self):
        # cover.* roles exist but resolve to styles the shell does NOT carry:
        # lookup misses, synthesis declines, output byte-identical to today.
        ghost = _docx_cover_roles(
            title_style="GhostTitleXYZ", subtitle_style="GhostSubXYZ"
        )
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _docx_shell(tp)
            out_e4 = tp / "e4.docx"
            out_today = tp / "today.docx"
            f_e4 = _docx_generate(
                _docx_profile(kind=_NONE, roles=ghost), shell, self._idoc(), out_e4
            )
            _docx_generate(_docx_profile(roles=ghost), shell, self._idoc(), out_today)
            self.assertEqual(_sha(out_e4), _sha(out_today))
            self.assertFalse(any(f.check == "cover_synthesized" for f in f_e4))

    def test_kind_not_none_skips_synthesis_byte_identical(self):
        # The explicit guard: any recorded kind other than NONE never opens the
        # synthesis branch, even with fully resolvable cover.* roles.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _docx_shell(tp)
            prof = _docx_profile(roles=_docx_cover_roles())
            prof["anchors"] = {
                "cover": {
                    "kind": schema.AnchorKind.SDT_ANCHORED.value,
                    "slots_found": 1,
                }
            }
            out_anchored = tp / "anchored.docx"
            out_today = tp / "today.docx"
            f_anchored = _docx_generate(prof, shell, self._idoc(), out_anchored)
            _docx_generate(
                _docx_profile(roles=_docx_cover_roles()),
                shell,
                self._idoc(),
                out_today,
            )
            self.assertEqual(_sha(out_anchored), _sha(out_today))
            self.assertFalse(any(f.check == "cover_synthesized" for f in f_anchored))

    def test_kind_none_without_cover_content_is_noop(self):
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _docx_shell(tp)
            idoc = ir.IntermediateDocument(blocks=[], cover=None)
            out_e4 = tp / "e4.docx"
            out_today = tp / "today.docx"
            f_e4 = _docx_generate(
                _docx_profile(kind=_NONE, roles=_docx_cover_roles()),
                shell,
                idoc,
                out_e4,
            )
            _docx_generate(
                _docx_profile(roles=_docx_cover_roles()), shell, idoc, out_today
            )
            self.assertEqual(_sha(out_e4), _sha(out_today))
            self.assertFalse(any(f.check == "cover_synthesized" for f in f_e4))

    def test_kind_none_empty_title_text_is_noop(self):
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _docx_shell(tp)
            idoc = self._idoc(title="", subtitle="")
            out_e4 = tp / "e4.docx"
            out_today = tp / "today.docx"
            f_e4 = _docx_generate(
                _docx_profile(kind=_NONE, roles=_docx_cover_roles()),
                shell,
                idoc,
                out_e4,
            )
            _docx_generate(
                _docx_profile(roles=_docx_cover_roles()), shell, idoc, out_today
            )
            self.assertEqual(_sha(out_e4), _sha(out_today))
            self.assertFalse(any(f.check == "cover_synthesized" for f in f_e4))

    def test_title_never_dropped_when_only_subtitle_resolves(self):
        # Title content + unresolvable cover.title: synthesis must decline
        # ENTIRELY (placing only the subtitle would drop the authored title);
        # the deterministic fallback still places the title. Byte-identical.
        roles = _docx_cover_roles(title_style="GhostTitleXYZ")
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _docx_shell(tp)
            out_e4 = tp / "e4.docx"
            out_today = tp / "today.docx"
            f_e4 = _docx_generate(
                _docx_profile(kind=_NONE, roles=roles), shell, self._idoc(), out_e4
            )
            _docx_generate(_docx_profile(roles=roles), shell, self._idoc(), out_today)
            self.assertEqual(_sha(out_e4), _sha(out_today))
            self.assertFalse(any(f.check == "cover_synthesized" for f in f_e4))
            gen = Document(out_e4)
            self.assertIn("Brand Title", [p.text for p in gen.paragraphs])

    def test_subtitle_only_content_synthesizes_subtitle_role(self):
        # No authored title at all: today writes NOTHING; E4 may place the
        # resolvable subtitle (a cover built from the roles that resolve).
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _docx_shell(tp)
            out = tp / "out.docx"
            findings = _docx_generate(
                _docx_profile(kind=_NONE, roles=_docx_cover_roles()),
                shell,
                self._idoc(title=""),
                out,
            )
            gen = Document(out)
            regions, _ = _para_regions(gen)
            self.assertEqual(regions.get("Brand Subtitle"), "cover")
            synth = [f for f in findings if f.check == "cover_synthesized"]
            self.assertEqual(len(synth), 1)
            self.assertIn("cover.subtitle", synth[0].message)
            self.assertNotIn("cover.title", synth[0].message)

    def test_unplaced_extra_fields_are_surfaced_not_silent(self):
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _docx_shell(tp)
            idoc = ir.IntermediateDocument(
                blocks=[],
                cover=ir.Cover(
                    title=[{"t": "Brand Title"}], fields={"date": "2026-06-10"}
                ),
            )
            out = tp / "out.docx"
            findings = _docx_generate(
                _docx_profile(kind=_NONE, roles=_docx_cover_roles()), shell, idoc, out
            )
            self.assertTrue(any(f.check == "cover_synthesized" for f in findings))
            notes = "\n".join(
                f.message for f in findings if f.check == "cover_degraded"
            )
            self.assertIn("date", notes)

    def test_synthesis_is_byte_idempotent(self):
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _docx_shell(tp)
            prof = _docx_profile(kind=_NONE, roles=_docx_cover_roles())
            out1 = tp / "a.docx"
            out2 = tp / "b.docx"
            _docx_generate(prof, shell, self._idoc(), out1)
            _docx_generate(prof, shell, self._idoc(), out2)
            self.assertEqual(_sha(out1), _sha(out2))


# ---------------------------------------------------------------------------
# PPTX - the reconcile path gains the synthesis; deterministic gains the audit
# ---------------------------------------------------------------------------
class PptxCoverSynthesisTest(unittest.TestCase):
    def _idoc(self):
        return ir.IntermediateDocument(
            blocks=[
                ir.Heading(level=1, runs=[{"t": "Section"}]),
                ir.Paragraph(runs=[{"t": "Body."}]),
            ],
            cover=ir.Cover(
                title=[{"t": "Deck Title"}], subtitle=[{"t": "Deck Subtitle"}]
            ),
        )

    def test_reconciled_kind_none_with_resolvable_role_synthesizes(self):
        # Comprehension PRESENT but (necessarily) no cover_slots - the kind==NONE
        # deck surfaced an empty inventory. Pre-E4 the authored cover was silently
        # dropped here; now the role-resolved layout builds it, audited.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _pptx_shell(tp)
            prof = _pptx_profile(kind=_NONE, layout=_pptx_cover_layout_name())
            _present_empty_comp(prof)
            self.assertTrue(store.comprehension_is_present(prof))
            out = tp / "out.pptx"
            findings = _pptx_generate(prof, shell, self._idoc(), out)

            prs = Presentation(out)
            titles = [
                s.shapes.title.text
                for s in prs.slides
                if s.shapes.title is not None and s.shapes.title.text
            ]
            # Exactly ONE synthesized cover, on the ROLE-RESOLVED layout, with the
            # subtitle placed too (the default cover layout carries both slots).
            self.assertEqual(titles.count("Deck Title"), 1)
            cover_slide = next(
                s
                for s in prs.slides
                if s.shapes.title is not None and s.shapes.title.text == "Deck Title"
            )
            self.assertEqual(cover_slide.slide_layout.name, _pptx_cover_layout_name())
            sub = pptx_generate._subtitle_placeholder(cover_slide)
            self.assertIsNotNone(sub)
            self.assertEqual(sub.text, "Deck Subtitle")
            synth = [f for f in findings if f.check == "cover_synthesized"]
            self.assertEqual(len(synth), 1)
            self.assertEqual(synth[0].severity, schema.Severity.INFO.value)
            self.assertIn("kind=NONE", synth[0].message)
            self.assertIn("cover.title", synth[0].message)
            self.assertNotIn("Deck Title", synth[0].message)

    def test_reconciled_without_recorded_fact_stays_coverless(self):
        # The control proving the trigger IS the recorded fact: same profile but
        # anchors never recorded -> the reconcile path keeps today's behavior
        # (no cover slide, no finding).
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _pptx_shell(tp)
            prof = _pptx_profile(kind=None, layout=_pptx_cover_layout_name())
            _present_empty_comp(prof)
            out = tp / "out.pptx"
            findings = _pptx_generate(prof, shell, self._idoc(), out)
            prs = Presentation(out)
            titles = [
                s.shapes.title.text
                for s in prs.slides
                if s.shapes.title is not None and s.shapes.title.text
            ]
            self.assertNotIn("Deck Title", titles)
            self.assertFalse(any(f.check == "cover_synthesized" for f in findings))

    def test_reconciled_kind_none_with_stub_role_is_byte_identical(self):
        # A freshly extracted kind==NONE deck: cover.title is an honest stub
        # (no layout) -> synthesis never fires, output byte-identical to the
        # fact-absent profile, no finding.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _pptx_shell(tp)
            prof_e4 = _pptx_profile(kind=_NONE, layout=None)
            _present_empty_comp(prof_e4)
            prof_today = _pptx_profile(kind=None, layout=None)
            _present_empty_comp(prof_today)
            out_e4 = tp / "e4.pptx"
            out_today = tp / "today.pptx"
            f_e4 = _pptx_generate(prof_e4, shell, self._idoc(), out_e4)
            _pptx_generate(prof_today, shell, self._idoc(), out_today)
            self.assertEqual(_sha(out_e4), _sha(out_today))
            self.assertFalse(any(f.check == "cover_synthesized" for f in f_e4))

    def test_deterministic_kind_none_with_resolvable_role_audited(self):
        # The deterministic rebuild already used the role-resolved layout; E4
        # adds the audit finding when that happened on a kind==NONE deck. The
        # deck bytes stay identical to the fact-absent profile.
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _pptx_shell(tp)
            out_e4 = tp / "e4.pptx"
            out_today = tp / "today.pptx"
            f_e4 = _pptx_generate(
                _pptx_profile(kind=_NONE, layout=_pptx_cover_layout_name()),
                shell,
                self._idoc(),
                out_e4,
            )
            f_today = _pptx_generate(
                _pptx_profile(kind=None, layout=_pptx_cover_layout_name()),
                shell,
                self._idoc(),
                out_today,
            )
            self.assertEqual(_sha(out_e4), _sha(out_today))
            self.assertTrue(any(f.check == "cover_synthesized" for f in f_e4))
            self.assertFalse(any(f.check == "cover_synthesized" for f in f_today))

    def test_deterministic_kind_none_with_stub_role_no_finding(self):
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _pptx_shell(tp)
            out_e4 = tp / "e4.pptx"
            out_today = tp / "today.pptx"
            f_e4 = _pptx_generate(
                _pptx_profile(kind=_NONE, layout=None), shell, self._idoc(), out_e4
            )
            _pptx_generate(
                _pptx_profile(kind=None, layout=None), shell, self._idoc(), out_today
            )
            self.assertEqual(_sha(out_e4), _sha(out_today))
            self.assertFalse(any(f.check == "cover_synthesized" for f in f_e4))

    def test_reconciled_synthesis_is_byte_idempotent(self):
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = _pptx_shell(tp)
            prof = _pptx_profile(kind=_NONE, layout=_pptx_cover_layout_name())
            _present_empty_comp(prof)
            out1 = tp / "a.pptx"
            out2 = tp / "b.pptx"
            _pptx_generate(prof, shell, self._idoc(), out1)
            _pptx_generate(prof, shell, self._idoc(), out2)
            self.assertEqual(_sha(out1), _sha(out2))


# ---------------------------------------------------------------------------
# XLSX - E4 is N/A by design (named ranges carry no cover-page concept)
# ---------------------------------------------------------------------------
class XlsxCoverSynthesisNotApplicableTest(unittest.TestCase):
    def test_xlsx_never_emits_cover_synthesized(self):
        """XLSX cover anchors are named ranges (data geometry, no page/slide
        semantics): there is no cover writer to synthesize into, so E4 is
        documented N/A and the finding can never appear for xlsx."""
        with tempfile.TemporaryDirectory() as td:
            tp = Path(td)
            shell = tp / "book.xlsx"
            wb = Workbook()
            wb.active.title = "S"
            wb.active["A1"] = "x"
            wb.save(shell)
            prof = schema.build_envelope("xlsx", {"name": "book"})
            prof["surface"] = {"xlsx": {"named_regions": {}}}
            prof["anchors"] = {"cover": {"kind": _NONE, "slots_found": 0}}
            out = tp / "out.xlsx"
            findings: list[Finding] = []
            xlsx_generate.generate(
                prof, shell, GridDocument(regions={}, cells={}), out, findings=findings
            )
            self.assertTrue(out.is_file())
            self.assertFalse(any(f.check == "cover_synthesized" for f in findings))


# ---------------------------------------------------------------------------
# Membership gate regression - E4 must not loosen the cover_slots fail-close
# ---------------------------------------------------------------------------
class CoverSlotsMembershipStillFailClosedTest(unittest.TestCase):
    def test_cover_slot_ref_into_empty_inventory_rejected(self):
        """A kind==NONE shell surfaces an EMPTY cover_anchors inventory; a model
        cover_slots ref into it must still reject fail-closed at merge (this is
        exactly why the E4 synthesis trigger can never race a bound slot)."""
        prof = schema.build_envelope("docx", {"name": "t"})
        prof["provenance"]["shell"]["sha256"] = "x"
        prof["surface"] = {"docx": {"cover_anchors": []}}
        prof["anchors"] = {"cover": {"kind": _NONE, "slots_found": 0}}
        res = comp_mod.merge(
            prof, {"cover_slots": {"sdt.1": {"fill_rule": "in_place"}}}
        )
        self.assertFalse(res.ok)
        self.assertEqual(prof["comprehension"]["status"], "rejected")


if __name__ == "__main__":  # pragma: no cover
    unittest.main()
