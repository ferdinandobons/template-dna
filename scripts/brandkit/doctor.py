# SPDX-License-Identifier: MIT
from __future__ import annotations

import importlib.util
import shutil
import subprocess
import tempfile
from pathlib import Path


REQUIRED = ("docx", "pptx", "openpyxl", "lxml", "PIL")
OPTIONAL_PYTHON = {"fitz": "PyMuPDF PDF raster fallback"}
OPTIONAL_BINARIES = {
    "soffice": "visual DOCX/PPTX/XLSX render",
    "pdftoppm": "PDF to PNG visual proof",
}
OPTIONAL_OCR_BINARIES = {"tesseract": "optional OCR visible-text audit"}
PYTHON_INSTALL_HINT = "python -m pip install -r requirements.txt"
OPTIONAL_PYTHON_INSTALL_HINTS = {
    "fitz": "python -m pip install PyMuPDF  # optional PDF raster fallback",
}
OPTIONAL_INSTALL_HINTS = {
    "soffice": (
        "macOS: brew install --cask libreoffice; "
        "if installed but unusable, open LibreOffice once and approve any macOS security prompt; "
        "Debian/Ubuntu: sudo apt-get install -y libreoffice; "
        "Fedora: sudo dnf install -y libreoffice; "
        "Windows: winget install TheDocumentFoundation.LibreOffice"
    ),
    "pdftoppm": (
        "macOS: brew install poppler; "
        "Debian/Ubuntu: sudo apt-get install -y poppler-utils; "
        "Fedora: sudo dnf install -y poppler-utils; "
        "Windows: install Poppler and add its bin directory to PATH"
    ),
    "tesseract": (
        "macOS: brew install tesseract; "
        "Debian/Ubuntu: sudo apt-get install -y tesseract-ocr; "
        "Fedora: sudo dnf install -y tesseract; "
        "Windows: winget install UB-Mannheim.TesseractOCR"
    ),
}
OPTIONAL_BINARY_PROBES = {
    "soffice": ("--headless", "--version"),
    "pdftoppm": ("-v",),
    "tesseract": ("--version",),
}
BINARY_PROBE_TIMEOUT_S = 10
VISUAL_PIPELINE_TIMEOUT_S = 45


def probe(*, skip_visual_pipeline: bool = False) -> dict:
    """Probe required/optional dependencies and the visual render pipeline.

    ``skip_visual_pipeline`` skips the slow serial soffice render smoke-test
    (which can spawn LibreOffice up to 3x). When skipped, ``visual_qa`` is set to
    ``None`` ("not probed") instead of running it; the binary version probes for
    soffice/pdftoppm still run since they are cheap.
    """
    deps = {name: importlib.util.find_spec(name) is not None for name in REQUIRED}
    optional_deps = {
        name: importlib.util.find_spec(name) is not None for name in OPTIONAL_PYTHON
    }
    bins: dict[str, bool] = {}
    paths: dict[str, str | None] = {}
    errors: dict[str, str] = {}
    for name in OPTIONAL_BINARIES:
        ok, path, error = _probe_binary(name)
        bins[name] = ok
        paths[name] = path
        if error:
            errors[name] = error
    ocr_bins: dict[str, bool] = {}
    ocr_paths: dict[str, str | None] = {}
    ocr_errors: dict[str, str] = {}
    for name in OPTIONAL_OCR_BINARIES:
        ok, path, error = _probe_binary(name)
        ocr_bins[name] = ok
        ocr_paths[name] = path
        if error:
            ocr_errors[name] = error
    visual_ok: bool | None = bool(bins.get("soffice")) and (
        bool(bins.get("pdftoppm")) or bool(optional_deps.get("fitz"))
    )
    visual_error = None
    if skip_visual_pipeline:
        # "not probed": skip the slow render smoke-test entirely.
        visual_ok = None
    elif visual_ok:
        visual_ok, visual_error = _probe_visual_pipeline(
            paths,
            {
                "pdftoppm": bool(bins.get("pdftoppm")),
                "fitz": bool(optional_deps.get("fitz")),
            },
        )
        if visual_error:
            errors["visual_qa"] = visual_error
    return {
        "python_deps": deps,
        "optional_python_deps": optional_deps,
        "binaries": bins,
        "binary_paths": paths,
        "binary_errors": errors,
        "visual_qa": visual_ok,
        "visual_qa_probed": not skip_visual_pipeline,
        "ocr_binaries": ocr_bins,
        "ocr_binary_paths": ocr_paths,
        "ocr_binary_errors": ocr_errors,
        "ocr_qa": bool(ocr_bins.get("tesseract")),
    }


def required_ok(status: dict) -> bool:
    """Return True only if every REQUIRED python dep is present.

    Optional renderers/OCR tools only downgrade visual QA and never gate this.
    """
    return all((status.get("python_deps") or {}).get(name) for name in REQUIRED)


def _probe_binary(name: str) -> tuple[bool, str | None, str | None]:
    path = shutil.which(name)
    if path is None:
        return False, None, "not found on PATH"
    preflight_error = _preflight_binary_error(name, path)
    if preflight_error:
        return False, path, preflight_error
    args = [path, *OPTIONAL_BINARY_PROBES.get(name, ("--version",))]
    try:
        proc = subprocess.run(
            args,
            capture_output=True,
            timeout=BINARY_PROBE_TIMEOUT_S,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        return False, path, str(exc)
    if proc.returncode != 0:
        stderr = _short_output(proc.stderr)
        stdout = _short_output(proc.stdout)
        detail = stderr or stdout or f"exit code {proc.returncode}"
        return False, path, detail
    return True, path, None


def _preflight_binary_error(name: str, path: str) -> str | None:
    if name == "soffice":
        return _soffice_app_signature_error(path)
    return None


def _soffice_app_signature_error(path: str) -> str | None:
    """Return a macOS LibreOffice signature error without launching the app.

    In the Codex desktop environment a quarantined or invalidly signed
    ``LibreOffice.app`` can abort inside AppKit before headless conversion even
    starts. Checking the bundle signature is safer than probing by conversion,
    because it avoids spawning the crashing process.
    """
    app = _libreoffice_app_for_soffice(path)
    if app is None:
        return None
    try:
        proc = subprocess.run(
            ["codesign", "--verify", "--deep", "--strict", str(app)],
            capture_output=True,
            timeout=BINARY_PROBE_TIMEOUT_S,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        return f"could not verify LibreOffice.app signature: {exc}"
    if proc.returncode == 0:
        return None
    detail = (
        _short_output(proc.stderr)
        or _short_output(proc.stdout)
        or f"exit code {proc.returncode}"
    )
    return f"LibreOffice.app signature invalid: {detail}"


def _libreoffice_app_for_soffice(path: str) -> Path | None:
    candidates = [
        Path("/Applications/LibreOffice.app"),
        Path(path).resolve().parents[2]
        if len(Path(path).resolve().parents) > 2
        else None,
    ]
    for candidate in candidates:
        if candidate and candidate.name == "LibreOffice.app" and candidate.is_dir():
            return candidate
    default = Path("/Applications/LibreOffice.app")
    return default if default.is_dir() else None


def _short_output(data) -> str:
    if data is None:
        return ""
    if isinstance(data, bytes):
        text = data.decode("utf-8", errors="replace")
    else:
        text = str(data)
    return " ".join(text.strip().split())[:240]


def _probe_visual_pipeline(
    paths: dict[str, str | None],
    rasterizers: dict[str, bool] | None = None,
) -> tuple[bool, str | None]:
    """Smoke-test the actual DOCX/PPTX/XLSX -> PDF -> PNG render pipeline.

    Version probes catch missing executables, but they do not prove LibreOffice can
    run headless conversion in the current environment. The visual audit supports
    all three OOXML formats, so the doctor probe must prove the full render chain
    for each one instead of treating DOCX success as a proxy for PPTX/XLSX.
    """
    rasterizers = rasterizers or {
        "pdftoppm": bool(paths.get("pdftoppm")),
        "fitz": importlib.util.find_spec("fitz") is not None,
    }
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        documents, error = _create_visual_probe_documents(tmp)
        if error:
            return False, error

        for document in documents:
            suffix = document.suffix.lstrip(".")
            ok, error = _probe_visual_document(
                paths,
                document,
                pdf_dir=tmp / f"pdf-{suffix}",
                png_dir=tmp / f"png-{suffix}",
                lo_profile=tmp / f"lo-profile-{suffix}",
                rasterizers=rasterizers,
            )
            if not ok:
                return False, f"{suffix} {error}"
    return True, None


def _create_visual_probe_documents(tmp: Path) -> tuple[list[Path], str | None]:
    """Create one tiny valid OOXML document for every visual-audit format."""
    try:
        from docx import Document
    except Exception as exc:
        return [], f"cannot create probe docx: {exc}"
    try:
        from pptx import Presentation
    except Exception as exc:
        return [], f"cannot create probe pptx: {exc}"
    try:
        from openpyxl import Workbook
    except Exception as exc:
        return [], f"cannot create probe xlsx: {exc}"

    try:
        docx_path = tmp / "probe.docx"
        doc = Document()
        doc.add_paragraph("BrandDocs visual QA probe")
        doc.save(docx_path)

        pptx_path = tmp / "probe.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        if slide.shapes.title is not None:
            slide.shapes.title.text = "BrandDocs visual QA probe"
        prs.save(pptx_path)

        xlsx_path = tmp / "probe.xlsx"
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "BrandDocs visual QA probe"
        wb.save(xlsx_path)
    except Exception as exc:
        return [], f"cannot create visual probe document: {exc}"

    return [docx_path, pptx_path, xlsx_path], None


def _probe_visual_document(
    paths: dict[str, str | None],
    document: Path,
    *,
    pdf_dir: Path,
    png_dir: Path,
    lo_profile: Path,
    rasterizers: dict[str, bool] | None = None,
) -> tuple[bool, str | None]:
    pdf_dir.mkdir()
    png_dir.mkdir()

    soffice_path = paths.get("soffice") or "soffice"
    try:
        soffice = subprocess.run(
            _soffice_convert_cmd(soffice_path, document, pdf_dir, lo_profile),
            capture_output=True,
            timeout=VISUAL_PIPELINE_TIMEOUT_S,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        return False, f"soffice convert failed: {exc}"
    if soffice.returncode != 0:
        return False, "soffice convert failed: " + (
            _short_output(soffice.stderr)
            or _short_output(soffice.stdout)
            or f"exit code {soffice.returncode}"
        )

    pdfs = list(pdf_dir.glob("*.pdf"))
    if not pdfs:
        return False, "soffice convert produced no PDF"

    rasterizers = rasterizers or {
        "pdftoppm": bool(paths.get("pdftoppm")),
        "fitz": False,
    }
    pdftoppm_path = paths.get("pdftoppm") or "pdftoppm"
    pdftoppm_error = None
    if rasterizers.get("pdftoppm"):
        ok, pdftoppm_error = _rasterize_pdf_with_pdftoppm(
            pdftoppm_path, pdfs[0], png_dir
        )
        if ok:
            return True, None
    if rasterizers.get("fitz"):
        ok, pymupdf_error = _rasterize_pdf_with_pymupdf(pdfs[0], png_dir, dpi=50)
        if ok:
            return True, None
        if pdftoppm_error:
            return (
                False,
                f"pdftoppm failed: {pdftoppm_error}; PyMuPDF failed: {pymupdf_error}",
            )
        return False, f"PyMuPDF failed: {pymupdf_error}"
    return False, pdftoppm_error or "no PDF rasterizer available"


def _rasterize_pdf_with_pdftoppm(
    pdftoppm_path: str, pdf: Path, png_dir: Path
) -> tuple[bool, str | None]:
    try:
        toppm = subprocess.run(
            [
                pdftoppm_path,
                "-png",
                "-r",
                "50",
                str(pdf),
                str(png_dir / "page"),
            ],
            capture_output=True,
            timeout=VISUAL_PIPELINE_TIMEOUT_S,
            check=False,
        )
    except (OSError, subprocess.TimeoutExpired) as exc:
        return False, str(exc)
    if toppm.returncode != 0:
        return False, (
            _short_output(toppm.stderr)
            or _short_output(toppm.stdout)
            or f"exit code {toppm.returncode}"
        )
    if not list(png_dir.glob("page-*.png")):
        return False, "produced no PNG"
    return True, None


def _rasterize_pdf_with_pymupdf(
    pdf: Path, png_dir: Path, *, dpi: int
) -> tuple[bool, str | None]:
    try:
        import fitz  # type: ignore[import-not-found]
    except Exception as exc:
        return False, f"fitz import failed: {exc}"
    try:
        doc = fitz.open(str(pdf))
        zoom = dpi / 72.0
        matrix = fitz.Matrix(zoom, zoom)
        for index, page in enumerate(doc, start=1):
            pix = page.get_pixmap(matrix=matrix, alpha=False)
            pix.save(str(png_dir / f"page-{index}.png"))
        if hasattr(doc, "close"):
            doc.close()
    except Exception as exc:
        return False, str(exc)
    if not list(png_dir.glob("page-*.png")):
        return False, "produced no PNG"
    return True, None


def _soffice_convert_cmd(
    soffice_path: str,
    document: Path,
    pdf_dir: Path,
    lo_profile: Path,
) -> list[str]:
    """Build a headless conversion command isolated from the user's LO profile."""
    return [
        soffice_path,
        f"-env:UserInstallation={lo_profile.as_uri()}",
        "--headless",
        "--nologo",
        "--nodefault",
        "--nolockcheck",
        "--norestore",
        "--nofirststartwizard",
        "--convert-to",
        "pdf",
        "--outdir",
        str(pdf_dir),
        str(document),
    ]


def print_report(status: dict | None = None) -> None:
    if status is None:
        status = probe()
    for name, ok in status["python_deps"].items():
        print(f"python:{name}: {'ok' if ok else 'missing'}")
    for name, ok in status.get("optional_python_deps", {}).items():
        print(f"python:{name}: {'ok' if ok else 'missing'} ({OPTIONAL_PYTHON[name]})")
    for name, ok in status["binaries"].items():
        if ok:
            label = "ok"
        elif status.get("binary_paths", {}).get(name):
            label = "unusable"
        else:
            label = "missing"
        msg = f"binary:{name}: {label} ({OPTIONAL_BINARIES[name]})"
        error = status.get("binary_errors", {}).get(name)
        if error and label == "unusable":
            msg += f" - {error}"
        print(msg)
    for name, ok in (status.get("ocr_binaries") or {}).items():
        if ok:
            label = "ok"
        elif (status.get("ocr_binary_paths") or {}).get(name):
            label = "unusable"
        else:
            label = "missing"
        msg = f"ocr:{name}: {label} ({OPTIONAL_OCR_BINARIES[name]})"
        error = (status.get("ocr_binary_errors") or {}).get(name)
        if error and label == "unusable":
            msg += f" - {error}"
        print(msg)
    if status.get("visual_qa") is None:
        print(
            "visual QA: not probed (--fast skips the soffice render smoke-test); "
            "L0 deterministic QA remains available"
        )
    elif status["visual_qa"]:
        print("visual QA: L1 proxy + L2 manifest available")
    else:
        suffix = ""
        if status.get("binary_errors", {}).get("visual_qa"):
            suffix = f" ({status['binary_errors']['visual_qa']})"
        print(f"visual QA disabled; L0 deterministic QA remains available{suffix}")
    if status.get("ocr_qa"):
        print("OCR QA: optional visible-text scan available")
    else:
        print(
            "OCR QA disabled; install optional OCR engine for rendered residual-text checks"
        )
    for hint in install_hints(status):
        print(hint)


def install_hints(status: dict) -> list[str]:
    """Return actionable install/repair hints for unavailable dependencies."""
    hints: list[str] = []
    missing_python = [
        name for name, ok in status.get("python_deps", {}).items() if not ok
    ]
    if missing_python:
        hints.append(
            "install:python: "
            f"{PYTHON_INSTALL_HINT}  # missing: {', '.join(sorted(missing_python))}"
        )

    def _binary_hints(bins: dict, paths: dict) -> None:
        for name, ok in (bins or {}).items():
            if ok:
                continue
            path = (paths or {}).get(name)
            action = "repair" if path else "install"
            detail = f" ({path})" if path else ""
            hint = OPTIONAL_INSTALL_HINTS.get(name)
            if hint:
                hints.append(f"{action}:{name}{detail}: {hint}")

    _binary_hints(status.get("binaries", {}), status.get("binary_paths") or {})
    _binary_hints(
        status.get("ocr_binaries") or {}, status.get("ocr_binary_paths") or {}
    )
    optional = status.get("optional_python_deps") or {}
    binaries = status.get("binaries") or {}
    if not binaries.get("pdftoppm") and not optional.get("fitz"):
        hint = OPTIONAL_PYTHON_INSTALL_HINTS.get("fitz")
        if hint:
            hints.append(f"install:fitz: {hint}")
    return hints
