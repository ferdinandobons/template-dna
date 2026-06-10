# SPDX-License-Identifier: MIT
"""Deterministic L0 checks for M1."""

from __future__ import annotations

import hashlib
import re
from contextlib import contextmanager
from typing import NamedTuple
from pathlib import Path

from brandkit.common import color as colorutil
from brandkit.common import text as textutil
from brandkit.ooxml import names, pack
from brandkit.profile import comprehension as comprehensionmod
from brandkit.profile import schema
from brandkit.profile.reconcile import confidence_clears_floor
from brandkit.qa.model import Finding

# WordprocessingML qualified-name builder (the spec-fixed namespace), so the shell
# readers below never hand-copy the namespace URI.
_W = names.make_qn("w")
# DrawingML qualified-name builder (the pptx run/theme namespace), shared with the
# pptx appearance collector below.
_A = names.make_qn("a")


# ---------------------------------------------------------------------------
# Per-run_qa artifact-load memo.
#
# One run_qa pass opens the SAME shell/output several times across independent
# checks (the workbook alone was loaded up to 9x per xlsx generate). Both files
# are immutable for the duration of one pass (the output is fully written before
# QA starts; the shell is never written), and every check is read-only over the
# loaded object, so loading each path at most once per pass is pure
# recomputation removal.
#
# Fail-closed scoping: the memo is ``None`` by default (every standalone check
# call loads fresh, exactly as before); ``gate.run_qa`` activates a fresh empty
# memo for the duration of ONE pass via :func:`load_memo` and always restores
# the previous state, so no loaded object can outlive its run_qa invocation. No
# fact is cached across runs - the shell-provenance sha check still reads file
# bytes straight from disk.
#
# The three loaders below are also the SINGLE import sites for the heavy Office
# libs (python-docx / openpyxl / python-pptx): importing them lazily keeps every
# CLI invocation from paying for all three when at most one format is active.
# ---------------------------------------------------------------------------
_LOAD_MEMO: dict | None = None


@contextmanager
def load_memo():
    """Scope a fresh artifact-load memo to one ``run_qa`` pass."""
    global _LOAD_MEMO
    previous = _LOAD_MEMO
    _LOAD_MEMO = {}
    try:
        yield
    finally:
        _LOAD_MEMO = previous


def _memo_load(kind: str, path, loader):
    if _LOAD_MEMO is None:
        return loader(path)
    key = (kind, str(path))
    if key not in _LOAD_MEMO:
        _LOAD_MEMO[key] = loader(path)
    return _LOAD_MEMO[key]


def _load_docx(path):
    from docx import Document

    return _memo_load("docx", path, Document)


def _load_pptx(path):
    from pptx import Presentation

    return _memo_load("pptx", path, Presentation)


def _load_xlsx(path):
    # Every QA workbook read uses the SAME flags (data_only=False), so the memo
    # never needs to key on loader options.
    from openpyxl import load_workbook

    return _memo_load("xlsx", path, lambda p: load_workbook(p, data_only=False))


def check_profile(profile: dict) -> list[Finding]:
    findings: list[Finding] = []
    for problem in schema.validate(profile):
        findings.append(Finding("schema", schema.Severity.ERROR.value, problem))
    for rid in (profile.get("roles") or {}).get("_index", []):
        entry = profile.get("roles", {}).get(rid)
        if not entry or not entry.get("resolver"):
            findings.append(
                Finding(
                    "every_role_resolves",
                    schema.Severity.ERROR.value,
                    f"{rid} has no resolver",
                )
            )
    return findings


def check_no_duplicate_parts(target) -> list[Finding]:
    """ERROR when the output OOXML package has duplicate ZIP part names.

    A valid OPC package is a ZIP whose part names are unique; a duplicate (e.g. a
    slide-partname collision after removing a non-last slide) makes PowerPoint/Word
    show a repair dialog. Cheap, format-agnostic integrity backstop that no
    text/structure scan catches; a missing/garbage file is handled by other checks.
    """
    import zipfile
    from collections import Counter

    if target is None:
        return []
    try:
        with zipfile.ZipFile(target) as z:
            names = z.namelist()
    except (OSError, zipfile.BadZipFile):
        return []
    dups = sorted(n for n, c in Counter(names).items() if c > 1)
    if not dups:
        return []
    shown = ", ".join(dups[:5]) + (" ..." if len(dups) > 5 else "")
    return [
        Finding(
            check="package_integrity",
            severity=schema.Severity.ERROR.value,
            message=f"output package has duplicate part name(s): {shown}",
        )
    ]


def check_shell_provenance(shell, profile: dict) -> list[Finding]:
    """ERROR when the saved shell no longer matches profile provenance.

    The Brand Profile records ``provenance.shell.sha256`` at extraction time. A
    later hand edit, corrupt copy, or tamper must be load-bearing in QA: generation
    from a drifted shell is not the same brand contract the profile verified.
    """
    if shell is None:
        return []
    recorded = ((profile.get("provenance") or {}).get("shell") or {}).get("sha256")
    if not recorded:
        return []

    shell_path = Path(shell)
    if not shell_path.is_file():
        return [
            Finding(
                "shell_provenance",
                schema.Severity.ERROR.value,
                f"recorded shell hash exists but shell is missing: {shell_path}",
                location=str(shell_path),
            )
        ]

    try:
        actual = _sha256_file(shell_path)
    except OSError as exc:
        return [
            Finding(
                "shell_provenance",
                schema.Severity.ERROR.value,
                f"could not hash shell {shell_path}: {exc}",
                location=str(shell_path),
            )
        ]

    if actual != recorded:
        return [
            Finding(
                "shell_provenance",
                schema.Severity.ERROR.value,
                f"shell hash drifted: recorded {recorded}, actual {actual}",
                location=str(shell_path),
            )
        ]
    return []


def _sha256_file(path: Path, *, chunk: int = 1 << 20) -> str:
    h = hashlib.sha256()
    with open(path, "rb") as fh:
        for block in iter(lambda: fh.read(chunk), b""):
            h.update(block)
    return h.hexdigest()


def check_blend_shell_provenance(shell, profile: dict) -> list[Finding]:
    """ERROR when a recorded SECONDARY (blend) shell no longer matches provenance.

    Multi-template blending records every donor in ``provenance.blended_shells``
    (sha256 + content-addressed path under the profile's ``template/`` dir) and
    keys its merged value-facts' evidence on those binaries. A missing, escaping,
    unreadable, or hash-drifted donor means that evidence is no longer auditable,
    so each is load-bearing (ERROR) - the secondary-shell peer of
    :func:`check_shell_provenance`. A never-blended profile (no
    ``blended_shells``, no ``blend`` block) yields zero findings at zero cost.
    """
    if shell is None:
        return []
    findings: list[Finding] = []
    prov = profile.get("provenance")
    entries = prov.get("blended_shells") if isinstance(prov, dict) else None
    if entries is not None and not isinstance(entries, list):
        findings.append(
            Finding(
                "blend_shell_provenance",
                schema.Severity.ERROR.value,
                "provenance.blended_shells is malformed (must be a list)",
            )
        )
        entries = None

    # The profile root: shell is <root>/template/shell.<ext>.
    profile_root = Path(shell).parent.parent
    root_resolved = profile_root.resolve()
    for entry in sorted(
        entries or [],
        key=lambda e: str(e.get("sha256")) if isinstance(e, dict) else "",
    ):
        if not isinstance(entry, dict):
            findings.append(
                Finding(
                    "blend_shell_provenance",
                    schema.Severity.ERROR.value,
                    f"provenance.blended_shells entry is malformed: {entry!r}",
                )
            )
            continue
        rel = entry.get("path")
        recorded = entry.get("sha256")
        if not isinstance(rel, str) or not rel or not isinstance(recorded, str):
            findings.append(
                Finding(
                    "blend_shell_provenance",
                    schema.Severity.ERROR.value,
                    f"provenance.blended_shells entry is malformed: {entry!r}",
                )
            )
            continue
        target = profile_root / rel
        try:
            target.resolve().relative_to(root_resolved)
        except ValueError:
            findings.append(
                Finding(
                    "blend_shell_provenance",
                    schema.Severity.ERROR.value,
                    f"recorded blend shell path escapes the profile dir: {rel!r}",
                    location=str(target),
                )
            )
            continue
        if not target.is_file():
            findings.append(
                Finding(
                    "blend_shell_provenance",
                    schema.Severity.ERROR.value,
                    f"recorded blend shell hash exists but shell is missing: {rel}",
                    location=str(target),
                )
            )
            continue
        try:
            actual = _sha256_file(target)
        except OSError as exc:
            findings.append(
                Finding(
                    "blend_shell_provenance",
                    schema.Severity.ERROR.value,
                    f"could not hash blend shell {rel}: {exc}",
                    location=str(target),
                )
            )
            continue
        if actual != recorded:
            findings.append(
                Finding(
                    "blend_shell_provenance",
                    schema.Severity.ERROR.value,
                    f"blend shell hash drifted: recorded {recorded}, actual {actual}",
                    location=str(target),
                )
            )

    # Stale-binding belt: the blend ledger must be bound to the LIVE primary
    # shell. Unreachable through the supported flow (a re-extract rewrites
    # profile.json wholesale), so this only fires on hand-editing - honest
    # fail-closed rather than a silent stale ledger.
    block = profile.get("blend")
    if isinstance(block, dict):
        recorded = block.get("source_shell_sha256")
        live = ((profile.get("provenance") or {}).get("shell") or {}).get("sha256")
        if recorded and live and recorded != live:
            findings.append(
                Finding(
                    "blend_shell_provenance",
                    schema.Severity.ERROR.value,
                    f"blend ledger bound to primary shell {recorded} but live "
                    f"shell is {live}",
                )
            )
    return findings


def check_resolver_targets(shell, profile: dict) -> list[Finding]:
    """Verify every role's resolver target actually exists in the ``shell``.

    This is the deterministic backstop for the core promise ("apply only artifacts
    the profile *proved* exist"): it opens the shell once and confirms each role's
    concrete target is present -
      - docx: ``style_id`` or ``style_name`` ∈ ``doc.styles``;
      - pptx: ``layout`` ∈ the presentation's slide-layout names;
      - xlsx: ``name`` ∈ the workbook's defined names.
    A role whose target is missing yields an ERROR ``resolver_targets_exist``
    finding. ``shell`` may be a path or None; when None the check is skipped (the
    caller is responsible for supplying the shell at verify time).
    """
    if shell is None:
        return []
    kind = profile.get("kind")
    roles = profile.get("roles") or {}
    role_ids = [r for r in roles if r != "_index"]
    findings: list[Finding] = []
    try:
        if kind == schema.Kind.DOCX.value:
            present = _docx_style_keys(shell)
            for rid in role_ids:
                resolver = (roles.get(rid) or {}).get("resolver") or {}
                if resolver.get("type") not in (
                    schema.ResolverType.NAMED_STYLE.value,
                    None,
                ):
                    continue
                sid = resolver.get("style_id")
                sname = resolver.get("style_name")
                if not sid and not sname:
                    continue
                if (sid and sid in present) or (sname and sname in present):
                    continue
                findings.append(
                    Finding(
                        "resolver_targets_exist",
                        schema.Severity.ERROR.value,
                        f"role {rid!r} target style {(sid or sname)!r} not found in shell",
                        location=rid,
                    )
                )
        elif kind == schema.Kind.PPTX.value:
            layout_names = _pptx_layout_names(shell)
            for rid in role_ids:
                resolver = (roles.get(rid) or {}).get("resolver") or {}
                if resolver.get("type") != schema.ResolverType.PLACEHOLDER.value:
                    continue
                layout = resolver.get("layout")
                if layout is None or layout in layout_names:
                    continue
                findings.append(
                    Finding(
                        "resolver_targets_exist",
                        schema.Severity.ERROR.value,
                        f"role {rid!r} layout {layout!r} not found in shell (have {sorted(layout_names)})",
                        location=rid,
                    )
                )
        elif kind == schema.Kind.XLSX.value:
            defined = _xlsx_defined_names(shell)
            for rid in role_ids:
                resolver = (roles.get(rid) or {}).get("resolver") or {}
                if resolver.get("type") != schema.ResolverType.NAMED_RANGE.value:
                    continue
                name = resolver.get("name")
                if name is None or name in defined:
                    continue
                findings.append(
                    Finding(
                        "resolver_targets_exist",
                        schema.Severity.ERROR.value,
                        f"role {rid!r} named range {name!r} not found in shell (have {sorted(defined)})",
                        location=rid,
                    )
                )
            # number_format roles must resolve to a mask the shell ACTUALLY uses
            # (the engine never fabricates a format) - the shell-backed peer of the
            # schema's intra-profile number_format consistency check.
            masks = _xlsx_number_format_masks(shell)
            for rid in role_ids:
                resolver = (roles.get(rid) or {}).get("resolver") or {}
                if resolver.get("type") != schema.ResolverType.NUMBER_FORMAT.value:
                    continue
                mask = resolver.get("number_format")
                if mask is None or mask in masks:
                    continue
                findings.append(
                    Finding(
                        "resolver_targets_exist",
                        schema.Severity.ERROR.value,
                        f"role {rid!r} number format {mask!r} not among the shell's "
                        "used formats",
                        location=rid,
                    )
                )
    except Exception as exc:  # opening the shell must never crash the gate
        findings.append(
            Finding(
                "resolver_targets_exist",
                schema.Severity.WARNING.value,
                f"could not verify resolver targets against shell: {exc}",
            )
        )
    return findings


def _docx_style_keys(shell) -> set:
    doc = _load_docx(shell)
    keys: set = set()
    for style in doc.styles:
        sid = getattr(style, "style_id", None)
        if sid:
            keys.add(sid)
        name = getattr(style, "name", None)
        if name:
            keys.add(name)
    return keys


def _docx_shell_fonts(shell, profile: dict) -> set:
    """Font names the shell makes available: every ``w:font`` in its fontTable, the
    theme major/minor latin typefaces, and the Arial ``docDefaults`` baseline."""
    fonts: set = {"Arial"}
    theme_fonts = (profile.get("theme") or {}).get("fonts") or {}
    for key in ("major", "minor"):
        latin = (theme_fonts.get(key) or {}).get("latin")
        if latin:
            fonts.add(latin)
    try:
        xml = pack.read_part(shell, "word/fontTable.xml")
    except KeyError:
        return fonts
    root = pack.parse_xml_bytes(xml)
    for font_el in root.findall(_W("font")):
        name = font_el.get(_W("name"))
        if name:
            fonts.add(name)
    return fonts


def _docx_shell_run_sizes_and_hexes(shell) -> tuple[set[int], set[str]]:
    """The PROVENANCE sets an applied size/color is validated against, read from the
    template's ``word/document.xml`` runs in a SINGLE parse:

      - sizes: every ``w:sz@w:val`` (half-points) the template's OWN runs carry;
      - hexes: every ``w:color@w:val`` (normalized hex), skipping ``auto``.

    A captured size/hex is only trusted when the template itself uses it (not a sanity
    bound). The hardened parser (``parse_xml_bytes``, ``resolve_entities=False``)
    matches every other shell-XML read. A missing/garbage document part yields empty
    sets - the caller then fails closed on any applied size/color."""
    sizes: set[int] = set()
    hexes: set[str] = set()
    try:
        xml = pack.read_part(shell, "word/document.xml")
    except KeyError:
        return sizes, hexes
    root = pack.parse_xml_bytes(xml)
    for sz in root.iter(_W("sz")):
        val = sz.get(_W("val"))
        if val is None:
            continue
        try:
            sizes.add(int(val))
        except (TypeError, ValueError):
            continue
    for color in root.iter(_W("color")):
        val = color.get(_W("val"))
        if not val or val.lower() == "auto":
            continue
        try:
            hexes.add(colorutil.normalize_hex(val))
        except (ValueError, AttributeError):
            continue
    return sizes, hexes


# The WordprocessingML themeColor -> clrScheme slot alias table, the single source of
# truth in brandkit.common.color (shared with the resolver's hex enrichment). A token
# outside it has no palette slot and is rejected (fail-closed).
_WML_THEME_TO_SLOT = colorutil.WML_THEME_TO_SLOT


def _docx_theme_palette(shell) -> dict[str, str]:
    """The template's parsed theme palette (slot -> hex). Reads the shell's own
    ``theme1.xml``; an absent part yields an empty palette (every theme-token color
    then fails closed)."""
    try:
        xml = pack.read_part(shell, "word/theme/theme1.xml")
    except KeyError:
        return {}
    return colorutil.parse_theme_colors(xml)


# ---------------------------------------------------------------------------
# Format-neutral shell appearance facts (A3): the PROVENANCE every applied
# font/size/color is re-validated against, collected per kind into ONE shape.
# ---------------------------------------------------------------------------
class ShellAppearanceFacts(NamedTuple):
    """The typographic provenance a shell proves it carries, format-neutral.

    The single shape every per-kind collector fills so the membership loops in
    :func:`check_appearance_targets` are format-agnostic:

      - ``fonts``: the font names the shell makes AVAILABLE (the allow-set an applied
        font is validated against - the fontTable/theme/Arial baseline, widened, never
        narrowed);
      - ``sizes``: every run size the shell's OWN runs carry, in HALF-POINTS (the unit
        :data:`theme.fonts.body.size_hp` and ``role.appearance.size_hp`` use);
      - ``hexes``: every explicit run hex the shell carries (normalized ``RRGGBB``);
      - ``palette``: the parsed ``clrScheme`` slot -> normalized hex (the theme palette a
        theme-token color resolves against).
    """

    fonts: set[str]
    sizes: set[int]
    hexes: set[str]
    palette: dict[str, str]


def _docx_collect_appearance_facts(shell, profile: dict) -> ShellAppearanceFacts:
    """The docx appearance facts, BYTE-IDENTICALLY composed from the existing docx
    shell readers (``_docx_shell_fonts`` + ``_docx_shell_run_sizes_and_hexes`` +
    ``_docx_theme_palette``). No new docx parsing: this is a pure assembly so the docx
    verify verdict is unchanged."""
    fonts = _docx_shell_fonts(shell, profile)
    sizes, hexes = _docx_shell_run_sizes_and_hexes(shell)
    palette = _docx_theme_palette(shell)
    return ShellAppearanceFacts(fonts=fonts, sizes=sizes, hexes=hexes, palette=palette)


# pptx ``a:rPr@sz`` / ``a:defRPr@sz`` are in CENTIPOINTS (1/100 pt); the engine's
# size unit is HALF-POINTS, so a value divides by 50 (e.g. 1800 -> 36 half-points).
_PPTX_SZ_PER_HALF_POINT = 50
_PPTX_SLIDE_PART = re.compile(r"ppt/slides/slide\d+\.xml")


def _pptx_collect_appearance_facts(shell, profile: dict) -> ShellAppearanceFacts:
    """The pptx appearance facts, read from the deck's own parts.

    - ``fonts``: every ``a:latin@typeface`` in the theme major/minor font scheme
      (``ppt/theme/theme1.xml``), widened by the Arial baseline and the profile's
      captured ``theme.fonts.major/minor.latin`` (the allow-set, never narrowed);
    - ``sizes``: every ``a:rPr@sz`` / ``a:defRPr@sz`` across the slide parts,
      converted from centipoints to half-points (``÷50``);
    - ``hexes``: every ``a:srgbClr@val`` under a run/defRPr ``a:solidFill`` on the
      slides, normalized;
    - ``palette``: the parsed clrScheme of the deck's theme.
    """
    fonts: set[str] = {"Arial"}
    theme_fonts = (profile.get("theme") or {}).get("fonts") or {}
    for key in ("major", "minor"):
        latin = (theme_fonts.get(key) or {}).get("latin")
        if latin:
            fonts.add(latin)
    try:
        theme_xml = pack.read_part(shell, "ppt/theme/theme1.xml")
    except KeyError:
        palette: dict[str, str] = {}
    else:
        palette = colorutil.parse_theme_colors(theme_xml)
        theme_root = pack.parse_xml_bytes(theme_xml)
        for latin in theme_root.iter(_A("latin")):
            face = latin.get("typeface")
            if face:
                fonts.add(face)
    sizes: set[int] = set()
    hexes: set[str] = set()
    for part in pack.list_parts(shell):
        if not _PPTX_SLIDE_PART.fullmatch(part):
            continue
        root = pack.parse_xml_bytes(pack.read_part(shell, part))
        for tag in ("rPr", "defRPr"):
            for rpr in root.iter(_A(tag)):
                val = rpr.get("sz")
                if val is None:
                    continue
                try:
                    # round() (not floor //) to match the capture side
                    # (pptx/typography: round(size.pt * 2) == round(centipoints / 50)),
                    # so a correctly applied fractional size never spuriously fails closed.
                    sizes.add(round(int(val) / _PPTX_SZ_PER_HALF_POINT))
                except (TypeError, ValueError):
                    continue
        for srgb in root.iter(_A("srgbClr")):
            val = srgb.get("val")
            if not val:
                continue
            try:
                hexes.add(colorutil.normalize_hex(val))
            except (ValueError, AttributeError):
                continue
    return ShellAppearanceFacts(fonts=fonts, sizes=sizes, hexes=hexes, palette=palette)


def _xlsx_collect_appearance_facts(shell, profile: dict) -> ShellAppearanceFacts:
    """The xlsx appearance facts, read from the workbook's own cells/styles/theme.

    - ``fonts``: every ``cell.font.name`` across the materialized cells, widened by
      each NamedStyle font name and the Arial baseline (the allow-set);
    - ``sizes``: ``round(font.sz * 2)`` for every explicit cell-font size (half-points,
      the same unit the xlsx capture records);
    - ``hexes``: every ``font.color.rgb`` that is an ``'rgb'`` color, 8-digit ARGB
      alpha stripped, normalized (non-rgb/indexed/auto colors contribute nothing);
    - ``palette``: the parsed clrScheme of the workbook's theme.
    """
    fonts: set[str] = {"Arial"}
    sizes: set[int] = set()
    hexes: set[str] = set()
    wb = _load_xlsx(shell)
    # ``wb.named_styles`` is a list of style NAMES (strings); the NamedStyle objects
    # (which carry the ``.font``) live on ``wb._named_styles``. Guard both for
    # openpyxl-version robustness - a NamedStyle font widens the allow-set.
    for style in getattr(wb, "_named_styles", []) or []:
        font = getattr(style, "font", None)
        name = getattr(font, "name", None) if font is not None else None
        if name:
            fonts.add(name)
    for ws in wb.worksheets:
        for cell in ws._cells.values():
            font = cell.font
            if font is None:
                continue
            if font.name:
                fonts.add(font.name)
            if font.sz is not None:
                try:
                    sizes.add(round(float(font.sz) * 2))
                except (TypeError, ValueError):
                    pass
            color = getattr(font, "color", None)
            if color is None or getattr(color, "type", None) != "rgb":
                continue
            rgb = color.rgb
            if not isinstance(rgb, str):
                continue
            hexpart = rgb[2:] if len(rgb) == 8 else rgb
            try:
                hexes.add(colorutil.normalize_hex(hexpart))
            except (ValueError, AttributeError):
                continue
    try:
        theme_xml = pack.read_part(shell, "xl/theme/theme1.xml")
    except KeyError:
        palette: dict[str, str] = {}
    else:
        palette = colorutil.parse_theme_colors(theme_xml)
    return ShellAppearanceFacts(fonts=fonts, sizes=sizes, hexes=hexes, palette=palette)


# The per-kind shell appearance collector (peer of :data:`_COMPONENT_COUNTERS`). Each
# fills the format-neutral :class:`ShellAppearanceFacts`; the lifted gate dispatches
# off ``profile['kind']`` and feeds the result into the unchanged membership loops.
_SHELL_APPEARANCE_COLLECTORS = {
    schema.Kind.DOCX.value: _docx_collect_appearance_facts,
    schema.Kind.PPTX.value: _pptx_collect_appearance_facts,
    schema.Kind.XLSX.value: _xlsx_collect_appearance_facts,
}


def _collect_applied_appearance(
    profile: dict,
) -> tuple[list[tuple[str, int]], list[tuple[str, dict]]]:
    """Gather every (where, size_hp) and (where, color-obj) the engine will APPLY.

    Scans per-role ``appearance`` plus the document body defaults
    (``theme.fonts.body.size_hp`` for size, ``theme.text.body.color`` for color),
    plus EVERY ``theme.palette[*].ref`` color object - a run carrying a palette
    color token is realized into that ref by the resolver, so each palette ref is an
    APPLIED color the shell must independently prove it carries. Folding the palette
    refs in here means ``check_appearance_targets`` re-validates them with the exact
    same loop (hex vs the palette UNION the observed ``w:color``; theme token vs the
    mapped ``clrScheme`` slot), fail-closed - the model can NAME a palette color but
    the deterministic shell still has to back it."""
    sizes: list[tuple[str, int]] = []
    colors: list[tuple[str, dict]] = []
    for rid, entry in (profile.get("roles") or {}).items():
        if rid == "_index" or not isinstance(entry, dict):
            continue
        appearance = entry.get("appearance") or {}
        size_hp = appearance.get("size_hp")
        if size_hp:
            sizes.append((f"role {rid!r}", int(size_hp)))
        color = appearance.get("color")
        if isinstance(color, dict):
            colors.append((f"role {rid!r}", color))
    theme = profile.get("theme") or {}
    body_size = ((theme.get("fonts") or {}).get("body") or {}).get("size_hp")
    if body_size:
        sizes.append(("theme.fonts.body", int(body_size)))
    body_color = ((theme.get("text") or {}).get("body") or {}).get("color")
    if isinstance(body_color, dict):
        colors.append(("theme.text.body", body_color))
    # Every palette ref is an applicable run color (a run color token resolves to
    # it), so re-validate it against the shell with the same loop. Sorted for a
    # deterministic finding order.
    for key in sorted((theme.get("palette") or {})):
        entry = (theme.get("palette") or {}).get(key)
        if not isinstance(entry, dict):
            continue
        ref = entry.get("ref")
        if isinstance(ref, dict):
            colors.append((f"theme.palette.{key}", ref))
    return sizes, colors


def check_appearance_targets(shell, profile: dict) -> list[Finding]:
    """Verify every typography value the engine will APPLY is one the ``shell``
    actually proves it carries (font family, run size, run color).

    Typography capture records the template's dominant direct run font/size/color
    into ``role.appearance`` and the document defaults (``theme.fonts.body`` for
    font/size, ``theme.text.body`` for color); this is their shell-backed peer of
    :func:`check_resolver_targets`. Each axis re-validates against true PROVENANCE,
    not a sanity bound, and fails closed (ERROR) on any applied value the shell does
    not prove it contains - off-brand output stays impossible by construction:

      - FONT: must be in the shell's available faces (fontTable/theme latin + the
        Arial baseline; the theme declarations widen the ALLOWED set, never the
        checked set).
      - SIZE: the applied ``size_hp`` must be a run size actually present on the
        template's own runs.
      - COLOR (hex): the applied hex must be a theme-palette hex OR a hex actually
        present on the template's runs.
      - COLOR (theme token): the token's mapped ``clrScheme`` slot
        (:data:`_WML_THEME_TO_SLOT`) must be present in the parsed palette.

    Format-neutral: the per-kind :data:`_SHELL_APPEARANCE_COLLECTORS` reduce the shell
    to one :class:`ShellAppearanceFacts` shape (docx/pptx/xlsx) and the membership loops
    below are kind-agnostic. A collector that cannot parse fails CLOSED - it emits a
    WARNING and yields empty fact sets, so every applied value is then rejected (ERROR).

    A no-op when no appearance value is present (every pre-capture profile), when the
    shell is absent, or when the kind has no registered collector.
    """
    if shell is None:
        return []
    collector = _SHELL_APPEARANCE_COLLECTORS.get(profile.get("kind"))
    if collector is None:
        return []
    findings: list[Finding] = []
    try:
        facts = collector(shell, profile)
    except Exception as exc:  # opening the shell must never crash the gate
        # Fail closed: surface the parse failure and continue with empty fact sets, so
        # any applied font/size/color is rejected below rather than silently passing.
        findings.append(
            Finding(
                "appearance_targets_exist",
                schema.Severity.WARNING.value,
                f"could not verify appearance targets against shell: {exc}",
            )
        )
        facts = ShellAppearanceFacts(fonts=set(), sizes=set(), hexes=set(), palette={})
    available = facts.fonts
    applied: list[tuple[str, str]] = []
    for rid, entry in (profile.get("roles") or {}).items():
        if rid == "_index" or not isinstance(entry, dict):
            continue
        latin = ((entry.get("appearance") or {}).get("font") or {}).get("latin")
        if latin:
            applied.append((f"role {rid!r}", latin))
    body_latin = (
        ((profile.get("theme") or {}).get("fonts") or {}).get("body") or {}
    ).get("latin")
    if body_latin:
        applied.append(("theme.fonts.body", body_latin))

    for where, font in applied:
        if font not in available:
            findings.append(
                Finding(
                    "appearance_targets_exist",
                    schema.Severity.ERROR.value,
                    f"{where} font {font!r} is not available in the shell "
                    f"(have {sorted(available)})",
                )
            )

    applied_sizes, applied_colors = _collect_applied_appearance(profile)
    if applied_sizes or applied_colors:
        shell_sizes, shell_hexes, palette = facts.sizes, facts.hexes, facts.palette
        for where, size_hp in applied_sizes:
            if size_hp not in shell_sizes:
                findings.append(
                    Finding(
                        "appearance_targets_exist",
                        schema.Severity.ERROR.value,
                        f"{where} size {size_hp} (half-points) is not present on any "
                        f"run in the shell (have {sorted(shell_sizes)})",
                    )
                )
        palette_hexes = {colorutil.normalize_hex(h) for h in palette.values()}
        allowed_hexes = palette_hexes | shell_hexes
        for where, color in applied_colors:
            kind = color.get("kind")
            if kind == "hex":
                try:
                    hexval = colorutil.normalize_hex(color.get("hex") or "")
                except (ValueError, AttributeError):
                    hexval = None
                if hexval is None or hexval not in allowed_hexes:
                    findings.append(
                        Finding(
                            "appearance_targets_exist",
                            schema.Severity.ERROR.value,
                            f"{where} color #{color.get('hex')!r} is not in the shell "
                            f"palette or its runs (allowed {sorted(allowed_hexes)})",
                        )
                    )
            elif kind == "theme":
                token = color.get("theme")
                # A theme ref may name the color in EITHER namespace and both must
                # validate against the parsed clrScheme palette:
                #   - a WML ``themeColor`` token (``text1`` / ``background1`` /
                #     ``hyperlink`` ...) captured off a run, mapped to its
                #     clrScheme slot via _WML_THEME_TO_SLOT;
                #   - a clrScheme slot name directly (``dk1`` / ``lt1`` / ``hlink``
                #     ... or ``accent1``), which ``theme.palette`` keys/seeds use
                #     verbatim (it has no WML alias for the dk/lt/hlink slots).
                # The token is present when it resolves to a parsed-palette slot
                # through either route; only a token that does neither fails closed.
                slot = _WML_THEME_TO_SLOT.get(token)
                present = (slot is not None and slot in palette) or (
                    token in colorutil.THEME_SLOTS and token in palette
                )
                if not present:
                    findings.append(
                        Finding(
                            "appearance_targets_exist",
                            schema.Severity.ERROR.value,
                            f"{where} theme color {token!r} is absent from the shell "
                            f"palette (have {sorted(palette)})",
                        )
                    )
    return findings


# ---------------------------------------------------------------------------
# Paragraph GEOMETRY check (Cluster D1, DOCX-ONLY).
# ---------------------------------------------------------------------------
# Geometry is captured NUMBERS (twips) and XML ELEMENTS (borders / shading), NOT
# symbolic refs into a shell inventory like a font name or a style id. So this check is
# NOT a membership-against-a-named-inventory like ``check_appearance_targets`` - it is
# the HONEST geometry analogue:
#   (1) SHAPE / SANITY: every captured value is well-formed - twips are integers in the
#       OOXML range, line_rule/border serialization are structurally valid, a shading
#       fill is a real ``RRGGBB`` hex;
#   (2) OBSERVED-FLOOR MEMBERSHIP: every applied geometry value is one the TEMPLATE'S
#       OWN paragraphs carried (byte-identical match against the shell's observed
#       ``w:pPr`` facts) - the captured floor is the template's own geometry, never a
#       synthesized value.
# Fail-closed: a malformed / out-of-range / un-observed (synthesized) value is ERROR;
# only the template's own observed geometry passes. No-op for any non-docx kind, an
# absent shell, or a profile with no captured geometry.

# OOXML ``w:spacing@w:before/@w:after/@w:line`` are ``ST_TwipsMeasure`` (unsigned);
# ``w:ind`` left/right are ``ST_SignedTwipsMeasure`` (may be negative); firstLine /
# hanging are ``ST_TwipsMeasure`` (unsigned). The page-dimension ceiling is 31680
# twips (22"); we use a generous symmetric sanity bound so a structurally-impossible
# value (a synthesized 999999) is rejected on SHAPE before the observed-floor check
# even runs. This is a structural floor, NOT a brand bound - the real gate is the
# observed-floor membership below.
_GEOMETRY_TWIPS_MIN = -31680
_GEOMETRY_TWIPS_MAX = 31680

# The geometry scalar (group, field) -> whether the field is a twips measure (vs the
# ``line_rule`` string). Mirrors common.typography._GEOMETRY_SCALAR_FIELDS.
_GEOMETRY_TWIPS_FIELDS: tuple[tuple[str, str], ...] = (
    ("spacing", "before_twips"),
    ("spacing", "after_twips"),
    ("spacing", "line_twips"),
    ("indentation", "left_twips"),
    ("indentation", "right_twips"),
    ("indentation", "first_line_twips"),
    ("indentation", "hanging_twips"),
)
_GEOMETRY_BORDER_SIDES: tuple[str, ...] = ("top", "bottom", "left", "right")
# The spec-fixed WordprocessingML line-rule tokens (closed vocabulary).
_GEOMETRY_LINE_RULES: frozenset[str] = frozenset({"auto", "exact", "atLeast"})


class ShellGeometryFacts(NamedTuple):
    """The paragraph-geometry facts a docx shell proves it carries on its OWN
    paragraphs, the OBSERVED FLOOR an applied geometry value is validated against:

      - ``twips``: ``{(group, field): set[int]}`` - every explicit twips value the
        shell's own ``w:pPr`` carry, per scalar field;
      - ``line_rules``: the set of ``w:spacing@w:lineRule`` tokens observed;
      - ``borders``: ``{side: set[str]}`` of serialized ``w:pBdr`` side elements
        observed (byte-identical copies);
      - ``shading``: the set of observed normalized ``w:shd@w:fill`` hexes.
    """

    twips: dict
    line_rules: set
    borders: dict
    shading: set


def _docx_collect_geometry_facts(shell) -> ShellGeometryFacts:
    """Read the docx shell's OWN paragraph geometry (every ``w:pPr`` in
    ``word/document.xml``) into the observed floor. A missing/garbage document part
    yields empty sets - the caller then fails closed on any applied geometry."""
    from lxml import etree as _etree

    twips: dict = {field: set() for field in _GEOMETRY_TWIPS_FIELDS}
    line_rules: set = set()
    borders: dict = {side: set() for side in _GEOMETRY_BORDER_SIDES}
    shading: set = set()
    try:
        xml = pack.read_part(shell, "word/document.xml")
    except KeyError:
        return ShellGeometryFacts(
            twips=twips, line_rules=line_rules, borders=borders, shading=shading
        )
    root = pack.parse_xml_bytes(xml)

    def _twips(el, attr):
        val = el.get(_W(attr))
        if val is None:
            return None
        try:
            return int(val)
        except (TypeError, ValueError):
            return None

    for ppr in root.iter(_W("pPr")):
        spacing = ppr.find(_W("spacing"))
        if spacing is not None:
            for field, attr in (
                (("spacing", "before_twips"), "before"),
                (("spacing", "after_twips"), "after"),
                (("spacing", "line_twips"), "line"),
            ):
                v = _twips(spacing, attr)
                if v is not None:
                    twips[field].add(v)
            rule = spacing.get(_W("lineRule"))
            if rule:
                line_rules.add(rule)
        ind = ppr.find(_W("ind"))
        if ind is not None:
            for field, attr in (
                (("indentation", "left_twips"), "left"),
                (("indentation", "right_twips"), "right"),
                (("indentation", "first_line_twips"), "firstLine"),
                (("indentation", "hanging_twips"), "hanging"),
            ):
                v = _twips(ind, attr)
                if v is not None:
                    twips[field].add(v)
        pbdr = ppr.find(_W("pBdr"))
        if pbdr is not None:
            for side in _GEOMETRY_BORDER_SIDES:
                el = pbdr.find(_W(side))
                if el is not None:
                    try:
                        borders[side].add(_etree.tostring(el, encoding="unicode"))
                    except Exception:
                        continue
        shd = ppr.find(_W("shd"))
        if shd is not None:
            fill = shd.get(_W("fill"))
            if fill and fill.lower() != "auto":
                try:
                    shading.add(colorutil.normalize_hex(fill))
                except (ValueError, AttributeError):
                    continue
    return ShellGeometryFacts(
        twips=twips, line_rules=line_rules, borders=borders, shading=shading
    )


def _collect_applied_geometry(profile: dict) -> list:
    """Gather every ``(where, geometry-dict)`` the engine will APPLY: each role's
    ``appearance.geometry`` and the document body geometry (``theme.geometry.body``).
    Sorted by ``where`` for a deterministic finding order."""
    applied: list = []
    for rid, entry in (profile.get("roles") or {}).items():
        if rid == "_index" or not isinstance(entry, dict):
            continue
        geometry = (entry.get("appearance") or {}).get("geometry")
        if isinstance(geometry, dict) and geometry:
            applied.append((f"role {rid!r}", geometry))
    body = ((profile.get("theme") or {}).get("geometry") or {}).get("body")
    if isinstance(body, dict) and body:
        applied.append(("theme.geometry.body", body))
    applied.sort(key=lambda item: item[0])
    return applied


def _geometry_border_is_valid(serialized) -> bool:
    """A captured border side is well-formed: a serialized ``w:top/.../w:right`` element
    whose required ``w:val`` (the border style) is present. A non-string, unparseable,
    or attribute-missing copy is rejected (fail-closed on a synthesized border)."""
    if not isinstance(serialized, str) or not serialized:
        return False
    try:
        el = pack.parse_xml_bytes(serialized.encode("utf-8"))
    except Exception:
        return False
    # A WordprocessingML border element requires ``w:val`` (the border line style).
    return el.get(_W("val")) is not None


def check_geometry_targets(shell, profile: dict) -> list[Finding]:
    """Verify every paragraph-geometry value the engine will APPLY is (1) WELL-FORMED
    and (2) a value the TEMPLATE'S OWN paragraphs carried (the captured floor) - the
    honest fail-closed peer of :func:`check_appearance_targets` for the geometry axis
    (Cluster D1, DOCX-ONLY).

    Geometry is captured NUMBERS and XML ELEMENTS, not symbolic refs into a shell
    inventory, so this is NOT a name-membership check. It proves instead:

      - SHAPE / SANITY: spacing/indent twips are integers in the OOXML range; a
        ``line_rule`` is a spec token; a border side is a structurally valid element
        carrying its required ``w:val``; a shading fill is a real ``RRGGBB`` hex.
      - OBSERVED FLOOR: every applied twips/border/shading value is byte-identical to a
        value the shell's own ``w:pPr`` carried (parsed by
        :func:`_docx_collect_geometry_facts`), so a captured value is only ever the
        template's own observed geometry - NEVER a synthesized value.

    Fail-closed: a malformed / out-of-range / un-observed value is ERROR. A no-op when
    the kind is not docx, the shell is absent, or no geometry is captured (every
    pre-D1 profile). A shell that cannot be parsed fails CLOSED (a WARNING plus empty
    observed sets, so every applied value is then rejected)."""
    if shell is None or profile.get("kind") != schema.Kind.DOCX.value:
        return []
    applied = _collect_applied_geometry(profile)
    if not applied:
        return []
    findings: list[Finding] = []
    try:
        facts = _docx_collect_geometry_facts(shell)
    except Exception as exc:  # opening the shell must never crash the gate
        findings.append(
            Finding(
                "appearance_geometry_targets",
                schema.Severity.WARNING.value,
                f"could not verify geometry targets against shell: {exc}",
            )
        )
        facts = ShellGeometryFacts(
            twips={field: set() for field in _GEOMETRY_TWIPS_FIELDS},
            line_rules=set(),
            borders={side: set() for side in _GEOMETRY_BORDER_SIDES},
            shading=set(),
        )

    def _err(where: str, msg: str) -> None:
        findings.append(
            Finding(
                "appearance_geometry_targets",
                schema.Severity.ERROR.value,
                msg,
                location=where,
            )
        )

    for where, geometry in applied:
        spacing = geometry.get("spacing") or {}
        indentation = geometry.get("indentation") or {}
        # (1) scalar twips: shape (int + in range) then observed-floor membership.
        for group, field in _GEOMETRY_TWIPS_FIELDS:
            sub = spacing if group == "spacing" else indentation
            if field not in sub:
                continue
            value = sub[field]
            if not isinstance(value, int) or isinstance(value, bool):
                _err(
                    where,
                    f"{where} geometry {group}.{field} {value!r} is not an integer twips value",
                )
                continue
            if not (_GEOMETRY_TWIPS_MIN <= value <= _GEOMETRY_TWIPS_MAX):
                _err(
                    where,
                    f"{where} geometry {group}.{field} {value} twips is out of the sane "
                    f"OOXML range [{_GEOMETRY_TWIPS_MIN}, {_GEOMETRY_TWIPS_MAX}]",
                )
                continue
            observed = facts.twips.get((group, field), set())
            if value not in observed:
                _err(
                    where,
                    f"{where} geometry {group}.{field} {value} twips is not observed on any "
                    f"paragraph in the shell (have {sorted(observed)})",
                )
        # line_rule: a spec token observed on the shell.
        line_rule = spacing.get("line_rule")
        if line_rule is not None:
            if line_rule not in _GEOMETRY_LINE_RULES:
                _err(
                    where,
                    f"{where} geometry spacing.line_rule {line_rule!r} is not a WordprocessingML line rule",
                )
            elif line_rule not in facts.line_rules:
                _err(
                    where,
                    f"{where} geometry spacing.line_rule {line_rule!r} is not observed in the shell "
                    f"(have {sorted(facts.line_rules)})",
                )
        # (2) borders: each side well-formed AND byte-identical to an observed side.
        borders = geometry.get("borders") or {}
        for side, serialized in borders.items():
            if side not in _GEOMETRY_BORDER_SIDES:
                _err(where, f"{where} geometry borders has unknown side {side!r}")
                continue
            if not _geometry_border_is_valid(serialized):
                _err(
                    where,
                    f"{where} geometry borders.{side} is not a valid WordprocessingML border element",
                )
                continue
            if serialized not in facts.borders.get(side, set()):
                _err(
                    where,
                    f"{where} geometry borders.{side} is not byte-identical to any border the shell "
                    "carries (synthesized geometry rejected)",
                )
        # (3) shading: a real hex AND observed on the shell.
        shading = geometry.get("shading") or {}
        fill = shading.get("fill_hex")
        if fill is not None:
            try:
                normalized = colorutil.normalize_hex(fill)
            except (ValueError, AttributeError):
                normalized = None
            if normalized is None:
                _err(
                    where,
                    f"{where} geometry shading.fill_hex {fill!r} is not a valid #RRGGBB hex",
                )
            elif normalized not in facts.shading:
                _err(
                    where,
                    f"{where} geometry shading.fill_hex #{normalized} is not observed on any paragraph "
                    f"in the shell (have {sorted(facts.shading)})",
                )
    return findings


# ---------------------------------------------------------------------------
# Cluster D2: TABLE conditional-format targets (DOCX-ONLY).
# The honest fail-closed peer of check_geometry_targets for the table axis. It proves
# THREE independent dimensions: the tblLook bitmask is WELL-FORMED (shape/sanity; the
# flags are spec-fixed bits, not template-derived names), the referenced table STYLE is
# a SYMBOLIC ref that the shell's styles part actually defines (name-membership, like
# check_appearance_targets does for fonts), and each cell margin is an intrinsic NUMBER
# (twips) the template's OWN tables carried (observed-floor, like geometry).
# ---------------------------------------------------------------------------

# The valid OR-able ``w:tblLook`` flag bits (the only bits the bitmask may set). Mirrors
# ``formats/docx/typography._TBLLOOK_FLAG_BITS``: firstRow / lastRow / firstColumn /
# lastColumn / noHBand / noVBand. A bitmask with any other bit set is malformed.
_TBLLOOK_VALID_BITS = 0x0020 | 0x0040 | 0x0080 | 0x0100 | 0x0200 | 0x0400
# The OOXML 16-bit ``w:tblLook@w:val`` range.
_TBLLOOK_MIN = 0
_TBLLOOK_MAX = 0xFFFF
# The four ``w:tblCellMar`` margin fields (the same twips range as paragraph geometry).
_TABLE_CELL_MARGIN_FIELDS: tuple[str, ...] = (
    "top_twips",
    "bottom_twips",
    "left_twips",
    "right_twips",
)
_TABLE_CELL_MARGIN_SIDE: dict[str, str] = {
    "top_twips": "top",
    "bottom_twips": "bottom",
    "left_twips": "left",
    "right_twips": "right",
}


class ShellTableFacts(NamedTuple):
    """The table facts a docx shell proves it carries on its OWN tables, the OBSERVED
    FLOOR / membership inventory an applied table value is validated against:

      - ``style_ids``: the ``@w:styleId`` of every ``w:style[@w:type='table']`` the
        shell declares (the symbolic table-style inventory a referenced style must be
        in);
      - ``cell_margins``: ``{field: set[int]}`` of every ``w:tblCellMar`` side twips the
        shell's own tables carry, per side.
    """

    style_ids: set
    cell_margins: dict


def _docx_table_style_ids(shell) -> set:
    """The ``@w:styleId`` (and ``w:name@w:val``) of every ``w:style[@w:type='table']``
    the shell's ``word/styles.xml`` declares. A referenced table style must be a member
    of this set (fail-closed name-membership, like fonts). A missing styles part yields
    an empty set (every applied style id then fails closed)."""
    ids: set = set()
    try:
        xml = pack.read_part(shell, "word/styles.xml")
    except KeyError:
        return ids
    root = pack.parse_xml_bytes(xml)
    for style in root.findall(_W("style")):
        if style.get(_W("type")) != "table":
            continue
        sid = style.get(_W("styleId"))
        if sid:
            ids.add(sid)
        name_el = style.find(_W("name"))
        if name_el is not None:
            name = name_el.get(_W("val"))
            if name:
                ids.add(name)
    return ids


def _docx_collect_table_facts(shell) -> ShellTableFacts:
    """Read the docx shell's OWN table facts: the table-style inventory (every
    ``w:style[@w:type='table']`` id/name) and the observed cell-margin floor (every
    ``w:tblPr/w:tblCellMar`` side twips the shell's tables carry). A missing/garbage
    document part yields empty margin sets - the caller then fails closed on any applied
    margin."""
    style_ids = _docx_table_style_ids(shell)
    cell_margins: dict = {field: set() for field in _TABLE_CELL_MARGIN_FIELDS}
    try:
        xml = pack.read_part(shell, "word/document.xml")
    except KeyError:
        return ShellTableFacts(style_ids=style_ids, cell_margins=cell_margins)
    root = pack.parse_xml_bytes(xml)
    for tblpr in root.iter(_W("tblPr")):
        cell_mar = tblpr.find(_W("tblCellMar"))
        if cell_mar is None:
            continue
        for field in _TABLE_CELL_MARGIN_FIELDS:
            side = _TABLE_CELL_MARGIN_SIDE[field]
            el = cell_mar.find(_W(side))
            if el is None:
                continue
            val = el.get(_W("w"))
            if val is None:
                continue
            try:
                cell_margins[field].add(int(val))
            except (TypeError, ValueError):
                continue
    return ShellTableFacts(style_ids=style_ids, cell_margins=cell_margins)


def _collect_applied_table(profile: dict) -> list:
    """Gather every ``(where, table-dict)`` the engine will APPLY: each role's
    ``appearance.table`` and the document body table default (``theme.table.body``).
    Sorted by ``where`` for a deterministic finding order."""
    applied: list = []
    for rid, entry in (profile.get("roles") or {}).items():
        if rid == "_index" or not isinstance(entry, dict):
            continue
        table = (entry.get("appearance") or {}).get("table")
        if isinstance(table, dict) and table:
            applied.append((f"role {rid!r}", table))
    body = ((profile.get("theme") or {}).get("table") or {}).get("body")
    if isinstance(body, dict) and body:
        applied.append(("theme.table.body", body))
    applied.sort(key=lambda item: item[0])
    return applied


def check_table_targets(shell, profile: dict) -> list[Finding]:
    """Verify every TABLE conditional-format value the engine will APPLY is well-formed
    and shell-backed - the honest fail-closed peer of :func:`check_geometry_targets` for
    the table axis (Cluster D2, DOCX-ONLY). It validates THREE INDEPENDENT dimensions:

      - tblLook SHAPE / SANITY: the bitmask is an integer in the 16-bit OOXML range
        ``[0, 0xFFFF]`` whose only set bits are the spec-fixed flags (firstRow / lastRow
        / firstColumn / lastColumn / noHBand / noVBand). The flags are spec-fixed bits
        (NOT template-derived names), so SHAPE is sufficient - no membership.
      - TABLE-STYLE REFERENCE MEMBERSHIP: the referenced ``style_id`` is the
        ``@w:styleId`` (or ``w:name``) of a ``w:style[@w:type='table']`` the shell's
        styles part actually declares. The style's ``w:tblStylePr`` conditional formats
        (the band fills / first-last emphasis) are AUTHORED IN THE SHELL - the engine
        only references the style; an undefined style is a brand breach (ERROR).
      - CELL-MARGINS SHAPE / OBSERVED-FLOOR: each margin twip is an integer in the OOXML
        range AND byte-identical to a margin the shell's OWN tables carried. Margins are
        intrinsic NUMBERS, so the floor is template-observed (like geometry), not a
        symbolic inventory.

    Fail-closed: a malformed bitmask / undefined style / malformed-or-un-observed margin
    is ERROR. A no-op when the kind is not docx, the shell is absent, or no table
    appearance is captured (every pre-D2 profile). A shell that cannot be parsed fails
    CLOSED (a WARNING plus empty inventories, so every applied value is then rejected)."""
    if shell is None or profile.get("kind") != schema.Kind.DOCX.value:
        return []
    applied = _collect_applied_table(profile)
    if not applied:
        return []
    findings: list[Finding] = []
    try:
        facts = _docx_collect_table_facts(shell)
    except Exception as exc:  # opening the shell must never crash the gate
        findings.append(
            Finding(
                "appearance_table_targets",
                schema.Severity.WARNING.value,
                f"could not verify table targets against shell: {exc}",
            )
        )
        facts = ShellTableFacts(
            style_ids=set(),
            cell_margins={field: set() for field in _TABLE_CELL_MARGIN_FIELDS},
        )

    def _err(where: str, msg: str) -> None:
        findings.append(
            Finding(
                "appearance_table_targets",
                schema.Severity.ERROR.value,
                msg,
                location=where,
            )
        )

    for where, table in applied:
        # (1) tblLook: SHAPE / SANITY - int, in 16-bit range, only spec-fixed flag bits.
        tbllook = table.get("tblLook")
        if tbllook is not None:
            if not isinstance(tbllook, int) or isinstance(tbllook, bool):
                _err(
                    where,
                    f"{where} table tblLook {tbllook!r} is not an integer bitmask",
                )
            elif not (_TBLLOOK_MIN <= tbllook <= _TBLLOOK_MAX):
                _err(
                    where,
                    f"{where} table tblLook {tbllook} is out of the OOXML 16-bit range "
                    f"[{_TBLLOOK_MIN}, {_TBLLOOK_MAX}]",
                )
            elif tbllook & ~_TBLLOOK_VALID_BITS:
                _err(
                    where,
                    f"{where} table tblLook {tbllook:#06x} sets a bit outside the valid "
                    f"flags (firstRow/lastRow/firstColumn/lastColumn/noHBand/noVBand)",
                )
        # (2) style_id: SYMBOLIC name-membership against the shell's table styles.
        style_id = table.get("style_id")
        if style_id is not None:
            if not isinstance(style_id, str) or not style_id:
                _err(
                    where,
                    f"{where} table style_id {style_id!r} is not a style reference",
                )
            elif style_id not in facts.style_ids:
                _err(
                    where,
                    f"{where} table style {style_id!r} is not a table style the shell "
                    f"defines (have {sorted(facts.style_ids)})",
                )
        # (3) cell_margins: SHAPE (int + range) then OBSERVED-FLOOR membership.
        cell_margins = table.get("cell_margins") or {}
        for field in _TABLE_CELL_MARGIN_FIELDS:
            if field not in cell_margins:
                continue
            value = cell_margins[field]
            if not isinstance(value, int) or isinstance(value, bool):
                _err(
                    where,
                    f"{where} table cell_margins.{field} {value!r} is not an integer "
                    "twips value",
                )
                continue
            if not (_GEOMETRY_TWIPS_MIN <= value <= _GEOMETRY_TWIPS_MAX):
                _err(
                    where,
                    f"{where} table cell_margins.{field} {value} twips is out of the "
                    f"sane OOXML range [{_GEOMETRY_TWIPS_MIN}, {_GEOMETRY_TWIPS_MAX}]",
                )
                continue
            observed = facts.cell_margins.get(field, set())
            if value not in observed:
                _err(
                    where,
                    f"{where} table cell_margins.{field} {value} twips is not observed "
                    f"on any table in the shell (have {sorted(observed)})",
                )
    return findings


# ---------------------------------------------------------------------------
# Cluster D3: LIST / NUMBERING definition targets (DOCX-ONLY).
# The honest fail-closed peer of check_geometry_targets / check_table_targets for the
# numbering axis. The numbering DEFINITION (w:abstractNum / w:num) is OWNED BY THE SHELL:
# the engine only REFERENCES it by id and at most CLONES the shell's own w:abstractNum -
# it NEVER synthesizes a numFmt / lvlText / indent from JSON. This check proves:
#   (1) num_id membership: the referenced w:numId is in the shell's w:num inventory;
#   (2) abstract_num_id membership: the resolved w:abstractNumId is in the shell's
#       w:abstractNum inventory (the def the engine clones by id);
#   (3) numFmt SHAPE: the per-level numFmt is a valid OOXML field code (closed enum);
#   (4) lvlText OBSERVED-FLOOR: the per-level lvlText byte-matches the shell's OWN
#       w:lvlText for that ilvl in the referenced abstractNum (never synthesized);
#   (5) indent OBSERVED-FLOOR: each per-level indent twip is an int in the OOXML range
#       AND byte-matches the shell's OWN w:ind value for that ilvl (never synthesized).
# Fail-closed: any undefined id / malformed numFmt / un-observed lvlText / out-of-range or
# un-observed indent is ERROR. A no-op for non-docx kinds and pre-D3 profiles.
# ---------------------------------------------------------------------------

# The closed set of OOXML ``w:numFmt@w:val`` field codes (ECMA-376 ST_NumberFormat). A
# captured numFmt MUST be one of these (SHAPE-only; the floor proper is the lvlText/indent
# observed-floor). This is a spec-fixed enum, NOT a template-derived inventory.
_NUMFMT_VALID_CODES: frozenset[str] = frozenset(
    {
        "decimal",
        "upperRoman",
        "lowerRoman",
        "upperLetter",
        "lowerLetter",
        "ordinal",
        "cardinalText",
        "ordinalText",
        "hex",
        "chicago",
        "ideographDigital",
        "japaneseCounting",
        "aiueo",
        "iroha",
        "decimalFullWidth",
        "decimalHalfWidth",
        "japaneseLegal",
        "japaneseDigitalTenThousand",
        "decimalEnclosedCircle",
        "decimalFullWidth2",
        "aiueoFullWidth",
        "irohaFullWidth",
        "decimalZero",
        "bullet",
        "ganada",
        "chosung",
        "decimalEnclosedFullstop",
        "decimalEnclosedParen",
        "decimalEnclosedCircleChinese",
        "ideographEnclosedCircle",
        "ideographTraditional",
        "ideographZodiac",
        "ideographZodiacTraditional",
        "taiwaneseCounting",
        "ideographLegalTraditional",
        "taiwaneseCountingThousand",
        "taiwaneseDigital",
        "chineseCounting",
        "chineseLegalSimplified",
        "chineseCountingThousand",
        "koreanDigital",
        "koreanCounting",
        "koreanLegal",
        "koreanDigital2",
        "vietnameseCounting",
        "russianLower",
        "russianUpper",
        "none",
        "numberInDash",
        "hebrew1",
        "hebrew2",
        "arabicAlpha",
        "arabicAbjad",
        "hindiVowels",
        "hindiConsonants",
        "hindiNumbers",
        "hindiCounting",
        "thaiLetters",
        "thaiNumbers",
        "thaiCounting",
        "bahtText",
        "dollarText",
        "custom",
    }
)
# The per-level indent attributes the check validates (mirrors the capture reader).
_NUM_INDENT_ATTRS: tuple[str, ...] = ("left", "right", "firstLine", "hanging")


class ShellNumberingFacts(NamedTuple):
    """The numbering facts a docx shell proves it carries, the membership inventory +
    observed floor an applied numbering value is validated against:

      - ``num_ids``: the ``@w:numId`` of every ``w:num`` the shell declares (the symbolic
        reference a captured ``num_id`` must be a member of);
      - ``abstract_num_ids``: the ``@w:abstractNumId`` of every ``w:abstractNum`` the
        shell declares (the def a captured ``abstract_num_id`` must be a member of);
      - ``per_level``: ``{abstract_num_id -> {ilvl -> {lvlText, indent}}}`` of the shell's
        OWN declared per-level ``w:lvlText`` / ``w:ind`` facts (the observed floor).
    """

    num_ids: set
    abstract_num_ids: set
    per_level: dict


def _docx_collect_numbering_facts(shell) -> ShellNumberingFacts:
    """Read the docx shell's OWN numbering facts from ``word/numbering.xml``: the w:num /
    w:abstractNum id inventories and the per-level lvlText/indent observed floor. A
    missing/garbage numbering part yields empty inventories - the caller then fails closed
    on every applied numbering value (a referenced id can never be a member of {})."""
    num_ids: set = set()
    abstract_num_ids: set = set()
    per_level: dict = {}
    try:
        xml = pack.read_part(shell, "word/numbering.xml")
    except KeyError:
        return ShellNumberingFacts(
            num_ids=num_ids, abstract_num_ids=abstract_num_ids, per_level=per_level
        )
    root = pack.parse_xml_bytes(xml)
    for num in root.findall(_W("num")):
        nid = num.get(_W("numId"))
        if nid is not None:
            num_ids.add(str(nid))
    for an in root.findall(_W("abstractNum")):
        aid = an.get(_W("abstractNumId"))
        if aid is None:
            continue
        aid = str(aid)
        abstract_num_ids.add(aid)
        levels: dict = {}
        for lvl in an.findall(_W("lvl")):
            try:
                ilvl = int(lvl.get(_W("ilvl")) or 0)
            except (TypeError, ValueError):
                continue
            facts: dict = {}
            lt = lvl.find(_W("lvlText"))
            if lt is not None and lt.get(_W("val")) is not None:
                facts["lvlText"] = lt.get(_W("val"))
            ppr = lvl.find(_W("pPr"))
            ind = ppr.find(_W("ind")) if ppr is not None else None
            if ind is not None:
                indent: dict = {}
                for attr in _NUM_INDENT_ATTRS:
                    val = ind.get(_W(attr))
                    if val is None:
                        continue
                    try:
                        indent[attr] = int(val)
                    except (TypeError, ValueError):
                        continue
                if indent:
                    facts["indent"] = indent
            levels[ilvl] = facts
        per_level[aid] = levels
    return ShellNumberingFacts(
        num_ids=num_ids, abstract_num_ids=abstract_num_ids, per_level=per_level
    )


def _collect_applied_numbering(profile: dict) -> list:
    """Gather every ``(where, numbering-dict)`` the engine will APPLY: each role's
    ``appearance.numbering`` and the document body numbering default
    (``theme.numbering.body``). Sorted by ``where`` for a deterministic finding order."""
    applied: list = []
    for rid, entry in (profile.get("roles") or {}).items():
        if rid == "_index" or not isinstance(entry, dict):
            continue
        numbering = (entry.get("appearance") or {}).get("numbering")
        if isinstance(numbering, dict) and numbering:
            applied.append((f"role {rid!r}", numbering))
    body = ((profile.get("theme") or {}).get("numbering") or {}).get("body")
    if isinstance(body, dict) and body:
        applied.append(("theme.numbering.body", body))
    applied.sort(key=lambda item: item[0])
    return applied


def check_numbering_targets(shell, profile: dict) -> list[Finding]:
    """Verify every LIST / NUMBERING value the engine will APPLY is shell-backed - the
    honest fail-closed peer of :func:`check_geometry_targets` / :func:`check_table_targets`
    for the numbering axis (Cluster D3, DOCX-ONLY). The numbering DEFINITION is owned by
    the shell; the engine only REFERENCES it by id and CLONES the shell's own
    ``w:abstractNum``. It validates:

      - NUM-ID MEMBERSHIP: the referenced ``num_id`` is a ``w:num`` the shell declares.
      - ABSTRACT-NUM-ID MEMBERSHIP: the resolved ``abstract_num_id`` is a ``w:abstractNum``
        the shell declares (the def the engine clones by id).
      - numFmt SHAPE: each per-level ``numFmt`` is a valid OOXML field code (closed enum;
        SHAPE only - the field code is spec-fixed, not template-derived).
      - lvlText OBSERVED-FLOOR: each per-level ``lvlText`` byte-matches the shell's OWN
        ``w:lvlText`` for that ilvl in the referenced abstractNum (never synthesized).
      - indent OBSERVED-FLOOR: each per-level indent twip is an int in the OOXML range AND
        byte-matches the shell's OWN ``w:ind`` value for that ilvl (never synthesized).

    Fail-closed: an undefined id / malformed numFmt / un-observed lvlText / malformed or
    un-observed indent is ERROR. A no-op when the kind is not docx, the shell is absent, or
    no numbering is captured (every pre-D3 profile). A shell that cannot be parsed fails
    CLOSED (a WARNING plus empty inventories, so every applied value is then rejected)."""
    if shell is None or profile.get("kind") != schema.Kind.DOCX.value:
        return []
    applied = _collect_applied_numbering(profile)
    if not applied:
        return []
    findings: list[Finding] = []
    try:
        facts = _docx_collect_numbering_facts(shell)
    except Exception as exc:  # opening the shell must never crash the gate
        findings.append(
            Finding(
                "appearance_numbering_targets",
                schema.Severity.WARNING.value,
                f"could not verify numbering targets against shell: {exc}",
            )
        )
        facts = ShellNumberingFacts(num_ids=set(), abstract_num_ids=set(), per_level={})

    def _err(where: str, msg: str) -> None:
        findings.append(
            Finding(
                "appearance_numbering_targets",
                schema.Severity.ERROR.value,
                msg,
                location=where,
            )
        )

    for where, numbering in applied:
        # (1) num_id: SYMBOLIC membership against the shell's w:num inventory.
        num_id = numbering.get("num_id")
        if num_id is None or not str(num_id):
            _err(
                where,
                f"{where} numbering num_id {num_id!r} is not a numbering reference",
            )
        elif str(num_id) not in facts.num_ids:
            _err(
                where,
                f"{where} numbering num_id {num_id!r} is not a w:num the shell defines "
                f"(have {sorted(facts.num_ids)})",
            )
        # (2) abstract_num_id: SYMBOLIC membership against the shell's w:abstractNum.
        abstract_num_id = numbering.get("abstract_num_id")
        if abstract_num_id is None or not str(abstract_num_id):
            _err(
                where,
                f"{where} numbering abstract_num_id {abstract_num_id!r} is not a "
                "numbering-definition reference",
            )
            continue  # without a resolvable def, the per-level floor cannot be checked
        if str(abstract_num_id) not in facts.abstract_num_ids:
            _err(
                where,
                f"{where} numbering abstract_num_id {abstract_num_id!r} is not a "
                f"w:abstractNum the shell defines (have {sorted(facts.abstract_num_ids)})",
            )
            continue
        shell_levels = facts.per_level.get(str(abstract_num_id), {})
        # (3)-(5) per-level facts: numFmt shape, lvlText + indent observed-floor.
        per_level = numbering.get("per_level_facts") or {}
        if not isinstance(per_level, dict):
            _err(where, f"{where} numbering per_level_facts is not a mapping")
            continue
        for raw_ilvl, level_facts in per_level.items():
            try:
                ilvl = int(raw_ilvl)
            except (TypeError, ValueError):
                _err(
                    where,
                    f"{where} numbering per_level_facts key {raw_ilvl!r} is not an "
                    "integer level",
                )
                continue
            if not isinstance(level_facts, dict):
                _err(
                    where,
                    f"{where} numbering per_level_facts[{ilvl}] is not a mapping",
                )
                continue
            shell_level = shell_levels.get(ilvl, {})
            # (3) numFmt SHAPE: a valid OOXML field code.
            numfmt = level_facts.get("numFmt")
            if numfmt is not None:
                if not isinstance(numfmt, str) or numfmt not in _NUMFMT_VALID_CODES:
                    _err(
                        where,
                        f"{where} numbering level {ilvl} numFmt {numfmt!r} is not a valid "
                        "OOXML field code",
                    )
            # (4) lvlText OBSERVED-FLOOR: byte-identical to the shell's own lvlText.
            lvltext = level_facts.get("lvlText")
            if lvltext is not None:
                if not isinstance(lvltext, str):
                    _err(
                        where,
                        f"{where} numbering level {ilvl} lvlText {lvltext!r} is not a string",
                    )
                elif "lvlText" not in shell_level or shell_level["lvlText"] != lvltext:
                    _err(
                        where,
                        f"{where} numbering level {ilvl} lvlText {lvltext!r} is not "
                        "byte-identical to the shell's own lvlText for that level "
                        "(synthesized numbering rejected)",
                    )
            # (5) indent OBSERVED-FLOOR: int + range, byte-identical to the shell's w:ind.
            indent = level_facts.get("indent") or {}
            if indent and not isinstance(indent, dict):
                _err(where, f"{where} numbering level {ilvl} indent is not a mapping")
                indent = {}
            shell_indent = shell_level.get("indent", {})
            for attr in _NUM_INDENT_ATTRS:
                if attr not in indent:
                    continue
                value = indent[attr]
                if not isinstance(value, int) or isinstance(value, bool):
                    _err(
                        where,
                        f"{where} numbering level {ilvl} indent.{attr} {value!r} is not an "
                        "integer twips value",
                    )
                    continue
                if not (_GEOMETRY_TWIPS_MIN <= value <= _GEOMETRY_TWIPS_MAX):
                    _err(
                        where,
                        f"{where} numbering level {ilvl} indent.{attr} {value} twips is out "
                        f"of the sane OOXML range [{_GEOMETRY_TWIPS_MIN}, "
                        f"{_GEOMETRY_TWIPS_MAX}]",
                    )
                    continue
                if attr not in shell_indent or shell_indent[attr] != value:
                    _err(
                        where,
                        f"{where} numbering level {ilvl} indent.{attr} {value} twips is not "
                        "byte-identical to the shell's own w:ind for that level "
                        "(synthesized numbering rejected)",
                    )
    return findings


def _pptx_layout_names(shell) -> set:
    prs = _load_pptx(shell)
    return {layout.name for layout in prs.slide_layouts}


def _xlsx_defined_names(shell) -> set:
    wb = _load_xlsx(shell)
    try:
        return set(wb.defined_names.keys())
    except AttributeError:
        # Older openpyxl exposes defined_names as a list-like of DefinedName.
        return {dn.name for dn in wb.defined_names}


def _xlsx_number_format_masks(shell) -> set:
    """The distinct number-format masks the workbook actually uses (``General`` dropped).

    Mirrors ``xlsx_structure.inventory_number_formats`` as a set, kept local so the
    QA layer needs no format-module import; iterates ``_cells`` (sparse-safe).
    """
    wb = _load_xlsx(shell)
    masks: set = set()
    for ws in wb.worksheets:
        for cell in ws._cells.values():
            if cell.value is None:
                continue
            fmt = cell.number_format
            if fmt and fmt != "General":
                masks.add(fmt)
    return masks


# ---------------------------------------------------------------------------
# Comprehension-aware deterministic checks (model-free; bind to captured facts)
# ---------------------------------------------------------------------------
def check_comprehension_targets(profile: dict) -> list[Finding]:
    """FAIL-CLOSED membership of every load-bearing comprehension ref (§3).

    Every ``anchor_ref`` / ``index_ref`` / ``region_ref`` / ``feeds_from_role_id``
    / ``role_annotations`` key must be a verbatim id from the surfaced inventories;
    a ref whose inventory is empty/absent is itself an ERROR (this is the SOLE gate
    for anchor/index/region refs, so it must reject, never skip). No-ops when the
    comprehension is absent (status != present), keeping the model-free CI path and
    pptx/xlsx green. Delegates to the single membership definition so the gate and
    the merge writer can never disagree.

    ``check_membership`` returns ALL membership problems, including the C1 audit-arm
    ``comprehension.audit`` keys. Those are attributed SOLELY to ``audit_targets_exist``
    by :func:`check_audit_targets`, so they are skipped here - else one bad audit key
    would double-report (one ``comprehension_targets_exist`` + one
    ``audit_targets_exist``). The skip is symmetric to ``check_audit_targets``' keep.
    """
    comp = profile.get("comprehension")
    if (
        not isinstance(comp, dict)
        or comp.get("status") != schema.ComprehensionStatus.PRESENT.value
    ):
        return []
    findings: list[Finding] = []
    for problem in comprehensionmod.check_membership(profile, comp):
        if problem.startswith("comprehension.audit"):
            continue  # attributed to audit_targets_exist (check_audit_targets)
        findings.append(
            Finding("comprehension_targets_exist", schema.Severity.ERROR.value, problem)
        )
    return findings


def check_audit_targets(profile: dict) -> list[Finding]:
    """FAIL-CLOSED membership of every persisted L2 visual-audit verdict key (C1).

    The audit peer of :func:`check_comprehension_targets`: every
    ``comprehension.audit`` key must be a verbatim id from the profile's derived
    visual checklist; a key not in that set (or any key when the derived checklist
    is empty) is an ERROR (this is the SOLE gate for audit keys, so it must reject,
    never skip). No-ops when the comprehension is absent (status != present),
    keeping the model-free CI path and pptx/xlsx green. Delegates to the single
    membership definition (``check_membership``'s audit arm, which itself binds to
    ``qa.visual.visual_checklist_ids``) so the gate and the merge writer can never
    disagree about which checklist ids exist.

    A surfaced finding here emits ``audit_targets_exist`` (the
    :data:`schema.DEFAULT_L0_INVARIANTS` id) rather than the
    ``comprehension_targets_exist`` id that ``check_membership`` problems carry under
    :func:`check_comprehension_targets`, so the verdict-key violation is attributable
    to its own invariant. ``check_membership`` returns ALL membership problems
    (anchor/index/region/role/palette + audit); we keep only the audit ones here so
    the two checks do not double-report the same anchor/index problem.
    """
    comp = profile.get("comprehension")
    if (
        not isinstance(comp, dict)
        or comp.get("status") != schema.ComprehensionStatus.PRESENT.value
    ):
        return []
    findings: list[Finding] = []
    for problem in comprehensionmod.check_membership(profile, comp):
        if problem.startswith("comprehension.audit"):
            findings.append(
                Finding("audit_targets_exist", schema.Severity.ERROR.value, problem)
            )
    return findings


def check_triage_targets(profile: dict) -> list[Finding]:
    """FAIL-CLOSED membership of every model-assisted QA-triage entry (Cluster C2).

    The triage peer of :func:`check_audit_targets`: every ``comprehension.triage``
    entry must name a check in the closed :data:`schema.AMBIGUOUS_TRIAGE_CHECKS` set,
    and each ``(check, location)`` pair must be unique across the proposal. A
    non-eligible check or a duplicate pair is an ERROR (this is the SOLE gate for
    triage entries, so it must reject, never skip). No-ops when the comprehension is
    absent (status != present), keeping the model-free CI path and pptx/xlsx green.
    Delegates to the single membership definition (``check_triage``) so the gate and
    the merge writer can never disagree about which checks are triage-eligible.

    Because the eligible set is WARNING-only, a triage entry can never even be aimed
    at an ERROR-emitting check - it would be rejected here as a non-member. That is
    the belt to ``qa.gate._apply_triage``'s ``severity == WARNING`` suspenders, so an
    ERROR can NEVER be demoted.
    """
    comp = profile.get("comprehension")
    if (
        not isinstance(comp, dict)
        or comp.get("status") != schema.ComprehensionStatus.PRESENT.value
    ):
        return []
    findings: list[Finding] = []
    for problem in comprehensionmod.check_triage(profile, comp):
        findings.append(
            Finding("triage_targets_exist", schema.Severity.ERROR.value, problem)
        )
    return findings


def _role_resolver_target_in_shell(kind, resolver: dict, present: set) -> bool:
    """True iff one role's resolver target is proven by the shell's ``present`` set.

    The membership predicate of ``check_override_targets``' reroute branch. It
    MIRRORS ``check_resolver_targets``' per-kind membership arms but is deliberately
    STRICTER, so it is not literally shared: ``check_resolver_targets`` SKIPS a role
    whose resolver is incomplete (a missing layout/name is another validator's
    problem), while a reroute target with the same incomplete resolver must be
    REJECTED here - a lesson may only re-point at a concrete, shell-backed op, never
    a stub. ``present`` is the kind-appropriate inventory (style keys / layout names
    / defined names+masks).
    """
    rtype = resolver.get("type")
    if kind == schema.Kind.DOCX.value:
        if rtype not in (schema.ResolverType.NAMED_STYLE.value, None):
            return False
        sid = resolver.get("style_id")
        sname = resolver.get("style_name")
        if not sid and not sname:
            return False
        return bool((sid and sid in present) or (sname and sname in present))
    if kind == schema.Kind.PPTX.value:
        if rtype != schema.ResolverType.PLACEHOLDER.value:
            return False
        layout = resolver.get("layout")
        return layout is not None and layout in present
    if kind == schema.Kind.XLSX.value:
        if rtype == schema.ResolverType.NAMED_RANGE.value:
            name = resolver.get("name")
            return name is not None and name in present
        if rtype == schema.ResolverType.NUMBER_FORMAT.value:
            mask = resolver.get("number_format")
            return mask is not None and mask in present
        return False
    return False


def check_override_targets(shell, profile: dict) -> list[Finding]:
    """FAIL-CLOSED membership of every LEARNED override target against the shell.

    The override peer of :func:`check_resolver_targets`, with the **reject-on-empty**
    discipline of :func:`check_comprehension_targets` (NOT the namespace-guarded
    skip-on-empty of ``_validate_resolver_consistency``): when a lesson is present,
    every re-point target must be proven by THIS shell, and a target whose inventory
    is empty/absent is itself an ERROR (this is the SOLE gate for override targets, so
    it must reject, never skip). Emits ERROR ``override_targets_exist`` when:

      - a ``reroute_role`` ``to`` is not a concrete declared role key (``_index`` is
        the order array, not a role), or that role's own resolver target is not
        shell-proven (so a reroute can only ever land on a real, shell-backed op);
      - a ``number_format`` swap mask is not among the shell's used masks
        (``surface.xlsx.number_formats`` via :func:`_xlsx_number_format_masks`);
      - a ``demo_clear`` value was not captured for this template
        (:func:`captured_template_texts`).

    No-ops when overrides are absent (``status != present``), keeping the model-free
    CI path and a pre-B3 profile green - exactly like ``check_comprehension_targets``.
    ``shell`` may be None at verify time only for the reroute/mask branches that need
    to open the package; the demo-clear branch reads the profile's captured set, so it
    fails closed even without a shell.
    """
    overrides = schema.overrides_block(profile)
    if overrides.get("status") != schema.ComprehensionStatus.PRESENT.value:
        return []
    kind = profile.get("kind")
    roles = profile.get("roles") or {}
    role_keys = {r for r in roles if r != "_index"}
    findings: list[Finding] = []

    reroutes = overrides.get("reroute_roles") or {}
    swaps = overrides.get("number_format_swaps") or {}
    clears = overrides.get("demo_clears") or []

    # Pure-profile membership (no shell needed, so it fails closed even at verify time
    # with shell=None or a broken shell): a reroute target must be a CONCRETE declared
    # role (``_index`` is the order array, not a role). Checked unconditionally; the
    # shell-backed resolver proof below only runs the resolver gate on the survivors.
    shell_backed_reroutes: dict[str, str] = {}
    for requested, target in reroutes.items():
        if target not in role_keys:
            findings.append(
                Finding(
                    "override_targets_exist",
                    schema.Severity.ERROR.value,
                    f"reroute target {target!r} (for role {requested!r}) is not a "
                    f"declared role (have {sorted(role_keys)})",
                )
            )
        else:
            shell_backed_reroutes[requested] = target

    # demo-clear values are membership-checked against the profile's OWN captured demo
    # set (no shell open needed) FIRST, so a later shell-open failure can never suppress
    # this fail-closed check. reject-on-empty: a clear value absent from the captured
    # set (including an empty set) is an ERROR.
    if clears:
        captured = set(captured_template_texts(profile))
        for value in clears:
            if value not in captured:
                findings.append(
                    Finding(
                        "override_targets_exist",
                        schema.Severity.ERROR.value,
                        f"demo_clear value {value!r} was not captured for this template",
                    )
                )

    # The resolver-target proof and the mask membership need the shell inventory; open
    # it ONCE. reject-on-empty (NOT skip-on-empty): when a lesson re-points to a target
    # the shell can't prove - including an EMPTY inventory - that is itself an ERROR.
    if (shell_backed_reroutes or swaps) and shell is not None:
        try:
            if kind == schema.Kind.DOCX.value:
                present = _docx_style_keys(shell)
            elif kind == schema.Kind.PPTX.value:
                present = _pptx_layout_names(shell)
            elif kind == schema.Kind.XLSX.value:
                present = _xlsx_defined_names(shell)
            else:
                present = set()
            masks = (
                _xlsx_number_format_masks(shell)
                if kind == schema.Kind.XLSX.value
                else set()
            )
        except Exception as exc:  # opening the shell must never crash the gate
            findings.append(
                Finding(
                    "override_targets_exist",
                    schema.Severity.WARNING.value,
                    f"could not verify override targets against shell: {exc}",
                )
            )
            return findings

        for requested, target in shell_backed_reroutes.items():
            target_resolver = (roles.get(target) or {}).get("resolver") or {}
            target_present = (
                masks
                if (
                    kind == schema.Kind.XLSX.value
                    and target_resolver.get("type")
                    == schema.ResolverType.NUMBER_FORMAT.value
                )
                else present
            )
            if not _role_resolver_target_in_shell(
                kind, target_resolver, target_present
            ):
                findings.append(
                    Finding(
                        "override_targets_exist",
                        schema.Severity.ERROR.value,
                        f"reroute target {target!r} (for role {requested!r}) does not "
                        "resolve to a shell-backed op",
                    )
                )

        for rid, mask in swaps.items():
            # reject-on-empty: a swap into an empty mask inventory is itself an ERROR.
            if mask not in masks:
                findings.append(
                    Finding(
                        "override_targets_exist",
                        schema.Severity.ERROR.value,
                        f"number_format swap mask {mask!r} (for role {rid!r}) is not "
                        f"among the shell's used formats (have {sorted(masks)})",
                    )
                )

    return findings


def check_overrides_applied(profile: dict) -> list[Finding]:
    """Surface every LIVE learned override as an INFO ``override_applied`` finding.

    The audit peer of :func:`check_override_targets`: where that PROVES a present
    lesson's targets against the shell (fail-closed ERROR), this makes a present lesson
    VISIBLE so a learned re-point is never silent and ``verify`` re-surfaces it. It
    gates on the EXACT presence+freeze predicate the resolver consumes on
    (:func:`store.overrides_are_present`), so it emits iff the resolver would actually
    take an override branch: an ``absent`` (advisory, not-yet ``--accept``ed) lesson, a
    drifted-shell stamp, or a pre-B3 profile emits NOTHING (byte-identical CI path).
    Reads ONLY the profile, so it runs in BOTH ``generate`` and ``verify`` QA - a live
    correction is auditable even at verify time when no document is produced.

    Emits one INFO ``override_applied`` per live entry, in deterministic order:
      - each ``reroute_roles`` ``requested -> target`` (location = the requested role
        id, a non-brand stable pointer; the message names the re-point);
      - each ``number_format_swaps`` ``rid -> mask`` (location = the role id; the mask
        is shell-backed so it is safe to name);
      - each ``demo_clears`` value (location = ``demo_clears[<i>]`` - NEVER the value,
        which is captured template/brand text; the message stays value-free).

    INFO severity, NOT in :data:`schema.DEFAULT_L0_INVARIANTS`, never ERROR, so it can
    never flip a verdict; deliberately NOT in :data:`schema.LEARNABLE_CHECKS` (an audit
    trail must not feed itself).
    """
    # Late import keeps this module free of the store layer at import time (mirrors the
    # late ``from brandkit.qa import checks_deterministic`` imports elsewhere).
    from brandkit.profile import store

    if not store.overrides_are_present(profile):
        return []
    block = schema.overrides_block(profile)
    findings: list[Finding] = []

    reroutes = block.get("reroute_roles") or {}
    if isinstance(reroutes, dict):
        for requested in sorted(reroutes):
            findings.append(
                Finding(
                    "override_applied",
                    schema.Severity.INFO.value,
                    f"live override re-points role {requested!r} -> "
                    f"{reroutes[requested]!r}",
                    location=requested,
                )
            )

    swaps = block.get("number_format_swaps") or {}
    if isinstance(swaps, dict):
        for rid in sorted(swaps):
            findings.append(
                Finding(
                    "override_applied",
                    schema.Severity.INFO.value,
                    f"live override swaps number_format of role {rid!r} to "
                    f"{swaps[rid]!r}",
                    location=rid,
                )
            )

    clears = block.get("demo_clears") or []
    if isinstance(clears, list):
        for i in range(len(clears)):
            findings.append(
                Finding(
                    "override_applied",
                    schema.Severity.INFO.value,
                    "live override clears a captured demo value",
                    location=f"demo_clears[{i}]",
                )
            )

    return findings


def check_color_token_targets(profile: dict) -> list[Finding]:
    """FAIL-CLOSED membership of every COLOR token against ``theme.palette``.

    The sibling of :func:`check_comprehension_targets` for model-driven color: a
    color token is only valid when it is a verbatim key of ``theme.palette`` (the
    SOLE namespace a run color resolves off - never ``palette_roles``). Today the
    only place a comprehension references a palette key is
    ``comprehension.palette_annotations`` (the model NAMES a captured color); each
    key must be a real palette entry, mirroring ``check_membership``'s rule for
    anchor/index/region refs - a key into an EMPTY/absent palette is itself an ERROR
    (the model can never invent a color the deterministic capture did not observe).

    Model-free and deterministic. No-ops when the comprehension is absent (status !=
    present), so the model-free CI path and a pre-palette profile stay green.
    """
    comp = _present_comprehension(profile)
    if comp is None:
        return []
    palette = (profile.get("theme") or {}).get("palette") or {}
    palette_keys = set(palette)
    findings: list[Finding] = []
    for key in comp.get("palette_annotations") or {}:
        if key not in palette_keys:
            findings.append(
                Finding(
                    "color_token_targets_exist",
                    schema.Severity.ERROR.value,
                    f"color token {key!r} is not a key of theme.palette "
                    f"(have {sorted(palette_keys)})",
                )
            )
    return findings


def check_palette_alias_targets(profile: dict) -> list[Finding]:
    """FAIL-CLOSED integrity of every model-NAMED palette ALIAS token (Cluster E1).

    The sibling of :func:`check_color_token_targets` for the alias bridge: a model
    NAMES an alias for a captured palette entry and the engine mints a dotted token
    whose ``ref`` is a BYTE-COPY of that entry's ``ref`` (an off-theme ``hex:RRGGBB``
    accent becomes addressable as a clean run-color token). The engine - never the
    model - authors the color, so the alias is valid ONLY when:

      1. the alias token is syntactically legal (a dotted role-id);
      2. the source ``palette_annotations`` key is a real ``theme.palette`` entry
         (already gated by ``check_color_token_targets`` / ``check_membership``);
      3. the minted alias token actually exists in ``theme.palette``;
      4. the minted ref is BYTE-IDENTICAL (``==`` dict equality) to the source
         entry's captured ref - an alias that invented a hex or diverged is an ERROR.

    A collision is impossible by construction (the mint refuses to shadow a non-alias
    key and ``check_membership`` rejects a colliding alias before the mint), but a
    minted token whose ref is not a byte-copy of its declared source is the paranoid
    ENGINE-error this check catches fail-closed.

    Model-free and deterministic. No-ops when the comprehension is absent (status !=
    present) and on annotations carrying no ``alias``, so the model-free CI path, a
    pre-palette profile, and the no-alias byte-identity path all stay green.
    """
    comp = _present_comprehension(profile)
    if comp is None:
        return []
    palette = (profile.get("theme") or {}).get("palette") or {}
    if not isinstance(palette, dict):
        return []
    findings: list[Finding] = []
    for key, ann in (comp.get("palette_annotations") or {}).items():
        if not isinstance(ann, dict):
            continue
        alias = ann.get("alias")
        if not alias:
            continue
        if not isinstance(alias, str) or not schema.is_valid_role_id(alias):
            findings.append(
                Finding(
                    "palette_alias_targets_exist",
                    schema.Severity.ERROR.value,
                    f"palette alias token {alias!r} (for source {key!r}) is not a "
                    "syntactically-legal dotted token",
                )
            )
            continue
        source = palette.get(key)
        source_ref = source.get("ref") if isinstance(source, dict) else None
        minted = palette.get(alias)
        if not isinstance(minted, dict):
            findings.append(
                Finding(
                    "palette_alias_targets_exist",
                    schema.Severity.ERROR.value,
                    f"palette alias token {alias!r} (for source {key!r}) was not "
                    "minted into theme.palette",
                )
            )
            continue
        minted_ref = minted.get("ref")
        if minted_ref != source_ref:
            findings.append(
                Finding(
                    "palette_alias_targets_exist",
                    schema.Severity.ERROR.value,
                    f"palette alias token {alias!r} ref {minted_ref!r} is not a "
                    f"byte-copy of source {key!r} ref {source_ref!r}",
                )
            )
    return findings


def _present_comprehension(profile: dict) -> dict | None:
    comp = profile.get("comprehension")
    if (
        isinstance(comp, dict)
        and comp.get("status") == schema.ComprehensionStatus.PRESENT.value
    ):
        return comp
    return None


def captured_template_texts(
    profile: dict, *, include_surface_prompts: bool = False
) -> list[str]:
    """Collect every demo/placeholder text the extractor captured for THIS template.

    Model-free and language-agnostic: the comparison strings come from the
    template's own captured facts (cover ``demo_value`` s, surfaced demo-region
    markers/text), never a fixed phrase in any language. Surface placeholder
    prompts are opt-in because some formats keep master/layout prompts inside the
    package even when they are not visibly rendered; L0 must not fail on those.
    """
    texts: list[str] = []
    comp = _present_comprehension(profile)
    if comp is not None:
        for slot in (comp.get("cover_slots") or {}).values():
            if isinstance(slot, dict) and slot.get("demo_value"):
                texts.append(str(slot["demo_value"]))
    kind = profile.get("kind")
    sub = ((profile.get("surface") or {}).get(kind) or {}) if kind else {}
    if isinstance(sub, dict):
        demo = sub.get("demo_region") or {}
        if isinstance(demo, dict):
            for m in demo.get("instruction_markers") or []:
                if m:
                    texts.append(str(m))
            if demo.get("start_text"):
                texts.append(str(demo["start_text"]))
        if not include_surface_prompts:
            return _dedup_texts(texts)
        for anchor in sub.get("cover_anchors") or []:
            if not isinstance(anchor, dict):
                continue
            for key in ("placeholder", "demo_value"):
                value = anchor.get(key)
                if value:
                    texts.append(str(value))
    return _dedup_texts(texts)


def _dedup_texts(texts: list[str]) -> list[str]:
    # De-dup preserving order.
    seen: set[str] = set()
    out: list[str] = []
    for t in texts:
        if t and t not in seen:
            seen.add(t)
            out.append(t)
    return out


def _captured_demo_texts(profile: dict) -> list[str]:
    return captured_template_texts(profile)


def check_residual_template_text(text: str, profile: dict) -> list[Finding]:
    """Generalized ``no_residual_template_text`` for ALL kinds (model-free).

    Scans the produced text for any demo/placeholder string captured at extract
    for THIS template (cover ``demo_value`` s + surfaced demo markers/text). No
    hardcoded phrases, language-agnostic, identical across formats.
    """
    findings: list[Finding] = []
    for marker in _captured_demo_texts(profile):
        if marker in text:
            findings.append(
                Finding(
                    "no_residual_template_text",
                    schema.Severity.ERROR.value,
                    f"residual template text: {marker!r}",
                    location=marker,
                )
            )
    return findings


def check_no_orphan_cover_placeholder(text: str, profile: dict) -> list[Finding]:
    """No bound cover slot still shows its captured ``demo_value`` in the output.

    A slot the comprehension marked ``fill_rule='in_place'`` (i.e. it should have
    been filled) whose captured ``demo_value`` still appears in the produced text
    is an ERROR (stale demo prompt left behind). A slot the model intentionally
    re-armed (``fill_rule='clear'``) is INFO. No-ops when comprehension is absent.
    """
    comp = _present_comprehension(profile)
    if comp is None:
        return []
    findings: list[Finding] = []
    for anchor_ref, slot in (comp.get("cover_slots") or {}).items():
        if not isinstance(slot, dict):
            continue
        demo_value = slot.get("demo_value")
        if not demo_value or str(demo_value) not in text:
            continue
        fill_rule = slot.get("fill_rule")
        if fill_rule == schema.FillRule.IN_PLACE.value:
            findings.append(
                Finding(
                    "no_orphan_cover_placeholder",
                    schema.Severity.ERROR.value,
                    f"cover slot {anchor_ref!r} still shows placeholder {demo_value!r}",
                )
            )
        else:
            findings.append(
                Finding(
                    "no_orphan_cover_placeholder",
                    schema.Severity.INFO.value,
                    f"cover slot {anchor_ref!r} retains re-armed prompt {demo_value!r}",
                )
            )
    return findings


def check_index_matches_content(
    present_seq_ids: set[str], profile: dict
) -> list[Finding]:
    """No preserved CAPTION index lacks corresponding captionable content (WARNING).

    Ruling A: the comparison keys on the DETERMINISTIC ``seq_id`` captured opaquely
    from the index's ``\\c`` field switch (plan §7), NEVER on the advisory ``kind``
    token. ``present_seq_ids`` is the set of caption SEQ identifiers the generator
    actually emitted captions for (each is a verbatim ``seq_id`` it matched a
    preserved index against). A preserved CAPTION index (one carrying a non-null
    ``seq_id``) that is kept (``reconcile`` not ``clear``) yet whose ``seq_id`` was
    not emitted is flagged WARNING - a stale caption index the reconciliation should
    have cleared. An outline TOC (``seq_id`` is null) is never flagged: it is always
    refreshed, not content-matched. No-ops when comprehension is absent.
    Generator-driven (the generator knows which SEQ classes it emitted); model-free.
    """
    comp = _present_comprehension(profile)
    if comp is None:
        return []
    findings: list[Finding] = []
    for idx in (comp.get("conventions") or {}).get("indexes") or []:
        if not isinstance(idx, dict):
            continue
        if idx.get("reconcile") == schema.Reconcile.CLEAR.value:
            continue  # the reconciliation already removed it on purpose
        seq_id = idx.get("seq_id")
        # Only a CAPTION index (a \c-switched index, identified by its opaque
        # seq_id) is content-matched; a bare outline TOC has no seq_id and is
        # always refreshed, never flagged.
        if seq_id and seq_id not in present_seq_ids:
            findings.append(
                Finding(
                    "index_matches_content",
                    schema.Severity.WARNING.value,
                    f"preserved caption index {idx.get('index_ref')!r} "
                    f"(seq_id {seq_id!r}) has no matching captionable content",
                )
            )
    return findings


def check_no_net_structure_loss(
    removed_refs: set[str],
    profile: dict,
    *,
    confidence: float | None = None,
) -> list[Finding]:
    """The destructive-action floor (§6): no preserved anchor/index removed unless
    the deterministic layer independently classified it placeholder/demo.

    ``removed_refs`` is the set of anchor/index refs the reconciliation actually
    deleted (supplied by the generator). For each, the comprehension must carry a
    deterministically-corroborated destructive verdict (``fill_rule='clear'`` for a
    cover anchor, ``reconcile='clear'`` for an index, or a ``verdict='demo'`` region);
    otherwise the deletion is an ERROR. Model-free (it reads the frozen verdicts).

    ``confidence`` (optional) is the model's single ``comprehension.confidence``, the
    SAME value the reconcile sites gate on. When supplied, the backstop also
    re-verifies the destructive-action confidence floor: a sanctioned removal whose
    confidence does NOT clear the floor is an ERROR, because under the uniform policy
    such a removal should have been downgraded to KEEP at the reconcile site (a wrong
    delete is unrecoverable). When ``confidence`` is ``None`` the floor re-check is
    skipped (additive, back-compatible) and only the verdict-corroboration gate runs.
    """
    comp = _present_comprehension(profile)
    findings: list[Finding] = []
    cleared_anchors = (
        {
            ref
            for ref, slot in (comp.get("cover_slots") or {}).items()
            if isinstance(slot, dict)
            and slot.get("fill_rule") == schema.FillRule.CLEAR.value
        }
        if comp
        else set()
    )
    cleared_indexes = (
        {
            idx.get("index_ref")
            for idx in ((comp.get("conventions") or {}).get("indexes") or [])
            if isinstance(idx, dict)
            and idx.get("reconcile") == schema.Reconcile.CLEAR.value
        }
        if comp
        else set()
    )
    demo_regions = (
        {
            r.get("region_ref")
            for r in ((comp.get("demo_classification") or {}).get("regions") or [])
            if isinstance(r, dict) and r.get("verdict") == schema.Verdict.DEMO.value
        }
        if comp
        else set()
    )
    sanctioned = cleared_anchors | cleared_indexes | demo_regions
    floor_blocks = confidence is not None and not confidence_clears_floor(confidence)
    for ref in sorted(removed_refs):
        if ref not in sanctioned:
            findings.append(
                Finding(
                    "no_net_structure_loss",
                    schema.Severity.ERROR.value,
                    f"reconciliation removed {ref!r} without a corroborated "
                    f"destructive verdict",
                )
            )
        elif floor_blocks:
            # The verdict is sanctioned, but the model's confidence is below the
            # destructive floor: under the uniform policy the reconcile site should
            # have downgraded this to KEEP. A removal that still happened is a floor
            # breach (a wrong delete is unrecoverable), so surface it LOUDLY.
            findings.append(
                Finding(
                    "no_net_structure_loss",
                    schema.Severity.ERROR.value,
                    f"reconciliation removed {ref!r} below the destructive "
                    f"confidence floor ({confidence:.2f})",
                )
            )
    return findings


# ---------------------------------------------------------------------------
# Shell-vs-output structural diffs (catch silent corruption the text scans miss)
# ---------------------------------------------------------------------------
def _xlsx_formula_map(path) -> dict[str, str]:
    """Map ``Sheet!Coord`` -> formula string for every formula cell in a workbook.

    Reads the live file (``data_only=False``) so the comparison reflects what is
    actually on disk, not a possibly-stale catalog snapshot. A formula is identified
    by ``data_type == "f"`` (authoritative) - NOT by a leading ``=`` in the string,
    so a cell deliberately neutralized to TEXT (an author ``=...`` value written as a
    string literal) is correctly NOT counted as a formula. Iterating ``_cells`` keeps
    this O(populated) on sparse corporate models.
    """
    out: dict[str, str] = {}
    wb = _load_xlsx(path)
    for ws in wb.worksheets:
        for cell in ws._cells.values():
            if cell.data_type == "f" and isinstance(cell.value, str):
                out[f"{ws.title}!{cell.coordinate}"] = cell.value
    return out


def check_formula_preservation(shell, output, profile: dict) -> list[Finding]:
    """ERROR on any formula the generation LOST or MUTATED vs the brand shell.

    The brand guarantee for xlsx is that the shell's own formulas are preserved
    verbatim (only the shell's inputs are refilled, never its formulas). A region
    fill that overwrites a formula cell is silent data corruption that no text
    scan can see, so this diffs the shell's ``address->formula`` set against the
    output's and emits an ERROR for every entry that disappeared or changed.

    Model-free and deterministic. No-ops (returns ``[]``) when either file is
    absent (e.g. verify time, when there is no output yet) so it is always safe to
    call from the gate. The extractor's ``artifact_catalog.formulas`` is the same
    baseline; here we read the live shell so the check stands even if the catalog
    drifts.
    """
    if shell is None or output is None:
        return []
    if profile.get("kind") != schema.Kind.XLSX.value:
        return []
    try:
        shell_formulas = _xlsx_formula_map(shell)
        output_formulas = _xlsx_formula_map(output)
    except Exception as exc:  # opening a workbook must never crash the gate
        return [
            Finding(
                "formula_preservation",
                schema.Severity.WARNING.value,
                f"could not verify formula preservation: {exc}",
            )
        ]
    findings: list[Finding] = []
    for address in sorted(shell_formulas):
        shell_formula = shell_formulas[address]
        out_formula = output_formulas.get(address)
        if out_formula is None:
            findings.append(
                Finding(
                    "formula_preservation",
                    schema.Severity.ERROR.value,
                    f"shell formula at {address} ({shell_formula!r}) was erased in the output",
                    location=address,
                )
            )
        elif out_formula != shell_formula:
            findings.append(
                Finding(
                    "formula_preservation",
                    schema.Severity.ERROR.value,
                    f"shell formula at {address} was mutated: "
                    f"{shell_formula!r} -> {out_formula!r}",
                    location=address,
                )
            )
    # The engine NEVER authors formulas: any formula in the OUTPUT that the shell did
    # not have was introduced by author content (a string starting with '='), which is
    # both an invariant break and a formula-injection risk. Fail closed even if the
    # write-path neutralization ever regresses.
    for address in sorted(set(output_formulas) - set(shell_formulas)):
        findings.append(
            Finding(
                "formula_preservation",
                schema.Severity.ERROR.value,
                f"output has a formula at {address} ({output_formulas[address]!r}) that "
                "the shell did not: the engine never authors formulas (author '='-led "
                "content must be neutralized to text)",
                location=address,
            )
        )
    return findings


def _docx_component_counts(path) -> dict[str, int]:
    # Only count components the generator preserves rather than re-authors. Tables
    # are a meaningful survival signal (a flattened table is a real defect); raw
    # list-paragraph counts are NOT compared, because docx generation rewrites the
    # body from new content, so a different list length is expected, not a loss.
    doc = _load_docx(path)
    return {"tables": len(doc.tables)}


def _pptx_component_counts(path) -> dict[str, int]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = _load_pptx(path)
    tables = charts = pictures = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_table", False):
                tables += 1
            if getattr(shape, "has_chart", False):
                charts += 1
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                pictures += 1
    return {"tables": tables, "charts": charts, "pictures": pictures}


def _xlsx_component_counts(path) -> dict[str, int]:
    wb = _load_xlsx(path)
    tables = charts = 0
    for ws in wb.worksheets:
        try:
            tables += len(ws.tables)
        except Exception:
            pass
        try:
            charts += len(ws._charts)
        except Exception:
            pass
    return {"tables": tables, "charts": charts}


_COMPONENT_COUNTERS = {
    schema.Kind.DOCX.value: _docx_component_counts,
    schema.Kind.PPTX.value: _pptx_component_counts,
    schema.Kind.XLSX.value: _xlsx_component_counts,
}


def check_component_survival(shell, output, profile: dict) -> list[Finding]:
    """WARN when a native component present in the shell has no counterpart in the output.

    Catches the silent down-render class (e.g. pptx flattening a native table to
    text, or a docx list losing its numbering) that no text scan detects: it
    counts native components (tables/charts/lists/pictures, per format) in the
    shell and in the output and emits a WARNING for every component family whose
    output count dropped below the shell's. A WARNING (not ERROR) because some
    drops are legitimate (the new content genuinely has fewer of a component); the
    signal is "you may have lost a native object", surfaced rather than silent.

    Model-free, deterministic, no-ops when either file is absent.
    """
    if shell is None or output is None:
        return []
    counter = _COMPONENT_COUNTERS.get(profile.get("kind"))
    if counter is None:
        return []
    try:
        shell_counts = counter(shell)
        output_counts = counter(output)
    except Exception as exc:  # opening the package must never crash the gate
        return [
            Finding(
                "component_survival",
                schema.Severity.WARNING.value,
                f"could not verify component survival: {exc}",
            )
        ]
    findings: list[Finding] = []
    for family in sorted(shell_counts):
        before = shell_counts[family]
        after = output_counts.get(family, 0)
        if before > 0 and after < before:
            findings.append(
                Finding(
                    "component_survival",
                    schema.Severity.WARNING.value,
                    f"native {family} count dropped {before} -> {after} "
                    f"between shell and output (possible down-render)",
                    # Cluster C2: the dropped FAMILY is the addressable location, so
                    # two dropped families no longer collide on (check, None) and a
                    # model-assisted triage entry can name a single family precisely.
                    location=family,
                )
            )
    return findings


def check_docx(path, profile: dict, shell=None) -> list[Finding]:
    findings = check_profile(profile)
    doc = _load_docx(path)
    text = "\n".join(
        [p.text for p in doc.paragraphs]
        + [cell.text for t in doc.tables for row in t.rows for cell in row.cells]
    )
    for hit in textutil.find_markdown_literals(text):
        findings.append(
            Finding(
                "no_literal_markdown",
                schema.Severity.ERROR.value,
                f"literal markdown leaked: {hit['match']!r}",
            )
        )
    findings.extend(check_residual_template_text(text, profile))
    findings.extend(check_no_orphan_cover_placeholder(text, profile))
    findings.extend(check_component_survival(shell, path, profile))
    return findings


def check_pptx(path, profile: dict, shell=None) -> list[Finding]:
    findings = check_profile(profile)
    prs = _load_pptx(path)
    text = "\n".join(
        shape.text
        for slide in prs.slides
        for shape in slide.shapes
        if hasattr(shape, "text")
    )
    for hit in textutil.find_markdown_literals(text):
        findings.append(
            Finding(
                "no_literal_markdown",
                schema.Severity.ERROR.value,
                f"literal markdown leaked: {hit['match']!r}",
            )
        )
    # Uniform, model-free residual check: compare against THIS template's captured
    # demo/placeholder text, never a fixed phrase (de-literalized).
    findings.extend(check_residual_template_text(text, profile))
    findings.extend(check_no_orphan_cover_placeholder(text, profile))
    findings.extend(check_component_survival(shell, path, profile))
    return findings


def check_xlsx(path, profile: dict, shell=None) -> list[Finding]:
    findings = check_profile(profile)
    wb = _load_xlsx(path)
    cell_texts: list[str] = []
    for ws in wb.worksheets:
        for row in ws.iter_rows():
            for cell in row:
                if isinstance(cell.value, str):
                    cell_texts.append(cell.value)
                    for hit in textutil.find_markdown_literals(cell.value):
                        findings.append(
                            Finding(
                                "no_literal_markdown",
                                schema.Severity.ERROR.value,
                                f"literal markdown leaked: {hit['match']!r}",
                            )
                        )
    # Uniform, model-free residual check across all kinds.
    text = "\n".join(cell_texts)
    findings.extend(check_residual_template_text(text, profile))
    findings.extend(check_no_orphan_cover_placeholder(text, profile))
    # Deterministic structural diffs that the text scan above is blind to: a fill
    # that erased a shell formula, or a native component lost in the output.
    findings.extend(check_formula_preservation(shell, path, profile))
    findings.extend(check_component_survival(shell, path, profile))
    return findings
